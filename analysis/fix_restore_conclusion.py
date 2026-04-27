"""
Slide 25 (결론) 복구:
- 중복된 Slide 25 (국가별 광고집행율) 삭제
- pptx_unpacked3의 원본 결론 슬라이드를 복원
"""
import sys, io, copy, zipfile, os
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE      = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX_IN   = BASE / "MobileGame_Market_Analysis_2022-2026_v2.pptx"
UNPACKED  = BASE / "pptx_unpacked3"
TMP_ORIG  = BASE / "_tmp_orig.pptx"


def pack_orig():
    with zipfile.ZipFile(TMP_ORIG, "w", zipfile.ZIP_DEFLATED) as zf:
        for fp in UNPACKED.rglob("*"):
            if fp.is_file():
                zf.write(fp, fp.relative_to(UNPACKED))


def delete_slide(prs, index: int):
    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(xml_slides[index])


def copy_slide_to_end(src_prs, src_idx: int, dst_prs):
    """src_prs의 src_idx 슬라이드를 dst_prs 맨 끝에 복사."""
    src_slide = src_prs.slides[src_idx]
    blank = dst_prs.slide_layouts[0]
    new_slide = dst_prs.slides.add_slide(blank)

    # spTree 통째로 복사
    new_slide.shapes._spTree.clear()
    for elem in src_slide.shapes._spTree:
        new_slide.shapes._spTree.append(copy.deepcopy(elem))

    # 배경 복사
    try:
        src_bg = src_slide.background.fill
        dst_bg = new_slide.background.fill
        if src_bg.type is not None:
            import lxml.etree as etree
            src_xml = src_slide._element.find(
                './/{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}bg')
            bg_elem = src_slide._element.find(
                './/{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
            if bg_elem is not None:
                new_slide._element.insert(2, copy.deepcopy(bg_elem))
    except Exception:
        pass

    return new_slide


def insert_slide_at(prs, slide, position: int):
    xml_slides = prs.slides._sldIdLst
    last = xml_slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(position, last)


def main():
    # 원본 패킹
    print("▶ 원본 pptx_unpacked3 패킹...")
    pack_orig()
    orig_prs = Presentation(str(TMP_ORIG))
    print(f"  원본 슬라이드 수: {len(orig_prs.slides)}")

    # 결론 슬라이드 = 원본 마지막 (index 16)
    conclusion_idx = len(orig_prs.slides) - 1
    conclusion_text = ""
    for sh in orig_prs.slides[conclusion_idx].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            conclusion_text = sh.text_frame.text.strip()[:50]
            break
    print(f"  결론 슬라이드 (index {conclusion_idx}): '{conclusion_text}'")

    # 현재 PPT 로드
    print(f"\n▶ 현재 PPT 로드: {PPTX_IN.name}")
    prs = Presentation(str(PPTX_IN))
    print(f"  총 {len(prs.slides)}개 슬라이드")

    # Slide 25 (index 24) 현재 내용 확인
    s25_title = ""
    for sh in prs.slides[24].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            s25_title = sh.text_frame.text.strip()[:50]
            break
    print(f"  현재 Slide 25: '{s25_title}'")

    # 중복 Slide 25 삭제
    delete_slide(prs, 24)
    print(f"  Slide 25 삭제 (현재 {len(prs.slides)}개)")

    # 결론 슬라이드 복사해서 끝에 추가
    print("  결론 슬라이드 복원 중...")
    copy_slide_to_end(orig_prs, conclusion_idx, prs)
    # 마지막 위치에 추가됐으므로 이미 올바른 위치 (index 24 = Slide 25)

    # 저장
    prs.save(str(PPTX_IN))
    print(f"\n▶ 저장 완료: {PPTX_IN.name}")
    print(f"  총 {len(prs.slides)}개 슬라이드")

    # 검증
    prs2 = Presentation(str(PPTX_IN))
    for i, s in enumerate(prs2.slides):
        texts = [sh.text_frame.text.strip()[:50] for sh in s.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        if i >= 22:  # Slide 23 이후만 출력
            print(f"  Slide {i+1}: {texts[0] if texts else '(빈 슬라이드)'}")

    TMP_ORIG.unlink(missing_ok=True)


if __name__ == "__main__":
    main()
