"""
두 PPT 파일에 광고집행율 분석 슬라이드 추가/업데이트:
  v3.pptx      → Slide 25 신규 추가 (결론 앞)
  exec report  → Slide 6 업데이트 (국가별 CN 제거, 퍼블리셔별 ad_active 추가)

국가별: KR·JP·US (CN 제외)
퍼블리셔별: 중화권·서구권·일본·한국 (ad_active 비율 포함)
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

rcParams["font.family"]        = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE      = Path("C:/Users/NHN/Documents/sensortower_api")
V3_PPTX   = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"
EXEC_PPTX = BASE / "MobileGame_Executive_Report_2026.pptx"

# ── 데이터 ────────────────────────────────────────────────────────────────
ctr_df = pd.read_csv(BASE / "ad_rate_by_country.csv")
tags   = pd.read_csv(BASE / "app_tags.csv")

# 국가별 (CN 제외)
CTR_ORDER  = ["KR", "JP", "US"]
CTR_LABEL  = {"KR": "KR\n한국", "JP": "JP\n일본", "US": "US\n미국"}
CTR_COLORS = {"KR": "#E74C3C", "JP": "#2ECC71", "US": "#3498DB"}

# 퍼블리셔별 집계 (중화권 포함)
PUB_ORDER  = ["중화권", "서구권", "일본", "한국"]
PUB_COLORS = {"중화권": "#F39C12", "서구권": "#3498DB", "일본": "#2ECC71", "한국": "#E74C3C"}

tags["total_ww"] = tags["paid_display_pct_ww"].fillna(0) + tags["paid_search_pct_ww"].fillna(0)
tags["organic_ww"] = 100 - tags["total_ww"]

pub_agg = (
    tags.groupby("publisher_group")
        .agg(n=("app_id","count"),
             ad_active_n=("ad_active","sum"),
             display=("paid_display_pct_ww","mean"),
             search=("paid_search_pct_ww","mean"),
             total=("total_ww","mean"),
             organic=("organic_ww","mean"))
        .round(1)
)
pub_agg["ad_active_pct"] = (pub_agg["ad_active_n"] / pub_agg["n"] * 100).round(0).astype(int)

ctr_idx = ctr_df.set_index("country")


def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ══════════════════════════════════════════════════════════════════════════
# 차트: 국가별(KR·JP·US) + 퍼블리셔별 광고집행율 + ad_active 비율
# ══════════════════════════════════════════════════════════════════════════
def make_ad_rate_chart(wide: bool = True) -> bytes:
    figsize = (11, 4.8) if wide else (10, 4.5)
    fig, (ax_c, ax_p) = plt.subplots(1, 2, figsize=figsize)
    fig.patch.set_alpha(0)

    legend_handles = [
        mpatches.Patch(facecolor="#888", alpha=0.92, edgecolor="white",
                       label="Paid Display (배너·영상·플레이어블)"),
        mpatches.Patch(facecolor="#888", alpha=0.45, edgecolor="white",
                       hatch="///", label="Paid Search (앱스토어 검색광고)"),
    ]

    # ── 좌: 국가별 (CN 제외) ──────────────────────────────────────────────
    ax_c.set_facecolor("#F8FAFF")
    x   = np.arange(len(CTR_ORDER))
    w   = 0.5
    disp_c = [float(ctr_idx.loc[c, "paid_display_pct_ww"]) for c in CTR_ORDER]
    srch_c = [float(ctr_idx.loc[c, "paid_search_pct_ww"])  for c in CTR_ORDER]
    tot_c  = [d + s for d, s in zip(disp_c, srch_c)]
    org_c  = [100 - t for t in tot_c]
    colors_c = [CTR_COLORS[c] for c in CTR_ORDER]

    b1 = ax_c.bar(x, disp_c, w, color=colors_c, alpha=0.92, edgecolor="white", lw=1.2)
    b2 = ax_c.bar(x, srch_c, w, bottom=disp_c, color=colors_c,
                  alpha=0.45, edgecolor="white", lw=1.2, hatch="///")

    for bar, v in zip(b1, disp_c):
        if v >= 4:
            ax_c.text(bar.get_x() + bar.get_width()/2, bar.get_height()/2,
                      f"{v:.1f}%", ha="center", va="center",
                      fontsize=8.5, fontweight="bold", color="white")
    for bar1, bar2, v in zip(b1, b2, srch_c):
        if v >= 2:
            ax_c.text(bar2.get_x() + bar2.get_width()/2,
                      bar1.get_height() + bar2.get_height()/2,
                      f"{v:.1f}%", ha="center", va="center",
                      fontsize=8, color="white", fontweight="bold")
    for xi, t, org in zip(x, tot_c, org_c):
        ax_c.text(xi, t + 0.8, f"합계 {t:.1f}%", ha="center", va="bottom",
                  fontsize=8.5, fontweight="bold", color="#333")
        ax_c.text(xi, t + 4.2, f"Organic {org:.1f}%", ha="center", va="bottom",
                  fontsize=7.5, color="#888")

    ax_c.set_title("국가별 시장 기준 (앱스토어 TOP 25)", fontsize=11,
                   fontweight="bold", color="#1E2761", pad=6)
    ax_c.set_xticks(x)
    ax_c.set_xticklabels([CTR_LABEL[c] for c in CTR_ORDER], fontsize=10.5, fontweight="bold")
    ax_c.set_ylabel("다운로드 중 비중 (%)", fontsize=9)
    ax_c.set_ylim(0, max(tot_c) * 1.55)
    ax_c.yaxis.grid(True, linestyle="--", alpha=0.35)
    ax_c.spines[["top", "right"]].set_visible(False)
    ax_c.legend(handles=legend_handles, loc="upper right", fontsize=7.5, framealpha=0.85)

    # ── 우: 퍼블리셔별 (중화권 포함, ad_active 비율 표시) ─────────────────
    ax_p.set_facecolor("#F8FAFF")
    pubs    = [p for p in PUB_ORDER if p in pub_agg.index]
    xp      = np.arange(len(pubs))
    colors_p = [PUB_COLORS[p] for p in pubs]

    disp_p = [float(pub_agg.loc[p, "display"]) for p in pubs]
    srch_p = [float(pub_agg.loc[p, "search"])  for p in pubs]
    tot_p  = [float(pub_agg.loc[p, "total"])   for p in pubs]
    org_p  = [float(pub_agg.loc[p, "organic"]) for p in pubs]
    act_p  = [int(pub_agg.loc[p, "ad_active_pct"]) for p in pubs]
    n_p    = [int(pub_agg.loc[p, "n"])             for p in pubs]

    p1 = ax_p.bar(xp, disp_p, w, color=colors_p, alpha=0.92, edgecolor="white", lw=1.2)
    p2 = ax_p.bar(xp, srch_p, w, bottom=disp_p, color=colors_p,
                  alpha=0.45, edgecolor="white", lw=1.2, hatch="///")

    for bar, v in zip(p1, disp_p):
        if v >= 4:
            ax_p.text(bar.get_x() + bar.get_width()/2, bar.get_height()/2,
                      f"{v:.1f}%", ha="center", va="center",
                      fontsize=8.5, fontweight="bold", color="white")
    for bar1, bar2, v in zip(p1, p2, srch_p):
        if v >= 2:
            ax_p.text(bar2.get_x() + bar2.get_width()/2,
                      bar1.get_height() + bar2.get_height()/2,
                      f"{v:.1f}%", ha="center", va="center",
                      fontsize=8, color="white", fontweight="bold")
    for xi, t, org, act, n in zip(xp, tot_p, org_p, act_p, n_p):
        ax_p.text(xi, t + 0.8, f"합계 {t:.1f}%", ha="center", va="bottom",
                  fontsize=8.5, fontweight="bold", color="#333")
        ax_p.text(xi, t + 4.2, f"Organic {org:.1f}%", ha="center", va="bottom",
                  fontsize=7.5, color="#888")
        # ad_active 비율 — 막대 하단에 배지로 표시
        badge_color = "#0D9488" if act >= 80 else "#F59E0B" if act >= 50 else "#E74C3C"
        ax_p.text(xi, -5.5, f"광고집행 중\n{act}%",
                  ha="center", va="top", fontsize=7.5, fontweight="bold",
                  color="white",
                  bbox=dict(boxstyle="round,pad=0.3", facecolor=badge_color,
                            edgecolor="none", alpha=0.9))

    xlbls_p = [f"{p}\n(n={n})" for p, n in zip(pubs, n_p)]
    ax_p.set_title("퍼블리셔별 기준 (WW 평균)", fontsize=11,
                   fontweight="bold", color="#1E2761", pad=6)
    ax_p.set_xticks(xp)
    ax_p.set_xticklabels(xlbls_p, fontsize=10, fontweight="bold")
    ax_p.set_ylabel("다운로드 중 비중 (%)", fontsize=9)
    ax_p.set_ylim(-12, max(tot_p) * 1.55)
    ax_p.yaxis.grid(True, linestyle="--", alpha=0.35)
    ax_p.spines[["top", "right"]].set_visible(False)
    ax_p.legend(handles=legend_handles, loc="upper right", fontsize=7.5, framealpha=0.85)

    fig.text(0.01, 0.0,
             "광고집행 중(%) = 해당 퍼블리셔 그룹 내 현재 광고 활성 앱 비율  |  "
             "국가별: 앱스토어 TOP 25 기준 (2026.01)  |  퍼블리셔별: app_tags WW 평균  |  신규게임 전용 데이터 미수집",
             fontsize=7, color="#888")
    fig.tight_layout(rect=[0, 0.05, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# v3.pptx — Slide 25 신규 추가 (결론 앞)
# ══════════════════════════════════════════════════════════════════════════
def add_slide_at(prs, position: int):
    idx    = min(0, len(prs.slide_layouts) - 1)
    layout = prs.slide_layouts[idx]
    slide  = prs.slides.add_slide(layout)
    sl_lst = prs.slides._sldIdLst
    el     = sl_lst[-1]
    sl_lst.remove(el)
    sl_lst.insert(position, el)
    return slide


def build_v3_slide(prs, img_bytes: bytes):
    """Slide 25(index 24) 위치에 신규 슬라이드 삽입."""
    slide = add_slide_at(prs, position=24)

    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFF)

    # 제목 바
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(9.1), Inches(0.62))
    tf = tb.text_frame; tf.word_wrap = False
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "광고 집행율 & Paid Install — 국가별(KR·JP·US) · 퍼블리셔별"
    run.font.size = Pt(19); run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)

    # 부제목
    sub = slide.shapes.add_textbox(Inches(0.3), Inches(0.74), Inches(9.1), Inches(0.26))
    r   = sub.text_frame.paragraphs[0].add_run()
    r.text = ("국가별: 각 앱스토어 TOP 25 기준 (2026.01, CN 제외)  |  "
              "퍼블리셔별: app_tags WW 평균  |  '광고집행 중 %' = 현재 광고 활성 앱 비율")
    r.font.size = Pt(8); r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # 차트 이미지
    slide.shapes.add_picture(io.BytesIO(img_bytes),
        Inches(0.25), Inches(1.05), Inches(9.2), Inches(5.15))


# ══════════════════════════════════════════════════════════════════════════
# exec report — Slide 6 업데이트 (기존 이미지 교체)
# ══════════════════════════════════════════════════════════════════════════
def update_exec_slide6(prs, img_bytes: bytes):
    """Slide 6 (index 5) 의 퍼블리셔 테이블 부분을 새 통합 차트로 교체."""
    slide = prs.slides[5]
    # 기존 도형 중 퍼블리셔 테이블 관련 도형들 제거 후 차트 이미지 추가
    # → 실제로는 기존 Slide 6 텍스트/도형 구조를 유지하면서 하단 퍼블리셔 테이블만 교체
    # 여기서는 슬라이드 전체를 새 차트 이미지로 덮는 방식 대신
    # 새로운 전용 슬라이드를 추가하는 방식으로 처리 (기존 Slide 6 보존)
    # → Slide 6 바로 뒤(index 6)에 신규 슬라이드 삽입
    pass  # 아래 build_exec_new_slide 사용


def build_exec_new_slide(prs, img_bytes: bytes):
    """Slide 6 다음(index 6 위치)에 신규 슬라이드 삽입."""
    from pptx.util import Emu
    NAVY  = "1E2761"
    TEAL  = "0D9488"
    WHITE = "FFFFFF"
    GRAY  = "6B7280"
    SILVER= "94A3B8"

    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # 슬라이드를 index 6 위치로 이동
    sl_lst = prs.slides._sldIdLst
    el     = sl_lst[-1]
    sl_lst.remove(el)
    sl_lst.insert(6, el)

    # 배경
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 헤더 바
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as RGB

    slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(10), Inches(1.1)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = RGB(0x1E, 0x27, 0x61)
    slide.shapes[-1].line.color.rgb      = RGB(0x1E, 0x27, 0x61)

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.08), Inches(9.2), Inches(0.58))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "광고 집행율 — 국가별(KR·JP·US) · 퍼블리셔별"
    run.font.size = Pt(22); run.font.bold = True
    run.font.color.rgb = RGB(0xFF, 0xFF, 0xFF)

    sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.68), Inches(9.2), Inches(0.30))
    sr  = sub.text_frame.paragraphs[0].add_run()
    sr.text = ("국가별: 앱스토어 TOP 25 기준 (2026.01, CN 제외)  |  "
               "퍼블리셔별: app_tags WW 평균  |  '광고집행 중 %' = 현재 광고 활성 앱 비율")
    sr.font.size = Pt(8.5); sr.font.color.rgb = RGB(0xCA, 0xDC, 0xFC)

    # 차트 이미지
    slide.shapes.add_picture(io.BytesIO(img_bytes),
        Inches(0.25), Inches(1.1), Inches(9.2), Inches(4.95))

    # 출처 텍스트
    ft = slide.shapes.add_textbox(Inches(0.3), Inches(5.18), Inches(9.4), Inches(0.22))
    fr = ft.text_frame.paragraphs[0].add_run()
    fr.text = "Source: Sensor Tower API  |  iOS TOP 100 기준  |  2022–2026.03"
    fr.font.size = Pt(8); fr.font.color.rgb = RGB(0x94, 0xA3, 0xB8)


# ══════════════════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════════════════
def main():
    print("▶ 차트 생성 중...")
    img = make_ad_rate_chart(wide=True)
    print("  ✔ 광고집행율 통합 차트 생성 완료")

    # ── v3.pptx 업데이트 ──────────────────────────────────────────────────
    print("\n▶ v3.pptx — Slide 25 신규 추가 중...")
    prs_v3 = Presentation(str(V3_PPTX))
    print(f"  현재 슬라이드 수: {len(prs_v3.slides)}")
    build_v3_slide(prs_v3, img)
    prs_v3.save(str(V3_PPTX))
    prs_v3 = Presentation(str(V3_PPTX))
    print(f"  저장 후 슬라이드 수: {len(prs_v3.slides)}")
    for i in range(22, min(27, len(prs_v3.slides))):
        s = prs_v3.slides[i]
        t = [sh.text_frame.text.strip()[:55] for sh in s.shapes
             if sh.has_text_frame and sh.text_frame.text.strip()]
        imgs = sum(1 for sh in s.shapes if sh.shape_type == 13)
        print(f"  Slide {i+1}: {t[0] if t else '(빈)'} | 이미지:{imgs}")

    # ── exec report 업데이트 ──────────────────────────────────────────────
    print(f"\n▶ Executive Report — 신규 슬라이드 추가 (Slide 7 위치)...")
    prs_ex = Presentation(str(EXEC_PPTX))
    print(f"  현재 슬라이드 수: {len(prs_ex.slides)}")
    build_exec_new_slide(prs_ex, img)
    prs_ex.save(str(EXEC_PPTX))
    prs_ex = Presentation(str(EXEC_PPTX))
    print(f"  저장 후 슬라이드 수: {len(prs_ex.slides)}")
    for i, s in enumerate(prs_ex.slides):
        t = [sh.text_frame.text.strip()[:50] for sh in s.shapes
             if sh.has_text_frame and sh.text_frame.text.strip()]
        imgs = sum(1 for sh in s.shapes if sh.shape_type == 13)
        print(f"  Slide {i+1}: {t[0] if t else '(빈)'} | 이미지:{imgs}")

    print("\n✔ 완료")


if __name__ == "__main__":
    main()
