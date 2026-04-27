"""
v2.pptx 슬라이드별 수정:
  Slide 15 - n 의미 주석
  Slide 16 - n 의미 + 중화권/서구권 포함 국가
  Slide 18 - KR/JP 겹침 해결 (선 스타일 구분) + 국가 설명
  Slide 19 - 중화권 데이터 없음 주석 + 포함 국가 설명
  Slide 21 - CN = 중국 시장 기준 설명
  Slide 22 - 중화권/서구권 포함 국가
  Slide 23 - Paid Display vs Paid Search 차이 설명
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

rcParams["font.family"]     = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE    = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX_IN = BASE / "MobileGame_Market_Analysis_2022-2026_v2.pptx"

NAVY = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x0D, 0x94, 0x88)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

COUNTRY_COLORS   = {"KR": "#E74C3C", "US": "#3498DB", "JP": "#2ECC71", "CN": "#F39C12"}
PUBLISHER_COLORS = {"한국": "#E74C3C", "일본": "#2ECC71", "중화권": "#F39C12",
                    "서구권": "#3498DB", "기타": "#95A5A6"}

# ── 헬퍼 ──────────────────────────────────────────────────────────────────────

def add_note_box(slide, text: str, top_inches: float = 6.55,
                 bg_color=None, font_size=8.5, bold=False,
                 left=0.25, width=9.2, height=0.5):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top_inches), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = GRAY
    if bg_color:
        tb.fill.solid()
        r, g, b = [int(bg_color[i:i+2], 16) for i in (1, 3, 5)]
        tb.fill.fore_color.rgb = RGBColor(r, g, b)
    return tb


def add_badge(slide, text: str, left, top, width=2.5, height=0.32,
              bg="#1E2761", fg="#FFFFFF", font_size=8):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = True
    run.font.color.rgb = RGBColor(*[int(fg[i:i+2], 16) for i in (1, 3, 5)])
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor(*[int(bg[i:i+2], 16) for i in (1, 3, 5)])
    return tb


def fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def replace_chart_image(slide, new_img_bytes: bytes):
    """슬라이드에서 첫 번째 그림 도형을 찾아 이미지를 교체 (위치/크기 유지)."""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            left   = shape.left
            top    = shape.top
            width  = shape.width
            height = shape.height
            # 기존 도형 제거
            sp = shape._element
            sp.getparent().remove(sp)
            # 새 이미지 삽입
            slide.shapes.add_picture(
                io.BytesIO(new_img_bytes), left, top, width, height)
            return True
    return False


# ── 차트 재생성 ───────────────────────────────────────────────────────────────

def make_retention_chart_fixed(df2: pd.DataFrame, group_col: str,
                               groups: list, colors: dict,
                               note_post: str = "") -> bytes:
    """잔존율 라인 차트 — KR/JP 겹침 구분, 빠진 그룹 주석 포함."""
    days_cols   = ["d1", "d7", "d14", "d30"]
    days_labels = ["D1", "D7", "D14", "D30"]

    pre  = df2[df2["period"] == "pre2025"]
    post = df2[df2["period"] == "post2025"]

    # 선 스타일: 겹칠 경우 구분용
    line_styles = ["-", "--", "-.", ":", (0,(3,1,1,1))]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.8))
    fig.patch.set_alpha(0)

    for ax, period_df, period_label in [
            (axes[0], pre,  "2025 이전  (2022–2024)"),
            (axes[1], post, "2025 이후  (2025–2026)")]:

        plotted = 0
        for i, grp in enumerate(groups):
            sub = period_df[period_df[group_col] == grp]
            if sub.empty:
                continue
            vals = [sub[c].mean() for c in days_cols]
            ls   = line_styles[i % len(line_styles)]
            ax.plot(days_labels, vals, marker="o", linewidth=2.2,
                    linestyle=ls, color=colors.get(grp, "#95A5A6"),
                    label=grp, zorder=10 - i)
            # 마지막 포인트에 레이블
            ax.annotate(f" {grp}({vals[-1]:.1f}%)",
                        xy=(days_labels[-1], vals[-1]),
                        fontsize=7.5, color=colors.get(grp, "#95A5A6"),
                        va="center")
            plotted += 1

        ax.set_title(period_label, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_ylabel("잔존율 (%)", fontsize=9)
        ax.set_ylim(0, 65)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.set_facecolor("#F8FAFF")
        ax.spines[["top", "right"]].set_visible(False)
        ax.legend(loc="upper right", fontsize=8, framealpha=0.7)

        # 2025이후 패널에 추가 주석
        if period_label.startswith("2025 이후") and note_post:
            ax.text(0.5, 0.97, note_post,
                    transform=ax.transAxes, ha="center", va="top",
                    fontsize=8, color="#E74C3C",
                    bbox=dict(boxstyle="round,pad=0.3",
                              facecolor="#FFF3F3", edgecolor="#E74C3C", alpha=0.85))

    fig.tight_layout()
    return fig_to_bytes(fig)


# ── 텍스트 상수 ───────────────────────────────────────────────────────────────

N_NOTE = "※ n = 해당 기간 TOP 100 차트에 신규 진입한 게임 수 (출시월 = 첫 차트 진입월 기준)"

CHINESE_NOTE = "중화권 = 중국(CN) · 홍콩(HK) · 대만(TW) · 마카오(MO)"
WESTERN_NOTE = "서구권 = 미국 · 캐나다 · 영국 · 독일 · 프랑스 · 핀란드 · 스웨덴 · 호주 · 아일랜드 외 서유럽·오세아니아"
REGION_NOTE  = f"{CHINESE_NOTE}    |    {WESTERN_NOTE}"

PAID_DIFF_NOTE = (
    "Paid Display: 배너·영상·플레이어블 광고를 보고 설치 — 광범위한 노출로 인지도 확보 목적\n"
    "Paid Search : 앱스토어 검색광고(Apple Search Ads 등)를 통해 설치 — 구매의향 높은 유저 타겟\n"
    "※ 두 합산 = 전체 다운로드 중 광고를 통한 유료 설치 비중 (Organic은 나머지 %)"
)


# ── MAIN ──────────────────────────────────────────────────────────────────────

def main():
    df2 = pd.read_csv(BASE / "slide2_v2.csv")

    print(f"▶ PPT 로드: {PPTX_IN.name}")
    prs = Presentation(str(PPTX_IN))
    print(f"  총 {len(prs.slides)}개 슬라이드")

    slides = prs.slides   # 0-based index

    # ── Slide 15 (index 14): n 의미 ──────────────────────────────────────────
    print("  [Slide 15] n 의미 주석 추가...")
    add_note_box(slides[14], N_NOTE, top_inches=6.6)

    # ── Slide 16 (index 15): n 의미 + 중화권/서구권 국가 ─────────────────────
    print("  [Slide 16] n 의미 + 포함 국가 주석 추가...")
    add_note_box(slides[15], N_NOTE, top_inches=6.3, height=0.28)
    add_note_box(slides[15], REGION_NOTE, top_inches=6.6, height=0.45)

    # ── Slide 18 (index 17): KR/JP 겹침 해결 ─────────────────────────────────
    print("  [Slide 18] 국가별 잔존율 차트 재생성 (KR/JP 겹침 구분)...")
    country_groups = ["KR", "US", "JP", "CN"]
    note18 = "KR · JP 2025이후 값 동일 → 선 스타일로 구분 (실선=JP, 점선=KR)"
    img18  = make_retention_chart_fixed(
        df2, group_col="country", groups=country_groups,
        colors=COUNTRY_COLORS, note_post=note18)
    replaced = replace_chart_image(slides[17], img18)
    print(f"    차트 교체: {'성공' if replaced else '이미지 없음 — 새로 추가'}")
    if not replaced:
        slides[17].shapes.add_picture(io.BytesIO(img18),
            Inches(0.25), Inches(1.1), Inches(9.2), Inches(5.4))
    add_note_box(slides[17],
        "※ KR과 JP의 2025 이후 잔존율이 동일한 값으로 집계됨 (선 스타일로 구분: JP=실선, KR=점선)",
        top_inches=6.6)

    # ── Slide 19 (index 18): 중화권 없음 주석 + 차트 재생성 ─────────────────
    print("  [Slide 19] 퍼블리셔별 잔존율 차트 재생성 + 중화권 주석...")
    pub_groups = ["한국", "일본", "중화권", "서구권"]
    note19 = "중화권: 2025 이후 신규진입 게임 잔존율 데이터 없음"
    img19  = make_retention_chart_fixed(
        df2, group_col="publisher_group", groups=pub_groups,
        colors=PUBLISHER_COLORS, note_post=note19)
    replaced19 = replace_chart_image(slides[18], img19)
    if not replaced19:
        slides[18].shapes.add_picture(io.BytesIO(img19),
            Inches(0.25), Inches(1.1), Inches(9.2), Inches(5.4))
    add_note_box(slides[18],
        "※ 중화권: 2025 이후 TOP 100 신규 진입 중화권 퍼블리셔 게임의 잔존율 수집 데이터 없음",
        top_inches=6.35, height=0.28)
    add_note_box(slides[18], REGION_NOTE, top_inches=6.62, height=0.45)

    # ── Slide 21 (index 20): CN = 중국 시장 기준 설명 ────────────────────────
    print("  [Slide 21] CN 기준 설명 추가...")
    add_note_box(slides[20],
        "※ CN = 중국 앱스토어 차트 기준 (퍼블리셔 국적이 아닌 서비스 국가 기준)  |  "
        "KR=한국 · US=미국 · JP=일본 · CN=중국 시장",
        top_inches=6.62)

    # ── Slide 22 (index 21): 중화권/서구권 포함 국가 ─────────────────────────
    print("  [Slide 22] 포함 국가 주석 추가...")
    add_note_box(slides[21], REGION_NOTE, top_inches=6.62, height=0.45)

    # ── Slide 23 (index 22): Paid Display vs Paid Search 차이 ───────────────
    print("  [Slide 23] Paid Display/Search 차이 설명 추가...")
    add_note_box(slides[22], PAID_DIFF_NOTE,
                 top_inches=6.1, height=0.85, font_size=8)

    # ── 저장 ─────────────────────────────────────────────────────────────────
    out = PPTX_IN   # 같은 파일 덮어쓰기 (백업이 필요하면 _v2b 로 변경)
    prs.save(str(out))
    print(f"\n▶ 저장 완료: {out.name}")


if __name__ == "__main__":
    main()
