"""
슬라이드 15~22 전면 재작성 — iOS·AOS 분리, pre2025/post2025 연평균 비교

  Slide 15: 국가별(KR/JP/US) 3개월 생존율 + 연평균 게임수
  Slide 16: 국가별 광고율 (3개월 생존 vs 미생존)
  Slide 17: 퍼블리셔별 3개월 생존율
  Slide 18: 퍼블리셔별 광고율 (생존 vs 미생존)
  Slide 19: 국가별 잔존율 M+1~M+12
  Slide 20: 퍼블리셔별 잔존율 M+1~M+12
  Slide 21: 국가별 순위별 월평균 매출 (1,10,20,50,100위)
  Slide 22: 퍼블리셔별 순위별 월평균 매출

집계: pre2025(2022~2024)÷3 vs post2025(2025~2026.3)÷1.25
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.lines as mlines
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

rcParams["font.family"]        = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"
DPI  = 200

# ── 색상 ────────────────────────────────────────────────────────────────
PUB_COLORS = {"중화권":"#F39C12","서구권":"#3498DB","일본":"#2ECC71","한국":"#E74C3C","기타":"#9B59B6"}
PUB_ORDER  = ["한국","일본","서구권","중화권"]
CTR_COLORS = {"KR":"#E74C3C","JP":"#2ECC71","US":"#3498DB"}
CTR_LABEL  = {"KR":"KR 한국","JP":"JP 일본","US":"US 미국"}
CTR_ORDER  = ["KR","JP","US"]
PERIOD_LBL = {"pre2025":"2022–2024\n(연평균)","post2025":"2025~\n(연평균)"}
IOS_ALPHA  = 0.88;  AOS_ALPHA = 0.55
NOTE_FS    = 7.5

# ── 데이터 로드 ─────────────────────────────────────────────────────────
ios_s1   = pd.read_csv(BASE / "slide1_v2.csv")
aod_s1   = pd.read_csv(BASE / "slide1_v2_android.csv")
ios_ad   = pd.read_csv(BASE / "newgame_ad_rates.csv")
aod_ad   = pd.read_csv(BASE / "newgame_ad_rates_android.csv")
ios_ret  = pd.read_csv(BASE / "slide2_v2.csv")
ios_tags = pd.read_csv(BASE / "app_tags.csv")
aod_tags = pd.read_csv(BASE / "app_tags_android.csv")
ios_rev  = pd.read_csv(BASE / "slide3_v2.csv")

# Android revenue — load if exists
_aod_rev_path = BASE / "slide3_rank_revenue_android.csv"
aod_rev = pd.read_csv(_aod_rev_path) if _aod_rev_path.exists() else pd.DataFrame()

# Merge survival + ad rates
ios_s1_ad = ios_s1.merge(ios_ad[["app_id","total_ww"]], on="app_id", how="left")
aod_s1_ad = aod_s1.merge(aod_ad[["app_id","total_ww"]], on="app_id", how="left")
ios_s1_ad["has_ad"] = ios_s1_ad["total_ww"].fillna(0) > 0
aod_s1_ad["has_ad"] = aod_s1_ad["total_ww"].fillna(0) > 0


def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=DPI, bbox_inches="tight",
                facecolor="white", transparent=False)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ── 연평균 환산 ─────────────────────────────────────────────────────────
PRE_YEARS  = 3     # 2022~2024
POST_YEARS = 1.25  # 2025~2026.3


def annual_avg(total, period):
    return total / PRE_YEARS if period == "pre2025" else total / POST_YEARS


# ── 슬라이드 유틸 ───────────────────────────────────────────────────────
FOOTNOTE_KW = [
    "출처:", "Source:", "신규진입 기준", "App Store TOP", "Google Play TOP",
    "중화권·서구권", "중화권=", "서구권=", "slide3_v2", "app_tags",
    "2025 전 =", "CN 제외", "KR=한국", "TOP 100", "Sensor Tower",
]

def _is_footnote(shape) -> bool:
    if not shape.has_text_frame:
        return False
    top_in = shape.top / 914400
    if top_in < 1.4:
        return False
    t = shape.text_frame.text
    return any(kw in t for kw in FOOTNOTE_KW)


def maximize_chart(slide, img_bytes, note_text="", prs_h=5.625):
    """각주 제거 + 차트 최대화 삽입 + 하단 해석·인사이트 텍스트 추가."""
    to_remove = []
    for sh in slide.shapes:
        if sh.shape_type == 13 or _is_footnote(sh):
            to_remove.append(sh)
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)

    title_bot = 0.85
    for sh in slide.shapes:
        if sh.has_text_frame and sh.top / 914400 < 1.5:
            b = (sh.top + sh.height) / 914400
            title_bot = max(title_bot, b)

    margin = 0.18
    top  = title_bot + 0.05
    left = margin
    w    = 10.0 - 2 * margin

    if note_text:
        note_h = 0.55
        h = prs_h - top - margin - note_h
        slide.shapes.add_picture(
            io.BytesIO(img_bytes), Inches(left), Inches(top), Inches(w), Inches(h))
        txb = slide.shapes.add_textbox(
            Inches(left), Inches(top + h + 0.02), Inches(w), Inches(note_h))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = note_text
        run = p.runs[0]
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    else:
        h = prs_h - top - margin
        slide.shapes.add_picture(
            io.BytesIO(img_bytes), Inches(left), Inches(top), Inches(w), Inches(h))


def find_slide(prs, keyword, fallback_idx=None):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and keyword in sh.text_frame.text:
                return i, s
    if fallback_idx is not None and fallback_idx < len(prs.slides):
        return fallback_idx, prs.slides[fallback_idx]
    return None, None


def survival_by(df, group_col):
    g = df.groupby([group_col, "period"])["survived_3m"].agg(["sum","count"])
    g["rate"] = (g["sum"] / g["count"] * 100).round(1)
    g["annual_n"] = g.apply(lambda r: annual_avg(r["count"],
                            r.name[1] if isinstance(r.name, tuple) else "pre2025"), axis=1).round(0)
    return g.reset_index()


# ══════════════════════════════════════════════════════════════════════════
# Slide 15: 국가별 3개월 생존율 + 연평균 게임수
# ══════════════════════════════════════════════════════════════════════════
def chart_s15() -> bytes:
    ios_g = survival_by(ios_s1[ios_s1["country"].isin(CTR_ORDER)], "country")
    aod_g = survival_by(aod_s1[aod_s1["country"].isin(CTR_ORDER)], "country")
    periods = [("pre2025","2022–2024"), ("post2025","2025~")]

    fig, axes = plt.subplots(1, 3, figsize=(13.5, 4.6), gridspec_kw={"wspace":0.44})
    fig.patch.set_facecolor("white")

    ios_patch = mpatches.Patch(facecolor="#2980B9", alpha=IOS_ALPHA, label="iOS")
    aod_patch = mpatches.Patch(facecolor="#27AE60", alpha=AOS_ALPHA, hatch="///", label="Android")

    W = 0.40; INNER = 0.54; OUTER = 1.45

    for ax, country in zip(axes, CTR_ORDER):
        col = CTR_COLORS[country]
        ax.set_facecolor("#F8FAFF")
        g_centers, g_labels = [], []

        for gi, (pk, pl) in enumerate(periods):
            xc_i = gi * OUTER
            xc_a = gi * OUTER + INNER

            si = ios_g[(ios_g["country"]==country)&(ios_g["period"]==pk)]
            sa = aod_g[(aod_g["country"]==country)&(aod_g["period"]==pk)]
            ri = float(si["rate"].values[0]) if len(si) else 0
            ni = int(si["annual_n"].values[0]) if len(si) else 0
            ra = float(sa["rate"].values[0]) if len(sa) else 0
            na = int(sa["annual_n"].values[0]) if len(sa) else 0

            ax.bar(xc_i, ri, W, color="#2980B9", alpha=IOS_ALPHA, edgecolor="white", lw=1.3)
            ax.bar(xc_a, ra, W, color="#27AE60", alpha=AOS_ALPHA, hatch="///", edgecolor="white", lw=1.3)

            for xc, rate, clr in [(xc_i,ri,"#2980B9"),(xc_a,ra,"#27AE60")]:
                if rate > 1:
                    ax.text(xc, rate+1.8, f"{rate:.0f}%", ha="center", va="bottom",
                            fontsize=12, fontweight="bold", color=clr)

            gc = (xc_i + xc_a) / 2
            g_centers.append(gc)
            g_labels.append(f"{pl}\niOS {ni}개/yr  AOS {na}개/yr")

        right0 = 0*OUTER + INNER + W/2
        left1  = 1*OUTER - W/2
        ax.axvline((right0+left1)/2, color="#CBD5E1", lw=1.2, ls="--")

        ax.set_title(CTR_LABEL[country], fontsize=13, fontweight="bold", color=col, pad=8)
        ax.set_xticks(g_centers)
        ax.set_xticklabels(g_labels, fontsize=8.5, fontweight="bold", ha="center")
        ax.set_xlim(0*OUTER - W/2 - 0.22, 1*OUTER + INNER + W/2 + 0.22)
        ax.set_ylim(0, 118)
        ax.set_ylabel("3개월 생존율 (%)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35)
        ax.spines[["top","right"]].set_visible(False)

    fig.legend(handles=[ios_patch, aod_patch], loc="upper right",
               bbox_to_anchor=(0.995,0.98), fontsize=10, framealpha=0.9)
    fig.text(0.5, 0.005,
             "매출 TOP 100 신규진입 기준 (첫 차트 진입월)  ·  iOS=App Store / AOS=Google Play  ·  출처: Sensor Tower",
             ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 16: 국가별 광고율 (생존 vs 미생존)
# ══════════════════════════════════════════════════════════════════════════
def chart_s16() -> bytes:
    fig, axes = plt.subplots(1, 3, figsize=(13.5, 4.6), gridspec_kw={"wspace":0.42})
    fig.patch.set_facecolor("white")

    surv_patch = mpatches.Patch(facecolor="#2ECC71", alpha=0.85, label="3개월 생존")
    churn_patch = mpatches.Patch(facecolor="#E74C3C", alpha=0.55, hatch="///", label="3개월 미생존")

    for ax, country in zip(axes, CTR_ORDER):
        ax.set_facecolor("#F8FAFF")
        col = CTR_COLORS[country]

        W = 0.35; INNER = 0.45; OUTER = 1.40
        g_centers, g_labels = [], []

        for gi, (pk, pl) in enumerate([("pre2025","2022–2024"),("post2025","2025~")]):
            for plat_df, plat_name, offset in [(ios_s1_ad,"iOS",0), (aod_s1_ad,"AOS",2)]:
                sub = plat_df[(plat_df["country"]==country)&(plat_df["period"]==pk)&(plat_df["total_ww"].notna())]
                surv = sub[sub["survived_3m"]==True]
                churn = sub[sub["survived_3m"]==False]
                r_surv  = (surv["has_ad"].mean()*100) if len(surv) else 0
                r_churn = (churn["has_ad"].mean()*100) if len(churn) else 0

                xbase = gi * OUTER * 2 + offset * 0.50
                xs = xbase
                xc = xbase + INNER

                ax.bar(xs, r_surv, W, color="#2ECC71", alpha=0.85, edgecolor="white", lw=1.2)
                ax.bar(xc, r_churn, W, color="#E74C3C", alpha=0.55, hatch="///", edgecolor="white", lw=1.2)

                for x, r, c in [(xs,r_surv,"#2ECC71"),(xc,r_churn,"#E74C3C")]:
                    if r > 1:
                        ax.text(x, r+2, f"{r:.0f}%", ha="center", va="bottom",
                                fontsize=9, fontweight="bold", color=c)

                mid = (xs+xc)/2
                g_centers.append(mid)
                g_labels.append(f"{plat_name}\n{pl}")

        ax.set_title(CTR_LABEL[country], fontsize=13, fontweight="bold", color=col, pad=8)
        ax.set_xticks(g_centers)
        ax.set_xticklabels(g_labels, fontsize=7.5, ha="center")
        ax.set_ylim(0, 120)
        ax.set_ylabel("광고집행율 (%)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35)
        ax.spines[["top","right"]].set_visible(False)

    fig.legend(handles=[surv_patch, churn_patch], loc="upper right",
               bbox_to_anchor=(0.995,0.98), fontsize=10, framealpha=0.9)
    fig.text(0.5, 0.005,
             "광고집행율 = Paid Install(Display+Search) 비중이 0% 초과인 게임 비율  ·  출처: Sensor Tower",
             ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 17: 퍼블리셔별 3개월 생존율
# ══════════════════════════════════════════════════════════════════════════
def chart_s17() -> bytes:
    ios_g = survival_by(ios_s1, "publisher_group")
    aod_g = survival_by(aod_s1, "publisher_group")
    periods = [("pre2025","2022–2024"), ("post2025","2025~")]

    fig, axes = plt.subplots(1, 2, figsize=(13, 4.5), gridspec_kw={"wspace":0.38})
    fig.patch.set_facecolor("white")

    ios_patch = mpatches.Patch(facecolor="#444", alpha=IOS_ALPHA, label="iOS")
    aod_patch = mpatches.Patch(facecolor="#444", alpha=AOS_ALPHA, hatch="///", label="Android")

    for ax, (pk, pl) in zip(axes, periods):
        ax.set_facecolor("#F8FAFF")
        ax.set_title(pl, fontsize=13, fontweight="bold", color="#1E2761", pad=7)

        pubs = [p for p in PUB_ORDER if p in ios_g["publisher_group"].values
                or p in aod_g["publisher_group"].values]
        n_pub = len(pubs)
        W_p = 0.32; INNER = 0.44; OUTER = 1.20

        xs_i = np.arange(n_pub, dtype=float) * OUTER
        xs_a = xs_i + INNER
        mid_x = (xs_i + xs_a) / 2

        for xi_i, xi_a, pub in zip(xs_i, xs_a, pubs):
            clr = PUB_COLORS.get(pub, "#9B59B6")

            def get_rv(df_g, _pub=pub, _pk=pk):
                sub = df_g[(df_g["publisher_group"]==_pub)&(df_g["period"]==_pk)]
                return (float(sub["rate"].values[0]) if len(sub) else 0,
                        int(sub["annual_n"].values[0]) if len(sub) else 0)

            ri, ni = get_rv(ios_g)
            ra, na = get_rv(aod_g)

            ax.bar(xi_i, max(ri,1.5), W_p, color=clr, alpha=IOS_ALPHA, edgecolor="white", lw=1.3)
            ax.bar(xi_a, max(ra,1.5), W_p, color=clr, alpha=AOS_ALPHA, hatch="///", edgecolor="white", lw=1.3)

            for xi, rate, n in [(xi_i,ri,ni),(xi_a,ra,na)]:
                if rate > 2:
                    ax.text(xi, rate+1.5, f"{rate:.0f}%", ha="center", va="bottom",
                            fontsize=10, fontweight="bold", color=clr)
                elif n == 0:
                    ax.text(xi, 5, "진입\n없음", ha="center", va="bottom",
                            fontsize=8, color="#aaa", style="italic")
                else:
                    ax.text(xi, 5, f"0%\n({n}탈락)", ha="center", va="bottom",
                            fontsize=8, color=clr, fontweight="bold")

            ax.text(xi_i, -1.5, f"iOS\n{ni}/yr", ha="center", va="top", fontsize=7, color="#2980B9", fontweight="bold")
            ax.text(xi_a, -1.5, f"AOS\n{na}/yr", ha="center", va="top", fontsize=7, color="#27AE60", fontweight="bold")

        ax.set_xticks(mid_x)
        ax.set_xticklabels(pubs, fontsize=11, fontweight="bold")
        ax.set_xlim(xs_i[0]-W_p/2-0.30, xs_a[-1]+W_p/2+0.30)
        ax.set_ylim(-4, 122)
        ax.set_ylabel("3개월 생존율 (%)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35)
        ax.spines[["top","right"]].set_visible(False)

    fig.legend(handles=[ios_patch, aod_patch], loc="upper right",
               bbox_to_anchor=(0.995,0.98), fontsize=10, framealpha=0.9)
    fig.text(0.5, 0.005,
             "매출 TOP 100 신규진입 기준  ·  서구권=미국·캐나다·영국·독일·프랑스·터키 등  ·  중화권=중국·대만·홍콩  ·  출처: Sensor Tower",
             ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 18: 퍼블리셔별 광고율 (생존 vs 미생존)
# ══════════════════════════════════════════════════════════════════════════
def chart_s18() -> bytes:
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.5), gridspec_kw={"wspace":0.38})
    fig.patch.set_facecolor("white")

    surv_patch = mpatches.Patch(facecolor="#2ECC71", alpha=0.85, label="3개월 생존")
    churn_patch = mpatches.Patch(facecolor="#E74C3C", alpha=0.55, hatch="///", label="3개월 미생존")

    for ax, (plat_name, plat_df) in zip(axes, [("iOS", ios_s1_ad), ("Android", aod_s1_ad)]):
        ax.set_facecolor("#F8FAFF")
        ax.set_title(plat_name, fontsize=13, fontweight="bold", color="#1E2761", pad=7)

        pubs = [p for p in PUB_ORDER if p in plat_df["publisher_group"].values]
        n_pub = len(pubs)
        W = 0.30; INNER = 0.38; OUTER = 1.10

        for pi, pub in enumerate(pubs):
            clr = PUB_COLORS.get(pub, "#9B59B6")
            xs = pi * OUTER
            xc = pi * OUTER + INNER

            sub = plat_df[(plat_df["publisher_group"]==pub)&(plat_df["total_ww"].notna())]
            surv = sub[sub["survived_3m"]==True]
            churn = sub[sub["survived_3m"]==False]
            r_surv  = (surv["has_ad"].mean()*100) if len(surv) else 0
            r_churn = (churn["has_ad"].mean()*100) if len(churn) else 0

            ax.bar(xs, r_surv, W, color="#2ECC71", alpha=0.85, edgecolor="white", lw=1.2)
            ax.bar(xc, r_churn, W, color="#E74C3C", alpha=0.55, hatch="///", edgecolor="white", lw=1.2)

            for x, r, c in [(xs,r_surv,"#2ECC71"),(xc,r_churn,"#E74C3C")]:
                if r > 1:
                    ax.text(x, r+2, f"{r:.0f}%", ha="center", va="bottom",
                            fontsize=10, fontweight="bold", color=c)

        xticks = [pi * OUTER + INNER/2 for pi in range(n_pub)]
        ax.set_xticks(xticks)
        ax.set_xticklabels(pubs, fontsize=11, fontweight="bold")
        ax.set_ylim(0, 120)
        ax.set_ylabel("광고집행율 (%)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35)
        ax.spines[["top","right"]].set_visible(False)

    fig.legend(handles=[surv_patch, churn_patch], loc="upper right",
               bbox_to_anchor=(0.995,0.98), fontsize=10, framealpha=0.9)
    fig.text(0.5, 0.005,
             "광고집행율 = Paid Install 비중 0% 초과  ·  출처: Sensor Tower",
             ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 19: 국가별 잔존율 M+1~M+12 (app_tags 기반)
# ══════════════════════════════════════════════════════════════════════════
def chart_s19() -> bytes:
    months_lbl = ["M+1","M+2","M+3","M+6","M+12"]
    ios_cols = ["d30_ret_ww","d60_ret_ww","d90_ret_ww","d180_ret_ww","d365_ret_ww"]
    aod_cols = ["d30_ret_ww","d60_ret_ww","d90_ret_ww","d180_ret_ww","d365_ret_ww"]

    # slide2_v2 has per-app retention with country
    # app_tags doesn't have country, so use slide2_v2 for country-level
    ret_df = ios_ret.copy()
    ret_cols = ["d30","d60","d90","d180","d365"]

    fig, axes = plt.subplots(1, 2, figsize=(13.5, 4.5), gridspec_kw={"wspace":0.28})
    fig.patch.set_facecolor("white")

    x = np.arange(len(months_lbl))
    markers = {"KR":"o","JP":"s","US":"^"}

    for ax, (pk, pl) in zip(axes, [("pre2025","2022–2024"),("post2025","2025~")]):
        ax.set_facecolor("#F8FAFF")
        ax.set_title(f"iOS — {pl}", fontsize=13, fontweight="bold", color="#1E2761", pad=8)
        handles = []

        for ctr in CTR_ORDER:
            sub = ret_df[(ret_df["country"]==ctr)&(ret_df["period"]==pk)]
            vals = []
            for col in ret_cols:
                v = sub[col].mean() if col in sub.columns and sub[col].notna().any() else np.nan
                vals.append(v)
            clr = CTR_COLORS[ctr]
            mk = markers.get(ctr, "o")
            valid_x = [xi for xi, v in zip(x, vals) if not np.isnan(v)]
            valid_v = [v for v in vals if not np.isnan(v)]

            if valid_v:
                ax.plot(valid_x, valid_v, color=clr, lw=2.5, marker=mk, markersize=8, zorder=3)
                ax.text(valid_x[-1]+0.15, valid_v[-1], f"{valid_v[-1]:.1f}%",
                        color=clr, fontsize=9, va="center", fontweight="bold")
                for xi, v in zip(valid_x, valid_v):
                    ax.text(xi, v+0.5, f"{v:.1f}", ha="center", va="bottom",
                            fontsize=8, color=clr)
            n = len(sub)
            handles.append(mlines.Line2D([],[],color=clr,lw=2.5,marker=mk,markersize=7,
                           label=f"{CTR_LABEL[ctr]}  (n={n})"))

        ax.set_xticks(x)
        ax.set_xticklabels(months_lbl, fontsize=11, fontweight="bold")
        ax.set_ylabel("잔존율 (%)", fontsize=10)
        ax.set_ylim(0, max(15, ax.get_ylim()[1]*1.15))
        ax.set_xlim(-0.4, len(months_lbl)-0.3)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(handles=handles, loc="upper right", fontsize=9, framealpha=0.9)

    fig.text(0.5, 0.005,
             "매출 TOP 100 신규진입 게임 기준  ·  WW 잔존율  ·  출처: Sensor Tower",
             ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 20: 퍼블리셔별 잔존율 M+1~M+12
# ══════════════════════════════════════════════════════════════════════════
def chart_s20() -> bytes:
    months_lbl = ["M+1","M+2","M+3","M+6","M+12"]
    pub_markers = {"한국":"D","일본":"^","서구권":"s","중화권":"o"}

    fig, axes = plt.subplots(1, 2, figsize=(13.5, 4.5), gridspec_kw={"wspace":0.28})
    fig.patch.set_facecolor("white")

    x = np.arange(len(months_lbl))

    ios_mcols = ["d30_ret_ww","d60_ret_ww","d90_ret_ww","d180_ret_ww","d365_ret_ww"]
    aod_mcols = ["d30_ret_ww","d60_ret_ww","d90_ret_ww","d180_ret_ww","d365_ret_ww"]

    for ax, (plat, tags_df, mcols, ls) in zip(axes, [
        ("iOS", ios_tags, ios_mcols, "-"),
        ("Android", aod_tags, aod_mcols, "--")
    ]):
        ax.set_facecolor("#F8FAFF")
        ax.set_title(f"{plat} — 퍼블리셔별 잔존율", fontsize=13, fontweight="bold", color="#1E2761", pad=8)
        handles = []

        for pub in PUB_ORDER:
            sub = tags_df[tags_df["publisher_group"]==pub]
            if len(sub) == 0:
                continue
            vals = []
            for col in mcols:
                v = sub[col].mean() if col in sub.columns and sub[col].notna().any() else np.nan
                vals.append(v)
            clr = PUB_COLORS[pub]
            mk = pub_markers.get(pub, "o")
            valid_x = [xi for xi, v in zip(x, vals) if not np.isnan(v)]
            valid_v = [v for v in vals if not np.isnan(v)]

            if valid_v:
                ax.plot(valid_x, valid_v, color=clr, lw=2.5, ls=ls, marker=mk, markersize=8, zorder=3)
                ax.text(valid_x[-1]+0.15, valid_v[-1], f"{valid_v[-1]:.1f}%",
                        color=clr, fontsize=9, va="center", fontweight="bold")
                for xi, v in zip(valid_x, valid_v):
                    if v > 0:
                        ax.text(xi, v+0.5, f"{v:.1f}", ha="center", va="bottom",
                                fontsize=8, color=clr)
            n = len(sub)
            handles.append(mlines.Line2D([],[],color=clr,lw=2.5,ls=ls,marker=mk,markersize=7,
                           label=f"{pub}  (n={n})"))

        ax.set_xticks(x)
        ax.set_xticklabels(months_lbl, fontsize=11, fontweight="bold")
        ax.set_ylabel("잔존율 (%)", fontsize=10)
        ax.set_ylim(0, 30)
        ax.set_xlim(-0.4, len(months_lbl)-0.3)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(handles=handles, loc="upper right", fontsize=9, framealpha=0.9)

    fig.add_artist(plt.Line2D([0.5,0.5],[0.06,0.96],
                              transform=fig.transFigure, color="#CBD5E1",lw=1.5,ls="--"))
    fig.text(0.5, 0.005,
             "WW 기준 잔존율  ·  서구권=US·EU 등  ·  중화권=CN·TW·HK  ·  출처: Sensor Tower",
             ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 21: 국가별 순위별 월평균 매출 (1,10,20,50,100위)
# ══════════════════════════════════════════════════════════════════════════
def chart_s21() -> bytes:
    ranks = [1, 10, 20, 50, 100]
    rank_labels = ["1위","10위","20위","50위","100위"]

    fig, axes = plt.subplots(1, 2, figsize=(13.5, 4.5), sharey=False,
                             gridspec_kw={"wspace":0.28})
    fig.patch.set_facecolor("white")

    # iOS data
    df_ios = ios_rev[ios_rev["country"].isin(CTR_ORDER)].copy()
    df_ios["revenue_m_usd"] = pd.to_numeric(df_ios["revenue_m_usd"], errors="coerce")

    for ax, (pk, pl) in zip(axes, [("pre2025","2025 이전 (2022–2024)"),("post2025","2025 이후 (2025~)")]):
        ax.set_facecolor("#F8FAFF")
        sub = df_ios[df_ios["period"]==pk] if "period" in df_ios.columns else df_ios
        x = np.arange(len(ranks))
        w = 0.25

        for k, ctr in enumerate(CTR_ORDER):
            vals = []
            for rank in ranks:
                r = sub[(sub["country"]==ctr)&(sub["rank"]==rank)]
                vals.append(r["revenue_m_usd"].mean() if len(r) else 0)
            bars = ax.bar(x + k*w, vals, w, label=CTR_LABEL[ctr],
                          color=CTR_COLORS[ctr], alpha=0.88, edgecolor="white", lw=1)
            for bar, v in zip(bars, vals):
                if v > 0.005:
                    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.01,
                            f"${v:.2f}M", ha="center", va="bottom",
                            fontsize=7, fontweight="bold", color=CTR_COLORS[ctr])

        ax.set_title(f"iOS — {pl}", fontsize=12, fontweight="bold", color="#1E2761", pad=8)
        ax.set_xticks(x + w)
        ax.set_xticklabels(rank_labels, fontsize=10.5)
        ax.set_ylabel("월 평균 매출 (백만 USD)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(fontsize=9, framealpha=0.9)

    fig.add_artist(plt.Line2D([0.5,0.5],[0.06,0.96],
                              transform=fig.transFigure, color="#CBD5E1",lw=1.5,ls="--"))
    fig.text(0.5, 0.005,
             "iOS App Store TOP 100  ·  KR=한국 / JP=일본 / US=미국  ·  출처: Sensor Tower",
             ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 22: 퍼블리셔별 순위별 월평균 매출
# ══════════════════════════════════════════════════════════════════════════
def chart_s22() -> bytes:
    ranks = [1, 10, 20, 50, 100]
    rank_labels = ["1위","10위","20위","50위","100위"]

    # slide3_v2 has publisher_group
    df = ios_rev[ios_rev["publisher_group"].isin(PUB_ORDER)].copy()
    df["revenue_m_usd"] = pd.to_numeric(df["revenue_m_usd"], errors="coerce")

    fig, axes = plt.subplots(1, 2, figsize=(13.5, 4.3), sharey=False,
                             gridspec_kw={"wspace":0.28})
    fig.patch.set_facecolor("white")

    for ax, (pk, pl) in zip(axes, [("pre2025","2025 이전 (2022–2024)"),("post2025","2025 이후 (2025~)")]):
        ax.set_facecolor("#F8FAFF")
        sub = df[df["period"]==pk] if "period" in df.columns else df
        x = np.arange(len(ranks))
        w = 0.2

        for k, pub in enumerate(PUB_ORDER):
            vals = []
            for rank in ranks:
                r = sub[(sub["publisher_group"]==pub)&(sub["rank"]==rank)]
                vals.append(r["revenue_m_usd"].mean() if len(r) else 0)
            bars = ax.bar(x + k*w, vals, w, label=pub,
                          color=PUB_COLORS[pub], alpha=0.88, edgecolor="white", lw=1)
            for bar, v in zip(bars, vals):
                if v > 0.005:
                    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.01,
                            f"${v:.2f}M", ha="center", va="bottom",
                            fontsize=7, fontweight="bold", color=PUB_COLORS[pub])

        ax.set_title(f"iOS — {pl}", fontsize=12, fontweight="bold", color="#1E2761", pad=8)
        ax.set_xticks(x + w*1.5)
        ax.set_xticklabels(rank_labels, fontsize=10.5)
        ax.set_ylabel("월 평균 매출 (백만 USD)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(fontsize=9, framealpha=0.9, ncol=2)

    fig.add_artist(plt.Line2D([0.5,0.5],[0.06,0.96],
                              transform=fig.transFigure, color="#CBD5E1",lw=1.5,ls="--"))
    fig.text(0.5, 0.005,
             "iOS App Store TOP 100  ·  서구권=미국·캐나다·영국·독일·프랑스·터키 등  ·  출처: Sensor Tower",
             ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════════════════
SLIDE_NOTES = {
    "s15": "해석: 2025년 이후 한국·일본 모두 생존율이 하락하며, 연평균 신규 진입 수도 감소 추세. | 인사이트: 시장 성숙기 진입으로 신규 게임의 정착 난이도가 높아지고 있음.",
    "s16": "해석: 생존 게임은 미생존 게임 대비 광고집행율이 높은 경향. | 인사이트: 적극적 UA 투자가 초기 생존에 유의미한 상관관계를 보이며, 특히 일본 시장에서 격차가 뚜렷.",
    "s17": "해석: 중화권·서구권 퍼블리셔가 전반적으로 높은 생존율을 기록, 한국은 플랫폼별 편차가 큼. | 인사이트: 글로벌 퍼블리셔의 다시장 런칭 경험이 생존율에 긍정적 영향.",
    "s18": "해석: 생존 게임의 광고집행율이 퍼블리셔 유형과 무관하게 일관되게 높음. | 인사이트: Paid UA는 퍼블리셔 규모와 관계없이 초기 생존의 필수 요소로 확인됨.",
    "s19": "해석: M+3 이후 잔존율 급감, 한국 시장이 전반적으로 높은 잔존율 유지. | 인사이트: 한국 게임의 강한 커뮤니티/라이브서비스 문화가 장기 잔존에 기여.",
    "s20": "해석: iOS에서 일본 퍼블리셔 잔존율이 가장 높고, Android에서 중화권이 빠르게 하락. | 인사이트: 플랫폼별 유저 충성도 차이가 퍼블리셔 전략에 영향을 미침.",
    "s21": "해석: 1위 매출은 국가 간 수배 차이, 50위 이하에서는 격차가 크게 줄어듦. | 인사이트: TOP 10 진입 여부가 수익성의 결정적 분기점이며, 미국 시장의 상위권 매출이 압도적.",
    "s22": "해석: 중화권 퍼블리셔가 상위권에서 높은 매출을 기록하나 하위권에서는 서구권과 유사. | 인사이트: 중화권의 매출 집중도가 높아 상위 랭킹 확보가 사업 성패의 핵심.",
}


def main():
    print("▶ 차트 재생성 중 (슬라이드 15~22 전면 재작성)...")

    imgs = {}
    imgs["s15"] = chart_s15(); print("  ✔ Slide 15 — 국가별 생존율")
    imgs["s16"] = chart_s16(); print("  ✔ Slide 16 — 국가별 광고율")
    imgs["s17"] = chart_s17(); print("  ✔ Slide 17 — 퍼블리셔별 생존율")
    imgs["s18"] = chart_s18(); print("  ✔ Slide 18 — 퍼블리셔별 광고율")
    imgs["s19"] = chart_s19(); print("  ✔ Slide 19 — 국가별 잔존율")
    imgs["s20"] = chart_s20(); print("  ✔ Slide 20 — 퍼블리셔별 잔존율")
    imgs["s21"] = chart_s21(); print("  ✔ Slide 21 — 국가별 매출")
    imgs["s22"] = chart_s22(); print("  ✔ Slide 22 — 퍼블리셔별 매출")

    prs = Presentation(str(PPTX))
    n_slides = len(prs.slides)
    print(f"\n슬라이드 수: {n_slides}장")

    # Slide 15~22 (0-indexed: 14~21)
    slide_map = {
        "s15": ("생존율 — 국가별",               14),
        "s16": ("광고율 — 국가별",               15),
        "s17": ("생존율 — 퍼블리셔",             16),
        "s18": ("광고율 — 퍼블리셔",             17),
        "s19": ("잔존율 — 국가별",               18),
        "s20": ("잔존율 — 퍼블리셔",             19),
        "s21": ("매출 — 국가별",                 20),
        "s22": ("매출 — 퍼블리셔",               21),
    }

    # slide title updates
    slide_titles = {
        "s15": "신규 게임 3개월 생존율 — 국가별 (KR·JP·US)",
        "s16": "광고 집행율 — 국가별 (생존 vs 미생존)",
        "s17": "신규 게임 3개월 생존율 — 퍼블리셔별",
        "s18": "광고 집행율 — 퍼블리셔별 (생존 vs 미생존)",
        "s19": "월별 잔존율 M+1~M+12 — 국가별",
        "s20": "월별 잔존율 M+1~M+12 — 퍼블리셔별",
        "s21": "순위별 월평균 매출 — 국가별 (1·10·20·50·100위)",
        "s22": "순위별 월평균 매출 — 퍼블리셔별",
    }

    for key, (kw, fb_idx) in slide_map.items():
        idx, slide = find_slide(prs, kw, fallback_idx=fb_idx)
        if slide is None:
            print(f"  ⚠ [{key}] Slide {fb_idx+1} 찾기 실패")
            continue

        # Update title
        title_text = slide_titles.get(key, "")
        if title_text:
            for sh in slide.shapes:
                if sh.has_text_frame and sh.top / 914400 < 1.5:
                    for para in sh.text_frame.paragraphs:
                        if para.text.strip():
                            for run in para.runs:
                                run.text = title_text
                            break
                    break

        note = SLIDE_NOTES.get(key, "")
        maximize_chart(slide, imgs[key], note_text=note)
        print(f"  ✔ [{key}] Slide {fb_idx+1} 업데이트 완료")

    prs.save(str(PPTX))
    print(f"\n▶ 저장 완료: {PPTX.name}")


if __name__ == "__main__":
    main()
