"""임원 보고용 Executive Summary v2 — 깔끔한 카드형 디자인"""
import sys
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)   # 와이드 16:9
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

NAVY = RGBColor(0x1E, 0x27, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xAA, 0xAA, 0xAA)
TEAL = RGBColor(0x0D, 0x94, 0x88)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
BLUE = RGBColor(0x34, 0x98, 0xDB)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFF)


def add_header(slide, title, subtitle=""):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12), Inches(0.65))
    p = txb.text_frame.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.78), Inches(12), Inches(0.35))
        p2 = txb2.text_frame.paragraphs[0]; p2.text = subtitle
        r2 = p2.runs[0]; r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)


def add_footer(slide):
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
    p = txb.text_frame.paragraphs[0]
    p.text = "Source: Sensor Tower API  |  매출 TOP 100 기준  |  2022-2026.Q1  |  랭킹 기준일: 매월 1일"
    p.runs[0].font.size = Pt(9); p.runs[0].font.color.rgb = LGRAY


def add_card(slide, left, top, w, h, fill=BG_LIGHT, border=RGBColor(0xDD, 0xDD, 0xDD)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = fill; card.line.color.rgb = border; card.line.width = Pt(1)
    return card


def add_text(slide, left, top, w, h, text, size=12, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return txb


def add_multiline(slide, left, top, w, h, lines, size=11, color=GRAY):
    """여러 줄 텍스트 (줄마다 별도 paragraph)"""
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(2)
        for r in p.runs:
            r.font.size = Pt(size); r.font.color.rgb = color


def add_badge(slide, left, top, w, h, text, fill_color):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    badge.fill.solid(); badge.fill.fore_color.rgb = fill_color; badge.line.fill.background()
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE


# ══════════════════════════════════════════════════════════════
# Slide 1: 표지
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

add_text(s, 1.5, 2.0, 10, 1.2, "모바일 게임 시장 변화 분석", 40, True, WHITE, PP_ALIGN.CENTER)
add_text(s, 1.5, 3.2, 10, 0.6, "Executive Summary", 24, False, TEAL, PP_ALIGN.CENTER)
add_text(s, 1.5, 4.2, 10, 0.5, "한국  |  일본  |  미국  |  2022 - 2026.Q1  |  매출 TOP 100 기준", 16, False, LGRAY, PP_ALIGN.CENTER)

# 하단 3개 키워드 배지
for i, (txt, clr) in enumerate([("iOS + AOS 분리 분석", TEAL), ("25년 전후 비교", ORANGE), ("Sensor Tower API", BLUE)]):
    add_badge(s, 2.5 + i * 3.0, 5.5, 2.5, 0.5, txt, clr)

add_text(s, 1.5, 6.5, 10, 0.3, "2026.04", 12, False, RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)
print("  S1 표지")


# ══════════════════════════════════════════════════════════════
# Slide 2: 핵심 메시지 3가지
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_header(s, "핵심 메시지 3가지", "2025년 이후 모바일 게임 시장의 구조적 변화")

msgs = [
    ("01", RED, "생존 난이도 급상승",
     ["2025년 이후 전 시장 생존율 하락",
      "한국 AOS 38.5% -> 22% (-16.5%p) 가장 큰 낙폭",
      "미국 AOS만 유일하게 +2.3%p 상승",
      "M+3 잔존율도 전 시장 20~30%대로 하락"]),
    ("02", ORANGE, "퍼블리셔 간 격차 확대",
     ["기타(Supercell, King 등) 생존율 42%로 최고",
      "한국 퍼블리셔 AOS 19%로 최하위",
      "일본 퍼블리셔 광고 없이도 높은 장기 잔존",
      "한국은 광고 의존도만 상승(41.7%)"]),
    ("03", BLUE, "미국 시장 매출 집중도 심화",
     ["iOS 1위 월매출: 미국 $52M vs 한국 $6M",
      "AOS 1위 월매출: 미국 $33M vs 한국 $21M",
      "TOP 10 진입이 수익성의 결정적 분기점",
      "중화권 AOS 상위권 매출 성장세"]),
]

for i, (num, color, title, bullets) in enumerate(msgs):
    x = 0.4 + i * 4.2
    add_card(s, x, 1.5, 3.9, 5.2)

    add_badge(s, x + 0.15, 1.7, 0.7, 0.5, num, color)
    add_text(s, x + 1.0, 1.7, 2.7, 0.5, title, 16, True, NAVY)

    add_multiline(s, x + 0.3, 2.5, 3.3, 4.0, bullets, 12, GRAY)

add_footer(s)
print("  S2 핵심 메시지")


# ══════════════════════════════════════════════════════════════
# Slide 3: 국가별 종합 스코어카드
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_header(s, "국가별 종합 스코어카드", "25년 이후 기준  |  iOS + AOS  |  매출 TOP 100 신규진입")

countries = [
    ("KR 한국", RED,
     [("생존율", "iOS 24%  /  AOS 22%", RED),
      ("M+3 잔존", "iOS 24%  /  AOS 22%", RED),
      ("1위 매출", "iOS $5.8M  /  AOS $21M", GRAY),
      ("종합", "경쟁 최고, 생존 최저", RED)]),
    ("JP 일본", GREEN,
     [("생존율", "iOS 25%  /  AOS 25%", ORANGE),
      ("M+3 잔존", "iOS 25%  /  AOS 25%", ORANGE),
      ("1위 매출", "iOS $19.6M  /  AOS $12.7M", ORANGE),
      ("종합", "장기 잔존 강자였으나 하락 추세", ORANGE)]),
    ("US 미국", BLUE,
     [("생존율", "iOS 33%  /  AOS 41%", GREEN),
      ("M+3 잔존", "iOS 33%  /  AOS 41%", GREEN),
      ("1위 매출", "iOS $52.4M  /  AOS $33M", GREEN),
      ("종합", "매출 최대, AOS 생존 유일 상승", GREEN)]),
]

for i, (name, color, metrics) in enumerate(countries):
    x = 0.4 + i * 4.2
    add_card(s, x, 1.5, 3.9, 5.2)

    # 국가 헤더
    hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.9), Inches(0.7))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
    tf = hdr.text_frame; p = tf.paragraphs[0]; p.text = name; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE

    for mi, (label, value, val_color) in enumerate(metrics):
        y = 2.5 + mi * 1.1
        add_text(s, x + 0.3, y, 1.2, 0.3, label, 11, True, LGRAY)
        is_last = mi == 3
        add_text(s, x + 0.3, y + 0.35, 3.3, 0.4, value, 14 if is_last else 13, is_last, val_color if is_last else NAVY)

add_footer(s)
print("  S3 국가별 스코어카드")


# ══════════════════════════════════════════════════════════════
# Slide 4: 퍼블리셔별 종합 스코어카드
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_header(s, "퍼블리셔별 종합 스코어카드", "25년 이후 기준  |  iOS + AOS 종합  |  매출 TOP 100 신규진입")

publishers = [
    ("한국", RED,
     ["생존율: 19~24%", "M+3 잔존: 19~21%", "광고 의존 급증"],
     "최하위"),
    ("일본", GREEN,
     ["생존율: 22~23%", "M+3 잔존: 23%", "광고 없이도 생존"],
     "안정적 잔존"),
    ("북미", BLUE,
     ["생존율: 22~30%", "M+3 잔존: 22~30%", "매출 규모 1위"],
     "매출 최강"),
    ("중화권", ORANGE,
     ["생존율: 25~27%", "M+3 잔존: 25~27%", "물량 공세 최다"],
     "진입 최다"),
    ("기타", PURPLE,
     ["생존율: 38~42%", "M+3 잔존: 38~42%", "Supercell, King"],
     "소수 정예 최강"),
]

for i, (name, color, bullets, tag) in enumerate(publishers):
    x = 0.3 + i * 2.55
    add_card(s, x, 1.5, 2.35, 5.2)

    hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(2.35), Inches(0.65))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
    tf = hdr.text_frame; p = tf.paragraphs[0]; p.text = name; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE

    add_multiline(s, x + 0.2, 2.5, 2.0, 2.5, bullets, 12, GRAY)

    # 태그 배지
    tag_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.3), Inches(5.5), Inches(1.75), Inches(0.55))
    tag_bg.fill.solid(); tag_bg.fill.fore_color.rgb = color; tag_bg.fill.fore_color.brightness = 0.7
    tag_bg.line.fill.background()
    tf2 = tag_bg.text_frame; p2 = tf2.paragraphs[0]; p2.text = tag; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.runs[0]; r2.font.size = Pt(12); r2.font.bold = True; r2.font.color.rgb = color

add_footer(s)
print("  S4 퍼블리셔 스코어카드")


# ══════════════════════════════════════════════════════════════
# Slide 5: NHN 시사점 & 대응 방향
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_header(s, "NHN 시사점 & 대응 방향", "시장 분석 기반 전략적 시사점")

items = [
    ("위기", RED, RGBColor(0xFD, 0xE8, 0xE8),
     ["한국 퍼블리셔 생존율 최하위 (AOS 19%)",
      "M+3 이후 차트 잔존율 급감",
      "광고 의존도만 상승, 장기 경쟁력 약화"],
     ["신규 론칭 시 초기 3개월 생존 전략 필수",
      "M+3 이후 잔존 개선 로드맵 선행 설계",
      "게임 품질과 IP 경쟁력 확보 우선"]),
    ("기회", ORANGE, RGBColor(0xFF, 0xF8, 0xE1),
     ["미국 AOS 유일한 생존율 상승 시장",
      "미국 1위 매출 $33M으로 최대 규모",
      "북미 퍼블리셔 안정적 생존력"],
     ["글로벌 타겟 시 미국 AOS 우선 진입 검토",
      "Action, Strategy 장르 중심 포트폴리오",
      "현지화 및 글로벌 퍼블리싱 역량 강화"]),
    ("과제", GREEN, RGBColor(0xE8, 0xF5, 0xE9),
     ["일본/기타(유럽) 퍼블리셔 장기 잔존 우위",
      "한국 M+3 잔존율 19~21%로 최하위",
      "광고만으로는 장기 생존 불가"],
     ["IP 기반 라이브서비스 강화",
      "커뮤니티 운영 고도화",
      "Whale 유저 리텐션 집중 전략"]),
]

for i, (badge_text, badge_color, bg_color, diagnosis, actions) in enumerate(items):
    y = 1.5 + i * 1.85

    add_badge(s, 0.5, y, 1.0, 0.5, badge_text, badge_color)

    # 현황 진단
    add_card(s, 1.7, y, 5.0, 1.55, bg_color)
    add_text(s, 1.9, y + 0.05, 2.0, 0.3, "현황 진단", 10, True, LGRAY)
    add_multiline(s, 1.9, y + 0.35, 4.6, 1.1, diagnosis, 11, NAVY)

    # 대응 방향
    add_card(s, 7.0, y, 5.8, 1.55, RGBColor(0xF0, 0xFF, 0xF0))
    add_text(s, 7.2, y + 0.05, 2.0, 0.3, "대응 방향", 10, True, LGRAY)
    add_multiline(s, 7.2, y + 0.35, 5.4, 1.1, actions, 11, GREEN)

add_footer(s)
print("  S5 NHN 시사점")


# ══════════════════════════════════════════════════════════════
# Slide 6: 결론
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_header(s, "결론", "")

conclusions = [
    ("01", "시장 성숙기 본격 진입",
     "2025년 이후 전 시장에서 신규 게임의 3개월 생존율과 장기 잔존율이 급락하여 시장 진입 난이도가 구조적으로 높아졌다."),
    ("02", "Paid UA만으로는 생존 보장 불가",
     "한국/일본에서 광고집행율과 생존의 상관이 있으나, 미국과 유럽 강자들은 게임 품질과 IP 경쟁력으로 생존하고 있다."),
    ("03", "미국 AOS가 유일한 성장 시장",
     "생존율 유일 상승, 1위 매출 $33M 최대 규모. 글로벌 매출 극대화를 위한 최우선 타겟 시장이다."),
    ("04", "한국 퍼블리셔 장기 잔존 전략 시급",
     "M+3 잔존율 19~21%로 전 퍼블리셔 중 최하위. IP 기반 라이브서비스와 커뮤니티 운영 강화가 핵심 과제이다."),
]

for i, (num, title, desc) in enumerate(conclusions):
    y = 1.6 + i * 1.35

    add_badge(s, 0.5, y, 0.7, 0.55, num, NAVY)
    add_text(s, 1.4, y, 4.0, 0.55, title, 16, True, NAVY)
    add_text(s, 1.4, y + 0.55, 11.0, 0.6, desc, 13, False, GRAY)

    if i < 3:
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y + 1.25), Inches(12.3), Inches(0.015))
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE); line.line.fill.background()

add_footer(s)
print("  S6 결론")


# 저장
out = "C:/Users/NHN/Documents/sensortower_api/output/pptx/MobileGame_Executive_Summary_v2.pptx"
prs.save(out)
print(f"\nDone: {out}")
