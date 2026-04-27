"""
PPTX 구조 수술:
- rId33가 slide25.xml(국가별 광고집행율 중복)을 가리키는 문제
- rId33 → slide26.xml(결론)로 변경
- slide26.xml, slide26.xml.rels, [Content_Types] 추가
백업: MobileGame_Market_Analysis_2022-2026_v2_backup.pptx 사용
"""
import sys, zipfile, re, shutil
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

from pathlib import Path

BASE     = Path("C:/Users/NHN/Documents/sensortower_api")
BACKUP   = BASE / "MobileGame_Market_Analysis_2022-2026_v2_backup.pptx"
UNPACKED = BASE / "pptx_unpacked3"
PPTX_OUT = BASE / "MobileGame_Market_Analysis_2022-2026_v2.pptx"
TMP      = BASE / "_rebuild2_tmp.pptx"


def main():
    # 결론 슬라이드 원본 읽기 (slide14.xml = 결론, slide17.xml = 월매출 차트)
    slide_xml  = (UNPACKED / "ppt/slides/slide14.xml").read_bytes()
    slide_rels = (UNPACKED / "ppt/slides/_rels/slide14.xml.rels").read_bytes()
    print("원본 결론 슬라이드 rels:")
    print(slide_rels.decode("utf-8"))

    # 백업에서 읽어서 수정된 버전으로 저장
    with zipfile.ZipFile(BACKUP, "r") as src, \
         zipfile.ZipFile(TMP, "w", zipfile.ZIP_DEFLATED) as dst:

        written = set()
        for item in src.infolist():
            name = item.filename
            data = src.read(name)

            # 중복 slide25.xml — 첫 번째만 통과 (국가별 광고집행율)
            if name in written:
                print(f"  [스킵 중복] {name}")
                continue

            # presentation.xml.rels: rId33의 Target을 slide26.xml로 변경
            if name == "ppt/_rels/presentation.xml.rels":
                text = data.decode("utf-8")
                # rId33 → slide26.xml
                text = re.sub(
                    r'(Id="rId33"[^>]+Target=")slides/slide25\.xml(")',
                    r'\1slides/slide26.xml\2',
                    text
                )
                print(f"  [수정] {name}")
                print(f"    rId33 → slides/slide26.xml")
                data = text.encode("utf-8")

            # [Content_Types].xml: slide26.xml Override 추가
            if name == "[Content_Types].xml":
                text = data.decode("utf-8")
                if "slide26.xml" not in text:
                    insert_before = "</Types>"
                    new_entry = ('  <Override PartName="/ppt/slides/slide26.xml" '
                                 'ContentType="application/vnd.openxmlformats-officedocument'
                                 '.presentationml.slide+xml"/>\n')
                    text = text.replace(insert_before, new_entry + insert_before)
                    print(f"  [수정] {name}: slide26.xml 추가")
                data = text.encode("utf-8")

            dst.writestr(item, data)
            written.add(name)

        # slide26.xml (결론) 추가
        dst.writestr("ppt/slides/slide26.xml", slide_xml)
        dst.writestr("ppt/slides/_rels/slide26.xml.rels", slide_rels)
        print("  [추가] ppt/slides/slide26.xml (결론)")
        print("  [추가] ppt/slides/_rels/slide26.xml.rels")

    # 최종 파일로 이동
    shutil.move(str(TMP), str(PPTX_OUT))
    print(f"\n▶ 저장 완료: {PPTX_OUT.name}")

    # 검증
    from pptx import Presentation
    prs = Presentation(str(PPTX_OUT))
    print(f"  총 {len(prs.slides)}개 슬라이드")
    for i, s in enumerate(prs.slides):
        texts = [sh.text_frame.text.strip()[:55] for sh in s.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        marker = " ◀" if i >= 22 else ""
        print(f"  Slide {i+1}: {texts[0] if texts else '(빈 슬라이드)'}{marker}")


if __name__ == "__main__":
    main()
