"""
v3로 저장:
- slide25.xml의 slideLayout 참조를 실제 존재하는 레이아웃으로 수정
- slide26.xml.rels의 notesSlide 참조 제거 (없는 파일 참조)
- MobileGame_Market_Analysis_2022-2026_v3.pptx로 저장
"""
import sys, zipfile, re, shutil
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")
from pathlib import Path

BASE    = Path("C:/Users/NHN/Documents/sensortower_api")
SRC     = BASE / "MobileGame_Market_Analysis_2022-2026_v2.pptx"
OUT     = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"
TMP     = BASE / "_v3_tmp.pptx"


def main():
    # 사용 가능한 slideLayout 목록 파악
    with zipfile.ZipFile(SRC, "r") as zf:
        layouts = sorted(set([n for n in zf.namelist()
                               if re.match(r"ppt/slideLayouts/slideLayout\d+\.xml$", n)]))
        print("사용 가능한 레이아웃:", layouts)
        # 빈 레이아웃으로 쓸 첫 번째 레이아웃
        blank_layout = layouts[0].split("ppt/slideLayouts/")[1]  # e.g. slideLayout1.xml
        print(f"사용할 레이아웃: {blank_layout}")

        # slide25.xml.rels 현재 내용
        rels25 = zf.read("ppt/slides/_rels/slide25.xml.rels").decode("utf-8")
        print(f"\n현재 slide25.xml.rels:\n{rels25}")

        # slide26.xml.rels 현재 내용
        rels26 = zf.read("ppt/slides/_rels/slide26.xml.rels").decode("utf-8")
        print(f"\n현재 slide26.xml.rels:\n{rels26}")

        # 존재하는 notesSlide 목록
        notes = [n for n in zf.namelist() if "notesSlide" in n]
        print(f"\n존재하는 notesSlide: {notes}")

    with zipfile.ZipFile(SRC, "r") as src, \
         zipfile.ZipFile(TMP, "w", zipfile.ZIP_DEFLATED) as dst:

        written = set()
        for item in src.infolist():
            name = item.filename
            if name in written:
                continue
            data = src.read(name)

            # slide25.xml.rels: slideLayout7 → 실제 존재하는 레이아웃으로 교체
            if name == "ppt/slides/_rels/slide25.xml.rels":
                text = data.decode("utf-8")
                # slideLayout 참조 교체
                text = re.sub(
                    r'Target="\.\./slideLayouts/slideLayout\d+\.xml"',
                    f'Target="../slideLayouts/{blank_layout}"',
                    text
                )
                # 없는 notesSlide 참조 제거
                text = re.sub(
                    r'<Relationship[^>]+notesSlide[^>]*/>\s*', '', text)
                print(f"  [수정] {name}")
                print(f"    → {text.strip()}")
                data = text.encode("utf-8")

            # slide26.xml.rels: notesSlide14 참조 제거 (없는 파일)
            if name == "ppt/slides/_rels/slide26.xml.rels":
                text = data.decode("utf-8")
                # notesSlide14.xml이 없으면 참조 제거
                if "notesSlide14.xml" in text and not any("notesSlide14" in n
                                                           for n in notes):
                    text = re.sub(
                        r'<Relationship[^>]+notesSlides/notesSlide14\.xml[^>]*/>\s*',
                        '', text)
                    print(f"  [수정] {name}: notesSlide14 참조 제거")
                data = text.encode("utf-8")

            dst.writestr(item, data)
            written.add(name)

    shutil.move(str(TMP), str(OUT))
    print(f"\n▶ 저장 완료: {OUT.name}")

    # 검증
    from pptx import Presentation
    prs = Presentation(str(OUT))
    print(f"  총 {len(prs.slides)}개 슬라이드")
    for i, s in enumerate(prs.slides):
        texts = [sh.text_frame.text.strip()[:55] for sh in s.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        marker = " ◀" if i >= 22 else ""
        print(f"  Slide {i+1}: {texts[0] if texts else '(빈 슬라이드)'}{marker}")


if __name__ == "__main__":
    main()
