"""
v2.pptx 추가 수정:
  Slide 18 - CN 제거(KR/US/JP만), KR/JP 겹침 선 스타일 완전 구분
  Slide 23 뒤에 새 슬라이드 삽입 - 광고집행율 국가별(KR/JP/CN/서구권)
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

rcParams["font.family"]     = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE    = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX_IN = BASE / "MobileGame_Market_Analysis_2022-2026_v2.pptx"

NAVY = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x0D, 0x94, 0x88)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

COUNTRY_COLORS   = {"KR": "#E74C3C", "US": "#3498DB", "JP": "#2ECC71", "CN": "#F39C12"}
PUBLISHER_COLORS = {"한국": "#E74C3C", "일본": "#2ECC71", "중화권": "#F39C12",
                    "서구권": "#3498DB", "기타": "#95A5A6"}
GROUP_TO_CODE    = {"한국": "KR", "일본": "JP", "중화권": "CN", "서구권": "US"}


# ── 헬퍼 ──────────────────────────────────────────────────────────────────────

def fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def replace_chart_image(slide, new_img_bytes: bytes):
    """슬라이드에서 첫 번째 그림 도형을 찾아 이미지를 교체 (위치/크기 유지)."""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            left   = shape.left
            top    = shape.top
            width  = shape.width
            height = shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(
                io.BytesIO(new_img_bytes), left, top, width, height)
            return True
    return False


def add_note_box(slide, text: str, top_inches: float = 6.55,
                 bg_color=None, font_size=8.5, bold=False,
                 left=0.25, width=9.2, height=0.5):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top_inches), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = GRAY
    if bg_color:
        tb.fill.solid()
        r, g, b = [int(bg_color[i:i+2], 16) for i in (1, 3, 5)]
        tb.fill.fore_color.rgb = RGBColor(r, g, b)
    return tb


def add_title_box(slide, text: str):
    txBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.1), Inches(9.1), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(20)
    run.font.bold    = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)


def insert_slide_at(prs: Presentation, slide, position: int):
    xml_slides = prs.slides._sldIdLst
    last = xml_slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(position, last)


# ── 차트: Slide 18 (국가별 잔존율, KR/JP 겹침 완전 구분) ─────────────────────

def make_retention_chart_slide18(df2: pd.DataFrame) -> bytes:
    """
    KR / US / JP 3국만.
    - 각 시리즈: 다른 선 스타일 + 다른 마커
    - KR·JP 값 동일 시: JP를 위에 그려서 dash-dot 패턴이 보이도록 zorder 조정
    - 끝점 레이블: 상하로 오프셋해서 겹침 방지
    """
    groups = ["KR", "US", "JP"]
    colors = {"KR": "#E74C3C", "US": "#3498DB", "JP": "#2ECC71"}
    # JP를 마지막에 그리되 zorder 가장 높게 → dash-dot이 KR solid 위에 보임
    line_styles = {"KR": "-",   "US": "--",  "JP": "-."}
    markers     = {"KR": "o",   "US": "s",   "JP": "^"}
    zorders     = {"KR": 8,     "US": 9,     "JP": 11}   # JP > KR so JP dash-dot visible
    lw          = {"KR": 2.8,   "US": 2.2,   "JP": 2.0}

    # 끝점 레이블 수직 오프셋 (포인트): KR 아래, JP 위
    label_va    = {"KR": "top",    "US": "center", "JP": "bottom"}
    label_off   = {"KR": -0.8,     "US": 0,        "JP": +0.8}

    days_cols   = ["d1", "d7", "d14", "d30"]
    days_labels = ["D1", "D7", "D14", "D30"]

    pre  = df2[df2["period"] == "pre2025"]
    post = df2[df2["period"] == "post2025"]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.8))
    fig.patch.set_alpha(0)

    for ax, period_df, period_label in [
            (axes[0], pre,  "2025 이전  (2022–2024)"),
            (axes[1], post, "2025 이후  (2025–2026)")]:

        for grp in groups:
            sub = period_df[period_df["country"] == grp]
            if sub.empty:
                continue
            vals = [sub[c].mean() for c in days_cols]
            ax.plot(days_labels, vals,
                    marker=markers[grp], markersize=7,
                    linewidth=lw[grp],
                    linestyle=line_styles[grp],
                    color=colors[grp],
                    label=grp,
                    zorder=zorders[grp])
            # 끝점 레이블 — 상하 오프셋으로 겹침 방지
            ax.annotate(
                f" {grp}  {vals[-1]:.1f}%",
                xy=(days_labels[-1], vals[-1]),
                xytext=(5, label_off[grp]),
                textcoords="offset points",
                fontsize=7.5, color=colors[grp],
                va=label_va[grp])

        ax.set_title(period_label, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_ylabel("잔존율 (%)", fontsize=9)
        ax.set_ylim(0, 65)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.set_facecolor("#F8FAFF")
        ax.spines[["top", "right"]].set_visible(False)
        ax.legend(loc="upper right", fontsize=8.5, framealpha=0.75,
                  handlelength=2.5)   # 핸들 길이 늘려 선 스타일 잘 보이게

        # post2025: KR·JP 동일값 안내
        if period_label.startswith("2025 이후"):
            ax.text(0.5, 0.97,
                    "※ KR·JP 잔존율 동일 (KR 실선 ●, JP 점선 ▲으로 구분)",
                    transform=ax.transAxes, ha="center", va="top",
                    fontsize=7.8, color="#444",
                    bbox=dict(boxstyle="round,pad=0.3",
                              facecolor="#FFF8DC", edgecolor="#CCA", alpha=0.9))

    fig.tight_layout()
    return fig_to_bytes(fig)


# ── 차트: 새 슬라이드 — 광고집행율 국가별 ─────────────────────────────────────

def make_ad_rate_country_chart(tags: pd.DataFrame) -> bytes:
    """
    publisher_group(한국/일본/중화권/서구권) 기준 광고집행율을
    국가코드(KR/JP/CN/US)로 표기한 가로 막대 차트.
    WW(전세계) vs US(미국 앱스토어) 나란히 비교.
    """
    groups = ["한국", "일본", "중화권", "서구권"]
    codes  = ["KR",   "JP",   "CN",    "US"]
    colors_bar = [COUNTRY_COLORS[c] for c in codes]

    agg = tags.groupby("publisher_group")[
        ["paid_display_pct_ww", "paid_search_pct_ww",
         "paid_display_pct_us", "paid_search_pct_us"]].mean().round(1)

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))
    fig.patch.set_alpha(0)

    titles = ["전세계 (WW) 기준", "미국 앱스토어 (US) 기준"]
    col_pairs = [("paid_display_pct_ww", "paid_search_pct_ww"),
                 ("paid_display_pct_us", "paid_search_pct_us")]

    for ax, title, (disp_col, srch_col) in zip(axes, titles, col_pairs):
        disp_vals = [agg.loc[g, disp_col] if g in agg.index else 0 for g in groups]
        srch_vals = [agg.loc[g, srch_col] if g in agg.index else 0 for g in groups]

        x = np.arange(len(groups))
        w = 0.38

        bars1 = ax.bar(x - w/2, disp_vals, w, label="Paid Display",
                       color=colors_bar, edgecolor="white", linewidth=1.1, alpha=0.9)
        bars2 = ax.bar(x + w/2, srch_vals, w, label="Paid Search",
                       color=colors_bar, edgecolor="white", linewidth=1.1, alpha=0.5,
                       hatch="///")

        # 값 레이블
        for bar, v in zip(bars1, disp_vals):
            if v > 0:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                        f"{v:.1f}%", ha="center", va="bottom", fontsize=8, fontweight="bold")
        for bar, v in zip(bars2, srch_vals):
            if v > 0:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                        f"{v:.1f}%", ha="center", va="bottom", fontsize=8, fontweight="bold")

        ax.set_title(title, fontsize=12, fontweight="bold", color="#1E2761", pad=8)
        ax.set_xticks(x)
        ax.set_xticklabels(codes, fontsize=11, fontweight="bold")
        ax.set_ylabel("다운로드 중 광고 설치 비중 (%)", fontsize=9)
        ax.set_ylim(0, max(max(disp_vals + srch_vals) * 1.3, 5))
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.set_facecolor("#F8FAFF")
        ax.spines[["top", "right"]].set_visible(False)

        # 범례: Display=진한색, Search=빗금
        from matplotlib.patches import Patch
        legend_handles = [
            Patch(facecolor="#888", edgecolor="white", label="Paid Display (배너·영상)"),
            Patch(facecolor="#888", edgecolor="white", hatch="///",
                  alpha=0.5, label="Paid Search (앱스토어 검색광고)")
        ]
        ax.legend(handles=legend_handles, loc="upper right", fontsize=8, framealpha=0.8)

    fig.tight_layout()
    return fig_to_bytes(fig)


# ── MAIN ──────────────────────────────────────────────────────────────────────

def main():
    df2  = pd.read_csv(BASE / "slide2_v2.csv")
    tags = pd.read_csv(BASE / "app_tags.csv")

    print(f"▶ PPT 로드: {PPTX_IN.name}")
    prs = Presentation(str(PPTX_IN))
    print(f"  총 {len(prs.slides)}개 슬라이드")

    slides = prs.slides

    # ── Slide 18 (index 17): CN 제거 + KR/JP 겹침 선 스타일 완전 구분 ─────────
    print("  [Slide 18] 차트 재생성 (CN 제거, KR/JP 선 스타일 명확 구분)...")
    img18 = make_retention_chart_slide18(df2)
    replaced = replace_chart_image(slides[17], img18)
    print(f"    차트 교체: {'성공' if replaced else '이미지 없음 — 새로 추가'}")
    if not replaced:
        slides[17].shapes.add_picture(io.BytesIO(img18),
            Inches(0.25), Inches(1.1), Inches(9.2), Inches(5.4))

    # 기존 Slide 18 하단 노트 박스 교체 (새 내용)
    add_note_box(slides[17],
        "※ KR 실선(●) · JP 점선(▲): 2025 이후 두 국가 잔존율 동일값으로 집계됨  |  "
        "KR=한국 · US=미국 · JP=일본 시장",
        top_inches=6.6, height=0.35)

    # ── Slide 24 새 삽입 (기존 Slide 23 = index 22 바로 뒤, Slide 24 결론 앞) ───
    print("  [Slide 24 신규] 광고집행율 국가별 차트 삽입...")
    img_ad = make_ad_rate_country_chart(tags)

    # 빈 슬라이드 생성
    blank_layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(blank_layout)

    # 배경
    bg = new_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFF)

    # 제목
    add_title_box(new_slide, "광고 집행율 & Paid Install 비중 — 국가별 (KR / JP / CN / US)")

    # 서브타이틀
    sub_tb = new_slide.shapes.add_textbox(
        Inches(0.3), Inches(0.78), Inches(9.1), Inches(0.28))
    sub_run = sub_tb.text_frame.paragraphs[0].add_run()
    sub_run.text = "퍼블리셔 출신 국가 기준  |  KR=한국 · JP=일본 · CN=중화권(중국+홍콩+대만) · US=서구권(미국+EU 등)"
    sub_run.font.size  = Pt(9)
    sub_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # 차트 이미지
    new_slide.shapes.add_picture(io.BytesIO(img_ad),
        Inches(0.25), Inches(1.1), Inches(9.2), Inches(5.3))

    # 하단 노트
    note_tb = new_slide.shapes.add_textbox(
        Inches(0.25), Inches(6.5), Inches(9.2), Inches(0.55))
    note_tf = note_tb.text_frame; note_tf.word_wrap = True
    note_p  = note_tf.paragraphs[0]; note_p.alignment = PP_ALIGN.LEFT
    note_run = note_p.add_run()
    note_run.text = (
        "Paid Display: 배너·영상·플레이어블 광고 통해 설치  |  "
        "Paid Search: 앱스토어 검색광고(Apple Search Ads 등) 통해 설치\n"
        "※ KR(한국) 퍼블리셔 앱 수가 적어 통계적 신뢰도 낮을 수 있음  |  Source: Sensor Tower API"
    )
    note_run.font.size = Pt(7.8)
    note_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # Slide 24 위치에 삽입 (index=23, 즉 기존 슬라이드 24인 '결론' 바로 앞)
    insert_slide_at(prs, new_slide, 23)

    # ── 저장 ─────────────────────────────────────────────────────────────────
    prs.save(str(PPTX_IN))
    print(f"\n▶ 저장 완료: {PPTX_IN.name}")
    print(f"  총 {len(prs.slides)}개 슬라이드")
    print("  Slide 18: KR/US/JP만, 선 스타일 명확 구분")
    print("  Slide 24: 광고집행율 국가별 (KR/JP/CN/US) — 신규 삽입")
    print("  Slide 25: 결론 (기존 24번)")


if __name__ == "__main__":
    main()
