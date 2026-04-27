"""
PPT 슬라이드 업데이트:
  슬라이드 15 → 15a(국가별) + 15b(퍼블리셔별)
  슬라이드 16 → 16a(국가별) + 16b(퍼블리셔별)
  슬라이드 17 → 17a(국가별) + 17b(퍼블리셔별)
  각 슬라이드: 2025 이전(2022-2024) vs 이후(2025~) 비교

디자인: navy #1E2761 / teal #0D9488 / light bg #F0F4FF
"""
import sys, zipfile, io, textwrap
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")
from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent))

BASE_DIR = Path("C:/Users/NHN/Documents/sensortower_api")
UNPACKED = BASE_DIR / "pptx_unpacked3"
OUT_PPTX = BASE_DIR / "MobileGame_Market_Analysis_2022-2026_v2.pptx"
TMP_PPTX = BASE_DIR / "_tmp_base.pptx"

# ── 색상 ──────────────────────────────────────────────────────────────────────
NAVY   = "#1E2761"
TEAL   = "#0D9488"
LIGHT  = "#F0F4FF"
WHITE  = "#FFFFFF"
GRAY   = "#6B7280"

PRE_COLOR  = "#1E2761"   # 2025 이전 = navy
POST_COLOR = "#0D9488"   # 2025 이후 = teal

COUNTRY_COLORS   = {"KR": "#E74C3C", "US": "#3498DB", "JP": "#2ECC71", "CN": "#F39C12"}
PUBLISHER_COLORS = {"한국": "#E74C3C", "일본": "#2ECC71", "중화권": "#F39C12", "서구권": "#3498DB", "기타": "#95A5A6"}

rcParams["font.family"]     = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False


# ─────────────────────────────────────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────────────────────────────────────

def pack_pptx():
    """pptx_unpacked3 → _tmp_base.pptx"""
    with zipfile.ZipFile(TMP_PPTX, "w", zipfile.ZIP_DEFLATED) as zf:
        for fp in UNPACKED.rglob("*"):
            if fp.is_file():
                zf.write(fp, fp.relative_to(UNPACKED))
    print(f"  [pack] {TMP_PPTX.name}")


def fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def add_slide_copy(prs: Presentation, source_idx: int) -> object:
    """source_idx 번째 슬라이드를 복사해 마지막에 추가."""
    from pptx.oxml.ns import qn
    import copy, lxml.etree as etree

    template = prs.slides[source_idx]
    blank_layout = prs.slide_layouts[0]   # blank
    new_slide = prs.slides.add_slide(blank_layout)

    # XML 내용 복사
    new_slide.shapes._spTree.clear()
    for elem in template.shapes._spTree:
        new_slide.shapes._spTree.append(copy.deepcopy(elem))

    return new_slide


def set_slide_title(slide, text: str, font_size=24):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 13:
            continue
        if shape.has_text_frame:
            tf = shape.text_frame
            if tf.paragraphs and len(tf.text) > 5:
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                tf.paragraphs[0].runs[0].text = text
                tf.paragraphs[0].runs[0].font.size = Pt(font_size)
                return


def clear_content_shapes(slide, keep_title=True):
    """차트 이미지 외 본문 도형 제거 (제목 제외)."""
    from pptx.oxml.ns import qn
    spTree = slide.shapes._spTree
    to_remove = []
    for shape in slide.shapes:
        if keep_title and shape.has_text_frame:
            # 가장 위쪽 텍스트 도형은 제목으로 보존
            continue
        to_remove.append(shape._element)
    for el in to_remove:
        try:
            spTree.remove(el)
        except Exception:
            pass


def add_title_box(slide, text: str):
    """navy 배경 제목 박스 추가."""
    txBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.1), Inches(9.1), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(20)
    run.font.bold    = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)


def add_label_box(slide, text: str, left, top, width, height,
                  bg="#1E2761", fg="#FFFFFF", font_size=11):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = True
    run.font.color.rgb = RGBColor(
        int(fg[1:3],16), int(fg[3:5],16), int(fg[5:7],16))
    fill = tb.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(
        int(bg[1:3],16), int(bg[3:5],16), int(bg[5:7],16))


def add_image(slide, img_bytes: bytes, left, top, width, height):
    slide.shapes.add_picture(
        io.BytesIO(img_bytes), left, top, width, height)


def add_source_note(slide):
    tb = slide.shapes.add_textbox(
        Inches(0.3), Inches(7.0), Inches(9.1), Inches(0.3))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Source: Sensor Tower API  |  iOS TOP 100 기준  |  2025 전 = 2022–2024, 2025 후 = 2025–2026"
    run.font.size  = Pt(8)
    run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)


# ─────────────────────────────────────────────────────────────────────────────
# 차트 생성 함수들
# ─────────────────────────────────────────────────────────────────────────────

def make_survival_chart(df1: pd.DataFrame, group_col: str,
                        groups: list, colors: dict, title: str) -> bytes:
    """
    생존/탈락 비율 — 2025 전후 비교 grouped bar chart
    Returns PNG bytes
    """
    pre  = df1[df1["period"] == "pre2025"]
    post = df1[df1["period"] == "post2025"]

    def survival_rate(df, grp):
        sub = df[df[group_col] == grp]
        if len(sub) == 0:
            return 0, 0
        rate = sub["survived_3m"].mean() * 100
        n    = len(sub)
        return round(rate, 1), n

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5), sharey=True)
    fig.patch.set_alpha(0)

    for ax, period_df, period_label in [
            (axes[0], pre, "2025 이전  (2022–2024)"),
            (axes[1], post, "2025 이후  (2025–2026)")]:

        rates, labels, clrs, ns = [], [], [], []
        for g in groups:
            rate, n = survival_rate(period_df, g)
            rates.append(rate)
            labels.append(g)
            clrs.append(colors.get(g, "#95A5A6"))
            ns.append(n)

        x   = np.arange(len(groups))
        bars = ax.bar(x, rates, color=clrs, width=0.55, edgecolor="white", linewidth=1.2)

        for bar, rate, n in zip(bars, rates, ns):
            ax.text(bar.get_x() + bar.get_width()/2,
                    bar.get_height() + 1.5, f"{rate:.1f}%\n(n={n})",
                    ha="center", va="bottom", fontsize=8.5, fontweight="bold")

        ax.set_title(period_label, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_xticks(x)
        ax.set_xticklabels(labels, fontsize=10)
        ax.set_ylim(0, 105)
        ax.set_ylabel("3개월 생존율 (%)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.set_facecolor("#F8FAFF")
        ax.spines[["top","right"]].set_visible(False)

    fig.suptitle(title, fontsize=13, fontweight="bold", color="#1E2761", y=1.01)
    fig.tight_layout()
    return fig_to_bytes(fig)


def make_retention_chart(df2: pd.DataFrame, group_col: str,
                         groups: list, colors: dict, title: str) -> bytes:
    """잔존율 D1/D7/D14/D30 — 2025 전후 비교 line chart"""
    days_cols = ["d1", "d7", "d14", "d30"]
    days_labels = ["D1", "D7", "D14", "D30"]

    pre  = df2[df2["period"] == "pre2025"]
    post = df2[df2["period"] == "post2025"]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))
    fig.patch.set_alpha(0)

    for ax, period_df, period_label in [
            (axes[0], pre, "2025 이전  (2022–2024)"),
            (axes[1], post, "2025 이후  (2025–2026)")]:

        for grp in groups:
            sub = period_df[period_df[group_col] == grp]
            if sub.empty:
                continue
            vals = [sub[c].mean() for c in days_cols]
            ax.plot(days_labels, vals, marker="o", linewidth=2,
                    color=colors.get(grp, "#95A5A6"), label=grp)
            ax.text(days_labels[-1], vals[-1], f" {grp}",
                    fontsize=8, color=colors.get(grp, "#95A5A6"), va="center")

        ax.set_title(period_label, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_ylabel("잔존율 (%)", fontsize=9)
        ax.set_ylim(0, max(60, ax.get_ylim()[1] + 5))
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.set_facecolor("#F8FAFF")
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(loc="upper right", fontsize=8, framealpha=0.7)

    fig.suptitle(title, fontsize=13, fontweight="bold", color="#1E2761", y=1.01)
    fig.tight_layout()
    return fig_to_bytes(fig)


def make_revenue_chart(df3: pd.DataFrame, group_col: str,
                       groups: list, colors: dict, title: str) -> bytes:
    """
    순위별 월 평균 매출 — 2025 전후 비교
    group_col = 'country' or 'publisher_group'
    """
    TARGET_RANKS = [1, 10, 20, 50, 100]

    pre  = df3[df3["period"] == "pre2025"]
    post = df3[df3["period"] == "post2025"]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))
    fig.patch.set_alpha(0)

    for ax, period_df, period_label in [
            (axes[0], pre,  "2025 이전  (2022–2024)"),
            (axes[1], post, "2025 이후  (2025–2026)")]:

        x   = np.arange(len(TARGET_RANKS))
        w   = 0.7 / max(len(groups), 1)
        for i, grp in enumerate(groups):
            sub   = period_df[period_df[group_col] == grp]
            means = [sub[sub["rank"] == r]["revenue_m_usd"].mean() for r in TARGET_RANKS]
            means = [v if not np.isnan(v) else 0 for v in means]
            offset = (i - len(groups)/2 + 0.5) * w
            bars = ax.bar(x + offset, means, width=w,
                          color=colors.get(grp, "#95A5A6"), label=grp,
                          edgecolor="white", linewidth=0.8)

        ax.set_title(period_label, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_xticks(x)
        ax.set_xticklabels([f"#{r}" for r in TARGET_RANKS], fontsize=9)
        ax.set_xlabel("차트 순위", fontsize=9)
        ax.set_ylabel("월 평균 매출 (백만 USD)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.set_facecolor("#F8FAFF")
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(loc="upper right", fontsize=8, framealpha=0.7)

    fig.suptitle(title, fontsize=13, fontweight="bold", color="#1E2761", y=1.01)
    fig.tight_layout()
    return fig_to_bytes(fig)


# ─────────────────────────────────────────────────────────────────────────────
# 슬라이드 빌더
# ─────────────────────────────────────────────────────────────────────────────

def build_new_slide(prs: Presentation, title: str, chart_bytes: bytes,
                    subtitle: str = ""):
    """새 슬라이드 추가 (blank layout + 제목 + 차트 이미지)."""
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)

    # 배경
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFF)

    # 제목 박스
    add_title_box(slide, title)

    # 서브타이틀
    if subtitle:
        tb = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.78), Inches(9.1), Inches(0.28))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = subtitle
        run.font.size  = Pt(9)
        run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # 차트 이미지
    top_offset = Inches(1.1) if subtitle else Inches(0.9)
    add_image(slide, chart_bytes,
              left=Inches(0.25), top=top_offset,
              width=Inches(9.2), height=Inches(5.6))

    add_source_note(slide)
    return slide


def insert_slide_at(prs: Presentation, slide, position: int):
    """슬라이드를 특정 위치에 삽입."""
    from pptx.oxml.ns import qn
    xml_slides = prs.slides._sldIdLst
    # 마지막에 추가된 슬라이드를 position으로 이동
    last = xml_slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(position, last)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    # ── 데이터 로드 ──────────────────────────────────────────────────────────
    print("▶ 데이터 로드 중...")

    df1_path = BASE_DIR / "slide1_v2.csv"
    df2_path = BASE_DIR / "slide2_v2.csv"
    df3_path = BASE_DIR / "slide3_v2.csv"

    # v2가 없으면 기존 파일로 fallback
    df1 = pd.read_csv(df1_path if df1_path.exists() else BASE_DIR/"slide1_classification_all.csv")
    df2 = pd.read_csv(df2_path if df2_path.exists() else BASE_DIR/"slide2_retention_all.csv")
    df3 = pd.read_csv(df3_path if df3_path.exists() else BASE_DIR/"slide3_rank_revenue.csv")

    # 공통 컬럼 보정
    if "period" not in df1.columns:
        df1["period"] = df1["first_chart_month"].apply(
            lambda x: "post2025" if str(x) >= "2025-01" else "pre2025")
    if "period" not in df2.columns:
        df2["period"] = df2["year"].apply(
            lambda x: "post2025" if int(x) >= 2025 else "pre2025")
    if "period" not in df3.columns:
        df3["period"] = df3["ym"].apply(
            lambda x: "post2025" if str(x) >= "2025-01" else "pre2025")
    if "publisher_group" not in df1.columns:
        df1["publisher_group"] = "미분류"
    if "publisher_group" not in df2.columns:
        df2["publisher_group"] = "미분류"
    if "publisher_group" not in df3.columns:
        df3["publisher_group"] = "미분류"

    print(f"  Slide1: {len(df1)}행 / Slide2: {len(df2)}행 / Slide3: {len(df3)}행")

    # ── PPT 로드 ──────────────────────────────────────────────────────────────
    print("▶ PPT 패킹 및 로드...")
    pack_pptx()
    prs = Presentation(str(TMP_PPTX))
    total = len(prs.slides)
    print(f"  총 {total}개 슬라이드")

    # ── 그룹 정의 ─────────────────────────────────────────────────────────────
    country_groups   = ["KR", "US", "JP"]
    publisher_groups = ["한국", "일본", "중화권", "서구권"]

    # 슬라이드 15 (index 14) → 생존/탈락
    # 슬라이드 16 (index 15) → 잔존율
    # 슬라이드 17 (index 16) → 월매출

    # 슬라이드 14=생존/탈락(idx 13), 15=잔존율(idx 14), 16=월매출(idx 15)
    # 각 뒤에 2개씩 삽입 → 삽입 순서대로 offset 보정
    new_slides_spec = [
        # (원본 0-based idx, 제목, 서브타이틀, 차트함수, 파라미터)
        (13, "신규 게임 3개월 생존율 — 국가별 × 2025 전후",
         "iOS TOP 100 차트 신규 진입 게임 기준 (시장별 분석)",
         make_survival_chart,
         dict(df1=df1, group_col="country", groups=country_groups,
              colors=COUNTRY_COLORS, title="")),

        (13, "신규 게임 3개월 생존율 — 퍼블리셔 출신별 × 2025 전후",
         "한국 / 일본 / 중화권 / 서구권 퍼블리셔 기준",
         make_survival_chart,
         dict(df1=df1, group_col="publisher_group", groups=publisher_groups,
              colors=PUBLISHER_COLORS, title="")),

        (14, "유저 잔존율 — 국가별 × 2025 전후",
         "D1 / D7 / D14 / D30 평균 잔존율 (시장별 비교)",
         make_retention_chart,
         dict(df2=df2, group_col="country", groups=country_groups,
              colors=COUNTRY_COLORS, title="")),

        (14, "유저 잔존율 — 퍼블리셔 출신별 × 2025 전후",
         "한국 / 일본 / 중화권 / 서구권 퍼블리셔 기준",
         make_retention_chart,
         dict(df2=df2, group_col="publisher_group", groups=publisher_groups,
              colors=PUBLISHER_COLORS, title="")),

        (15, "순위별 월 평균 매출 — 국가별 × 2025 전후",
         "차트 #1 / #10 / #20 / #50 / #100 기준 평균 월매출 (백만 USD)",
         make_revenue_chart,
         dict(df3=df3, group_col="country", groups=["KR","US","JP","CN"],
              colors=COUNTRY_COLORS, title="")),

        (15, "순위별 월 평균 매출 — 퍼블리셔 출신별 × 2025 전후",
         "한국 / 일본 / 중화권 / 서구권 퍼블리셔 기준",
         make_revenue_chart,
         dict(df3=df3, group_col="publisher_group", groups=publisher_groups,
              colors=PUBLISHER_COLORS, title="")),
    ]

    # 새 슬라이드 6개 생성 및 삽입
    # 삽입할 때마다 이후 인덱스가 밀리므로 누적 카운트로 보정
    print("▶ 새 슬라이드 생성 중...")
    total_inserted = 0
    for base_idx, title, subtitle, chart_fn, params in new_slides_spec:
        # 원본 슬라이드의 현재 실제 위치: base_idx + total_inserted
        # 삽입 위치: 그 바로 뒤
        insert_pos = base_idx + total_inserted + 1

        print(f"  [slide {insert_pos+1}] {title[:40]}...")
        params["title"] = ""
        chart_bytes = chart_fn(**params)

        slide = build_new_slide(prs, title, chart_bytes, subtitle)
        insert_slide_at(prs, slide, insert_pos)
        total_inserted += 1

    # 저장
    print(f"▶ 저장: {OUT_PPTX.name}")
    prs.save(str(OUT_PPTX))
    TMP_PPTX.unlink(missing_ok=True)
    print(f"  완료! 총 {len(prs.slides)}개 슬라이드")


if __name__ == "__main__":
    main()
