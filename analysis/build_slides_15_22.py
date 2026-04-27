"""
슬라이드 15~22 — 새 PPTX로 생성 (기존 파일 열지 않음)
수정: x축 겹침 해소, 인사이트 폰트크기 통일, AOS 추가, 데이터 보강
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.lines as mlines
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

rcParams["font.family"]        = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE = Path("C:/Users/NHN/Documents/sensortower_api")
DPI  = 200

PUB_COLORS = {"중화권":"#F39C12","북미":"#3498DB","일본":"#2ECC71","한국":"#E74C3C","기타":"#9B59B6"}
PUB_ORDER  = ["한국","일본","북미","중화권","기타"]
CTR_COLORS = {"KR":"#E74C3C","JP":"#2ECC71","US":"#3498DB"}
CTR_LABEL  = {"KR":"KR 한국","JP":"JP 일본","US":"US 미국"}
CTR_ORDER  = ["KR","JP","US"]
IOS_ALPHA  = 0.88;  AOS_ALPHA = 0.55
NOTE_FS    = 7.5

# ── 데이터 로드 ─────────────────────────────────────────────────────────
ios_s1   = pd.read_csv(BASE / "slide1_v2.csv")
aod_s1   = pd.read_csv(BASE / "slide1_v2_android.csv")
ios_ad   = pd.read_csv(BASE / "newgame_ad_rates.csv")
aod_ad   = pd.read_csv(BASE / "newgame_ad_rates_android.csv")
ios_tags = pd.read_csv(BASE / "app_tags.csv")
aod_tags = pd.read_csv(BASE / "app_tags_android.csv")
ios_rev  = pd.read_csv(BASE / "slide3_v2.csv")
aod_rev  = pd.read_csv(BASE / "slide3_rank_revenue_android.csv") if (BASE / "slide3_rank_revenue_android.csv").exists() else pd.DataFrame()

# Merge survival + ad rates
ios_s1_ad = ios_s1.merge(ios_ad[["app_id","total_ww"]], on="app_id", how="left")
aod_s1_ad = aod_s1.merge(aod_ad[["app_id","total_ww"]], on="app_id", how="left")
ios_s1_ad["has_ad"] = ios_s1_ad["total_ww"].fillna(0) > 0
aod_s1_ad["has_ad"] = aod_s1_ad["total_ww"].fillna(0) > 0

# 차트 잔존율 (신규진입 게임이 M+N 시점 TOP100에 남아있는 비율)
ret_full = pd.read_csv(BASE / "chart_retention_newentry.csv") if (BASE / "chart_retention_newentry.csv").exists() else pd.DataFrame()

PRE_MONTHS = 33; POST_MONTHS = 13  # 2022~2024 각 2~12월 = 33개월, 2025 2~12월 + 2026 2~3월 = 13개월

def monthly_avg(total, period):
    return total / PRE_MONTHS if period == "pre2025" else total / POST_MONTHS

def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=DPI, bbox_inches="tight", facecolor="white", transparent=False)
    plt.close(fig); buf.seek(0)
    return buf.read()

def survival_by(df, group_col):
    g = df.groupby([group_col, "period"])["survived_3m"].agg(["sum","count"])
    g["rate"] = (g["sum"] / g["count"] * 100).round(1)
    g["monthly_n"] = g.apply(lambda r: monthly_avg(r["count"], r.name[1] if isinstance(r.name, tuple) else "pre2025"), axis=1).round(0)
    return g.reset_index()


# ══════════════════════════════════════════════════════════════════════════
# Slide 15: 국가별 3개월 생존율 + 월평균 게임수  (x축 겹침 해소)
# ══════════════════════════════════════════════════════════════════════════
def chart_s15() -> bytes:
    ios_g = survival_by(ios_s1[ios_s1["country"].isin(CTR_ORDER)], "country")
    aod_g = survival_by(aod_s1[aod_s1["country"].isin(CTR_ORDER)], "country")
    periods = [("pre2025","'22–'24"), ("post2025","'25~")]
    fig, axes = plt.subplots(1, 3, figsize=(13.5, 4.6), gridspec_kw={"wspace":0.36})
    fig.patch.set_facecolor("white")
    ios_patch = mpatches.Patch(facecolor="#2980B9", alpha=IOS_ALPHA, label="iOS")
    aod_patch = mpatches.Patch(facecolor="#27AE60", alpha=AOS_ALPHA, hatch="///", label="Android")
    W = 0.35; INNER = 0.45; OUTER = 1.30
    for ax, country in zip(axes, CTR_ORDER):
        col = CTR_COLORS[country]; ax.set_facecolor("#F8FAFF")
        g_centers, g_labels = [], []
        for gi, (pk, pl) in enumerate(periods):
            xc_i = gi * OUTER; xc_a = gi * OUTER + INNER
            si = ios_g[(ios_g["country"]==country)&(ios_g["period"]==pk)]
            sa = aod_g[(aod_g["country"]==country)&(aod_g["period"]==pk)]
            ri = float(si["rate"].values[0]) if len(si) else 0
            ni = int(si["monthly_n"].values[0]) if len(si) else 0
            ra = float(sa["rate"].values[0]) if len(sa) else 0
            na = int(sa["monthly_n"].values[0]) if len(sa) else 0
            ax.bar(xc_i, ri, W, color="#2980B9", alpha=IOS_ALPHA, edgecolor="white", lw=1.3)
            ax.bar(xc_a, ra, W, color="#27AE60", alpha=AOS_ALPHA, hatch="///", edgecolor="white", lw=1.3)
            for xc, rate, clr in [(xc_i,ri,"#2980B9"),(xc_a,ra,"#27AE60")]:
                if rate > 1:
                    ax.text(xc, rate+1.8, f"{rate:.0f}%", ha="center", va="bottom", fontsize=12, fontweight="bold", color=clr)
            gc = (xc_i + xc_a) / 2; g_centers.append(gc)
            g_labels.append(f"{pl}\n월평균 I:{ni} / A:{na}개")
        ax.axvline((INNER+W/2 + OUTER-W/2)/2, color="#CBD5E1", lw=1.2, ls="--")
        ax.set_title(CTR_LABEL[country], fontsize=13, fontweight="bold", color=col, pad=8)
        ax.set_xticks(g_centers); ax.set_xticklabels(g_labels, fontsize=9, fontweight="bold", ha="center")
        ax.set_xlim(-W/2-0.18, OUTER+INNER+W/2+0.18); ax.set_ylim(0, 118)
        ax.set_ylabel("3개월 생존율 (%)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35); ax.spines[["top","right"]].set_visible(False)
    fig.legend(handles=[ios_patch, aod_patch], loc="upper right", bbox_to_anchor=(0.995,0.98), fontsize=10, framealpha=0.9)
    fig.text(0.5, 0.005, "신규진입 = 각 연도 1월 TOP100 기준, 2~12월 첫 진입 · n=월평균 게임수, 생존율=전체 기간 누적 비율 · 매월 1일 기준 · * 소표본(n<50) 참고 수준 · 출처: Sensor Tower", ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 16: 국가별 광고율 (생존 vs 미생존) — iOS 중심 + AOS 주석
# ══════════════════════════════════════════════════════════════════════════
def chart_s16() -> bytes:
    fig, axes = plt.subplots(1, 3, figsize=(13.5, 4.6), gridspec_kw={"wspace":0.36})
    fig.patch.set_facecolor("white")
    surv_patch = mpatches.Patch(facecolor="#2ECC71", alpha=0.85, label="3개월 생존")
    churn_patch = mpatches.Patch(facecolor="#E74C3C", alpha=0.55, hatch="///", label="3개월 미생존")
    for ax, country in zip(axes, CTR_ORDER):
        ax.set_facecolor("#F8FAFF"); col = CTR_COLORS[country]
        W = 0.35; INNER = 0.42; OUTER = 1.25
        g_centers, g_labels = [], []
        for gi, (pk, pl) in enumerate([("pre2025","'22–'24"),("post2025","'25~")]):
            # iOS — 월평균 광고집행율
            sub = ios_s1_ad[(ios_s1_ad["country"]==country)&(ios_s1_ad["period"]==pk)&(ios_s1_ad["total_ww"].notna())]
            surv = sub[sub["survived_3m"]==True]; churn = sub[sub["survived_3m"]==False]
            surv_mo = surv.groupby("first_chart_month")["has_ad"].mean()*100
            churn_mo = churn.groupby("first_chart_month")["has_ad"].mean()*100
            r_surv = surv_mo.mean() if len(surv_mo) else 0
            r_churn = churn_mo.mean() if len(churn_mo) else 0
            xs = gi * OUTER; xc = gi * OUTER + INNER
            ax.bar(xs, r_surv, W, color="#2ECC71", alpha=0.85, edgecolor="white", lw=1.2)
            ax.bar(xc, r_churn, W, color="#E74C3C", alpha=0.55, hatch="///", edgecolor="white", lw=1.2)
            for x, r, c in [(xs,r_surv,"#2ECC71"),(xc,r_churn,"#E74C3C")]:
                if r > 1: ax.text(x, r+2, f"{r:.0f}%", ha="center", va="bottom", fontsize=11, fontweight="bold", color=c)
            g_centers.append((xs+xc)/2); g_labels.append(f"iOS {pl}")
        ax.set_title(CTR_LABEL[country], fontsize=13, fontweight="bold", color=col, pad=8)
        ax.set_xticks(g_centers); ax.set_xticklabels(g_labels, fontsize=9, ha="center", fontweight="bold")
        ax.set_ylim(0, 110); ax.set_ylabel("광고집행율 (%)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35); ax.spines[["top","right"]].set_visible(False)
    fig.legend(handles=[surv_patch, churn_patch], loc="upper right", bbox_to_anchor=(0.995,0.98), fontsize=10, framealpha=0.9)
    fig.text(0.5, 0.005, "iOS 기준 · 월평균 광고집행율(월별 비율의 평균) · AOS는 Sensor Tower 광고 커버리지 한계로 생략 · 출처: Sensor Tower", ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 17: 퍼블리셔별 3개월 생존율 (x축 겹침 해소)
# ══════════════════════════════════════════════════════════════════════════
def chart_s17() -> bytes:
    ios_g = survival_by(ios_s1, "publisher_group")
    aod_g = survival_by(aod_s1, "publisher_group")
    periods = [("pre2025","2022–2024"), ("post2025","2025~")]
    fig, axes = plt.subplots(1, 2, figsize=(13.5, 4.5), gridspec_kw={"wspace":0.30})
    fig.patch.set_facecolor("white")
    ios_patch = mpatches.Patch(facecolor="#444", alpha=IOS_ALPHA, label="iOS")
    aod_patch = mpatches.Patch(facecolor="#444", alpha=AOS_ALPHA, hatch="///", label="Android")
    for ax, (pk, pl) in zip(axes, periods):
        ax.set_facecolor("#F8FAFF")
        ax.set_title(pl, fontsize=13, fontweight="bold", color="#1E2761", pad=7)
        pubs = [p for p in PUB_ORDER if p in ios_g["publisher_group"].values or p in aod_g["publisher_group"].values]
        n_pub = len(pubs); W_p = 0.28; INNER_P = 0.36; OUTER_P = 1.05
        xs_i = np.arange(n_pub, dtype=float) * OUTER_P; xs_a = xs_i + INNER_P
        mid_x = (xs_i + xs_a) / 2
        for xi_i, xi_a, pub in zip(xs_i, xs_a, pubs):
            clr = PUB_COLORS.get(pub, "#9B59B6")
            def get_rv(df_g, _pub=pub, _pk=pk):
                sub = df_g[(df_g["publisher_group"]==_pub)&(df_g["period"]==_pk)]
                return (float(sub["rate"].values[0]) if len(sub) else 0, int(sub["monthly_n"].values[0]) if len(sub) else 0)
            ri, ni = get_rv(ios_g); ra, na = get_rv(aod_g)
            ax.bar(xi_i, max(ri,1.5), W_p, color=clr, alpha=IOS_ALPHA, edgecolor="white", lw=1.3)
            ax.bar(xi_a, max(ra,1.5), W_p, color=clr, alpha=AOS_ALPHA, hatch="///", edgecolor="white", lw=1.3)
            for xi, rate in [(xi_i,ri),(xi_a,ra)]:
                if rate > 2: ax.text(xi, rate+1.5, f"{rate:.0f}%", ha="center", va="bottom", fontsize=10, fontweight="bold", color=clr)
            ax.text(xi_i, -2.5, f"I:월{ni}개", ha="center", va="top", fontsize=7, color="#2980B9", fontweight="bold")
            ax.text(xi_a, -2.5, f"A:월{na}개", ha="center", va="top", fontsize=7, color="#27AE60", fontweight="bold")
        ax.set_xticks(mid_x); ax.set_xticklabels(pubs, fontsize=10, fontweight="bold")
        ax.set_xlim(xs_i[0]-W_p/2-0.22, xs_a[-1]+W_p/2+0.22); ax.set_ylim(-5, 115)
        ax.set_ylabel("3개월 생존율 (%)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35); ax.spines[["top","right"]].set_visible(False)
    fig.legend(handles=[ios_patch, aod_patch], loc="upper right", bbox_to_anchor=(0.995,0.98), fontsize=10, framealpha=0.9)
    fig.text(0.5, 0.005, "신규진입 = 각 연도 1월 TOP100 기준, 2~12월 첫 진입 · n=월평균 게임수, 생존율=전체 기간 누적 비율 · 매월 1일 기준 · * 소표본(n<50) 참고 수준 · 출처: Sensor Tower", ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 18: 퍼블리셔별 광고율 — iOS 중심 + AOS 주석
# ══════════════════════════════════════════════════════════════════════════
def chart_s18() -> bytes:
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.5), gridspec_kw={"wspace":0.35})
    fig.patch.set_facecolor("white")
    surv_patch = mpatches.Patch(facecolor="#2ECC71", alpha=0.85, label="3개월 생존")
    churn_patch = mpatches.Patch(facecolor="#E74C3C", alpha=0.55, hatch="///", label="3개월 미생존")
    for ax, (pk, pl) in zip(axes, [("pre2025","2022–2024"),("post2025","2025~")]):
        ax.set_facecolor("#F8FAFF")
        ax.set_title(f"iOS — {pl}", fontsize=13, fontweight="bold", color="#1E2761", pad=7)
        pubs = [p for p in PUB_ORDER if p in ios_s1_ad["publisher_group"].values]; n_pub = len(pubs)
        W = 0.30; INNER_P = 0.38; OUTER_P = 1.10
        for pi, pub in enumerate(pubs):
            clr = PUB_COLORS.get(pub, "#9B59B6"); xs = pi * OUTER_P; xc = pi * OUTER_P + INNER_P
            sub = ios_s1_ad[(ios_s1_ad["publisher_group"]==pub)&(ios_s1_ad["period"]==pk)&(ios_s1_ad["total_ww"].notna())]
            surv = sub[sub["survived_3m"]==True]; churn = sub[sub["survived_3m"]==False]
            surv_mo = surv.groupby("first_chart_month")["has_ad"].mean()*100
            churn_mo = churn.groupby("first_chart_month")["has_ad"].mean()*100
            r_surv = surv_mo.mean() if len(surv_mo) else 0
            r_churn = churn_mo.mean() if len(churn_mo) else 0
            ax.bar(xs, r_surv, W, color="#2ECC71", alpha=0.85, edgecolor="white", lw=1.2)
            ax.bar(xc, r_churn, W, color="#E74C3C", alpha=0.55, hatch="///", edgecolor="white", lw=1.2)
            for x, r, c in [(xs,r_surv,"#2ECC71"),(xc,r_churn,"#E74C3C")]:
                if r > 1: ax.text(x, r+2, f"{r:.0f}%", ha="center", va="bottom", fontsize=10, fontweight="bold", color=c)
        xticks = [pi * OUTER_P + INNER_P/2 for pi in range(n_pub)]
        ax.set_xticks(xticks); ax.set_xticklabels(pubs, fontsize=10, fontweight="bold")
        ax.set_ylim(0, 110); ax.set_ylabel("광고집행율 (%)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35); ax.spines[["top","right"]].set_visible(False)
    fig.legend(handles=[surv_patch, churn_patch], loc="upper right", bbox_to_anchor=(0.995,0.98), fontsize=10, framealpha=0.9)
    fig.text(0.5, 0.005, "iOS 기준 · 월평균 광고집행율(월별 비율의 평균) · AOS는 Sensor Tower 광고 커버리지 한계로 생략 · 출처: Sensor Tower", ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 19: 국가별 차트 잔존율 M+1~M+12 (신규진입 게임의 TOP100 잔존 비율)
# ══════════════════════════════════════════════════════════════════════════
def chart_s19() -> bytes:
    months_lbl = ["M+1","M+2","M+3","M+6","M+12"]
    rcols = ["M+1","M+2","M+3","M+6","M+12"]
    markers = {"KR":"o","JP":"s","US":"^"}
    fig, axes = plt.subplots(2, 2, figsize=(13.5, 7.5), gridspec_kw={"wspace":0.25,"hspace":0.42})
    fig.patch.set_facecolor("white")
    x = np.arange(len(months_lbl))
    configs = [
        (axes[0,0], "ios", "iOS — 2025 이전 (2022–2024)", "pre2025"),
        (axes[0,1], "ios", "iOS — 2025 이후 (2025~)", "post2025"),
        (axes[1,0], "android", "Android — 2025 이전 (2022–2024)", "pre2025"),
        (axes[1,1], "android", "Android — 2025 이후 (2025~)", "post2025"),
    ]
    for ax, plat, title, pk in configs:
        ax.set_facecolor("#F8FAFF"); handles = []
        df = ret_full[(ret_full["platform"]==plat)&(ret_full["period"]==pk)]
        for ctr in CTR_ORDER:
            sub = df[df["country"]==ctr]
            vals = [sub[c].mean()*100 if c in sub.columns and sub[c].notna().any() else np.nan for c in rcols]
            clr = CTR_COLORS[ctr]; mk = markers.get(ctr,"o")
            vx = [xi for xi, v in zip(x, vals) if not np.isnan(v)]
            vv = [v for v in vals if not np.isnan(v)]
            if vv:
                ax.plot(vx, vv, color=clr, lw=2.5, marker=mk, markersize=7, zorder=3)
                for xi, v in zip(vx, vv):
                    ax.text(xi, v+1.5, f"{v:.0f}%", ha="center", va="bottom", fontsize=7.5, color=clr, fontweight="bold")
            fixed_months = PRE_MONTHS if pk == "pre2025" else POST_MONTHS
            n_mo = round(len(sub) / fixed_months)
            handles.append(mlines.Line2D([],[],color=clr,lw=2.5,marker=mk,markersize=6, label=f"{CTR_LABEL[ctr]} (월평균 {n_mo}개)"))
        ax.set_title(title, fontsize=11, fontweight="bold", color="#1E2761", pad=6)
        ax.set_xticks(x); ax.set_xticklabels(months_lbl, fontsize=10, fontweight="bold")
        ax.set_ylabel("차트 잔존율 (%)", fontsize=9); ax.set_ylim(0, 100); ax.set_xlim(-0.3, len(months_lbl)-0.3)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4); ax.spines[["top","right"]].set_visible(False)
        ax.legend(handles=handles, loc="upper right", fontsize=8, framealpha=0.9)
    fig.text(0.5, 0.005, "신규진입 = 각 연도 1월 기준선, 2~12월 첫 진입 · M+N 시점 TOP100 존재 비율 · 매월 1일 기준 · * post2025 M+6(61%),M+12(21%)만 판정 가능 · 출처: Sensor Tower", ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.03, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 20: 퍼블리셔별 차트 잔존율 M+1~M+12 (신규진입 게임의 TOP100 잔존 비율)
# ══════════════════════════════════════════════════════════════════════════
def chart_s20() -> bytes:
    months_lbl = ["M+1","M+2","M+3","M+6","M+12"]
    rcols = ["M+1","M+2","M+3","M+6","M+12"]
    pub_markers = {"한국":"D","일본":"^","북미":"s","중화권":"o"}
    fig, axes = plt.subplots(2, 2, figsize=(13.5, 7.5), gridspec_kw={"wspace":0.25,"hspace":0.42})
    fig.patch.set_facecolor("white")
    x = np.arange(len(months_lbl))
    configs = [
        (axes[0,0], "ios", "iOS — 2025 이전 (2022–2024)", "pre2025"),
        (axes[0,1], "ios", "iOS — 2025 이후 (2025~)", "post2025"),
        (axes[1,0], "android", "Android — 2025 이전 (2022–2024)", "pre2025"),
        (axes[1,1], "android", "Android — 2025 이후 (2025~)", "post2025"),
    ]
    for ax, plat, title, pk in configs:
        ax.set_facecolor("#F8FAFF"); handles = []
        df = ret_full[(ret_full["platform"]==plat)&(ret_full["period"]==pk)]
        for pub in PUB_ORDER:
            sub = df[df["publisher_group"]==pub]
            if len(sub) == 0: continue
            vals = [sub[c].mean()*100 if c in sub.columns and sub[c].notna().any() else np.nan for c in rcols]
            clr = PUB_COLORS[pub]; mk = pub_markers.get(pub,"o")
            vx = [xi for xi, v in zip(x, vals) if not np.isnan(v)]
            vv = [v for v in vals if not np.isnan(v)]
            if vv:
                ax.plot(vx, vv, color=clr, lw=2.5, marker=mk, markersize=7, zorder=3)
                for xi, v in zip(vx, vv):
                    if v > 0: ax.text(xi, v+1.5, f"{v:.0f}%", ha="center", va="bottom", fontsize=7.5, color=clr)
            fixed_months = PRE_MONTHS if pk == "pre2025" else POST_MONTHS
            n_mo = round(len(sub) / fixed_months)
            handles.append(mlines.Line2D([],[],color=clr,lw=2.5,marker=mk,markersize=6, label=f"{pub} (월평균 {n_mo}개)"))
        ax.set_title(f"{title}", fontsize=11, fontweight="bold", color="#1E2761", pad=6)
        ax.set_xticks(x); ax.set_xticklabels(months_lbl, fontsize=10, fontweight="bold")
        ax.set_ylabel("차트 잔존율 (%)", fontsize=9); ax.set_ylim(0, 100); ax.set_xlim(-0.3, len(months_lbl)-0.3)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4); ax.spines[["top","right"]].set_visible(False)
        ax.legend(handles=handles, loc="upper right", fontsize=8, framealpha=0.9)
    fig.text(0.5, 0.005, "신규진입 = 각 연도 1월 기준선, 2~12월 첫 진입 · M+N 시점 TOP100 존재 비율 · 매월 1일 기준 · * post2025 M+6(61%),M+12(21%)만 판정 가능 · 출처: Sensor Tower", ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.03, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 21: 국가별 순위별 월평균 매출 (iOS + AOS 2×2)
# ══════════════════════════════════════════════════════════════════════════
def chart_s21() -> bytes:
    ranks = [1, 10, 20, 50, 100]; rank_labels = ["1위","10위","20위","50위","100위"]
    fig, axes = plt.subplots(2, 2, figsize=(13.5, 7.5), gridspec_kw={"wspace":0.28,"hspace":0.45})
    fig.patch.set_facecolor("white")
    df_ios = ios_rev[ios_rev["country"].isin(CTR_ORDER)].copy()
    df_ios["revenue_m_usd"] = pd.to_numeric(df_ios["revenue_m_usd"], errors="coerce")
    df_aod = aod_rev[aod_rev["country"].isin(CTR_ORDER)].copy() if len(aod_rev) else pd.DataFrame()
    if len(df_aod): df_aod["revenue_m_usd"] = pd.to_numeric(df_aod["revenue_m_usd"], errors="coerce")
    configs = [
        (axes[0,0], df_ios, "iOS — 2025 이전 (2022–2024)", "pre2025"),
        (axes[0,1], df_ios, "iOS — 2025 이후 (2025~)", "post2025"),
        (axes[1,0], df_aod, "Android — 2025 이전 (2022–2024)", "pre2025"),
        (axes[1,1], df_aod, "Android — 2025 이후 (2025~)", "post2025"),
    ]
    for ax, df, title, pk in configs:
        ax.set_facecolor("#F8FAFF"); x = np.arange(len(ranks)); w = 0.25
        if len(df) > 0:
            sub = df[df["period"]==pk] if "period" in df.columns else df
            for k, ctr in enumerate(CTR_ORDER):
                vals = [sub[(sub["country"]==ctr)&(sub["rank"]==r)]["revenue_m_usd"].mean() if len(sub[(sub["country"]==ctr)&(sub["rank"]==r)]) else 0 for r in ranks]
                bars = ax.bar(x + k*w, vals, w, label=CTR_LABEL[ctr], color=CTR_COLORS[ctr], alpha=0.88, edgecolor="white", lw=1)
                for bar, v in zip(bars, vals):
                    if v > 0.005:
                        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.01, f"${v:.1f}M", ha="center", va="bottom", fontsize=6.5, fontweight="bold", color=CTR_COLORS[ctr])
        ax.set_title(title, fontsize=11, fontweight="bold", color="#1E2761", pad=6)
        ax.set_xticks(x + w); ax.set_xticklabels(rank_labels, fontsize=9.5)
        ax.set_ylabel("월 평균 매출 (백만 USD)", fontsize=8)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4); ax.spines[["top","right"]].set_visible(False)
        ax.legend(fontsize=8, framealpha=0.9)
    fig.text(0.5, 0.005, "TOP 100 기준 · KR=한국 / JP=일본 / US=미국 · 출처: Sensor Tower", ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.03, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 22: 퍼블리셔별 순위별 월평균 매출 (iOS + AOS)
# ══════════════════════════════════════════════════════════════════════════
def chart_s22() -> bytes:
    ranks = [1, 10, 20, 50, 100]; rank_labels = ["1위","10위","20위","50위","100위"]
    # iOS
    df_ios = ios_rev[ios_rev["publisher_group"].isin(PUB_ORDER)].copy()
    df_ios["revenue_m_usd"] = pd.to_numeric(df_ios["revenue_m_usd"], errors="coerce")
    # AOS
    df_aod = aod_rev[aod_rev["publisher_group"].isin(PUB_ORDER)].copy() if len(aod_rev) and "publisher_group" in aod_rev.columns else pd.DataFrame()
    if len(df_aod): df_aod["revenue_m_usd"] = pd.to_numeric(df_aod["revenue_m_usd"], errors="coerce")

    fig, axes = plt.subplots(2, 2, figsize=(13.5, 7.5), sharey=False, gridspec_kw={"wspace":0.28,"hspace":0.45})
    fig.patch.set_facecolor("white")
    configs = [
        (axes[0,0], df_ios, "iOS — 2025 이전 (2022–2024)", "pre2025"),
        (axes[0,1], df_ios, "iOS — 2025 이후 (2025~)", "post2025"),
        (axes[1,0], df_aod, "Android — 2025 이전 (2022–2024)", "pre2025"),
        (axes[1,1], df_aod, "Android — 2025 이후 (2025~)", "post2025"),
    ]
    for ax, df, title, pk in configs:
        ax.set_facecolor("#F8FAFF"); x = np.arange(len(ranks)); w = 0.2
        if len(df) > 0:
            sub = df[df["period"]==pk] if "period" in df.columns else df
            for k, pub in enumerate(PUB_ORDER):
                vals = [sub[(sub["publisher_group"]==pub)&(sub["rank"]==r)]["revenue_m_usd"].mean() if len(sub[(sub["publisher_group"]==pub)&(sub["rank"]==r)]) else 0 for r in ranks]
                bars = ax.bar(x + k*w, vals, w, label=pub, color=PUB_COLORS[pub], alpha=0.88, edgecolor="white", lw=1)
                for bar, v in zip(bars, vals):
                    if v > 0.005:
                        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.01, f"${v:.1f}M", ha="center", va="bottom", fontsize=6, fontweight="bold", color=PUB_COLORS[pub])
        ax.set_title(title, fontsize=11, fontweight="bold", color="#1E2761", pad=6)
        ax.set_xticks(x + w*1.5); ax.set_xticklabels(rank_labels, fontsize=9.5)
        ax.set_ylabel("월 평균 매출 (백만 USD)", fontsize=8)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4); ax.spines[["top","right"]].set_visible(False)
        ax.legend(fontsize=7.5, framealpha=0.9, ncol=2)
    fig.text(0.5, 0.005, "TOP 100 기준 · 북미=미국·캐나다·영국·독일 등 · 출처: Sensor Tower", ha="center", fontsize=NOTE_FS, color="#777")
    fig.tight_layout(rect=[0, 0.03, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# PPT 빌드 — 인사이트 폰트크기 = 해석과 동일
# ══════════════════════════════════════════════════════════════════════════
SLIDE_INFO = [
    ("s15", "신규 게임 3개월 생존율 — 국가별 (KR·JP·US)",
     "매출 TOP 100 신규진입 · iOS/AOS 분리 · 2025 전/후 월평균",
     "해석: 2025년 이후 대부분 시장에서 생존율 하락. KR이 iOS 31->24%, AOS 39->22%로 가장 큰 낙폭. US AOS만 유일하게 39->41%로 상승.",
     "인사이트: 한국 시장의 경쟁 강도가 가장 높으며, 미국 AOS는 상대적으로 신규 진입 기회가 남아있는 시장."),
    ("s16", "광고 집행율 — 국가별 (3개월 생존 vs 미생존)",
     "매출 TOP 100 신규진입 · iOS 기준 · 월평균 광고집행율",
     "해석: KR·JP에서 생존 게임의 광고집행율(11~16%)이 미생존(1~6%) 대비 확연히 높음. US는 생존/미생존 격차가 적음.",
     "인사이트: 한국·일본 시장에서는 Paid UA 투자가 초기 생존과 강한 상관관계. 미국은 오가닉 유입 비중이 높아 광고 외 요인이 중요."),
    ("s17", "신규 게임 3개월 생존율 — 퍼블리셔별",
     "매출 TOP 100 신규진입 · iOS/AOS 분리 · 2025 전/후 월평균",
     "해석: 기타(핀란드·영국·이스라엘 등)가 iOS 42%, AOS 59%로 최고 생존율. 중화권 iOS 36%, 북미 iOS 29%. 한국은 AOS post2025 19%로 최하위.",
     "인사이트: 기타 그룹에 Supercell(핀란드)·King(스웨덴) 등 글로벌 강자 포함. 중화권은 물량 공세, 북미는 안정적 생존. 한국 퍼블리셔 개선 시급."),
    ("s18", "광고 집행율 — 퍼블리셔별 (3개월 생존 vs 미생존)",
     "매출 TOP 100 신규진입 · iOS 기준 · 2025 전/후 비교",
     "해석: 생존 게임의 광고집행율이 퍼블리셔 유형과 무관하게 미생존 대비 높은 경향.",
     "인사이트: Paid UA는 퍼블리셔 규모·국적과 관계없이 초기 생존의 공통 요소. 다만 광고만으로 생존이 보장되지는 않음."),
    ("s19", "차트 잔존율 M+1~M+12 — 국가별 (KR·JP·US)",
     "신규진입 게임이 M+N 시점에 매출 TOP100에 남아있는 비율 · iOS/AOS · 2025 전/후 월평균",
     "해석: pre2025에서 M+12 잔존율이 JP·US 29~32%로 양호했으나, post2025에서 전 시장 4~12%로 급락. KR이 가장 낮음(4%).",
     "인사이트: 2025년 이후 신규 게임의 1년 후 TOP100 잔존이 극히 어려워짐. 기존 상위 게임의 방어력이 강화되고 있음을 시사."),
    ("s20", "차트 잔존율 M+1~M+12 — 퍼블리셔별",
     "신규진입 게임이 M+N 시점에 매출 TOP100에 남아있는 비율 · iOS/AOS · 2025 전/후 월평균",
     "해석: 기타 그룹이 M+12 iOS 29%, AOS 27%로 최고. 일본 iOS 25%, AOS 23%. 한국 iOS 16%, AOS 15%로 최하위.",
     "인사이트: 기타(Supercell·King·Playrix 등 유럽 강자)와 일본의 장기 잔존이 압도적. 한국 퍼블리셔는 M+3 이후 급감 — 라이브서비스 강화 필요."),
    ("s21", "순위별 월평균 매출 — 국가별 (1·10·20·50·100위)",
     "iOS + AOS · TOP 100 전체 게임 · 2025 전/후 월평균",
     "해석: iOS US 1위 월매출 $56M, AOS US 1위 $29M으로 타 시장 대비 압도적. 50위 이하에서는 국가 간 격차가 $1~2M 수준으로 급감.",
     "인사이트: 미국 시장의 상위권 매출 규모가 독보적이며, TOP 10 진입이 수익성의 결정적 분기점."),
    ("s22", "순위별 월평균 매출 — 퍼블리셔별",
     "iOS + AOS · TOP 100 전체 게임 · 2025 전/후 월평균",
     "해석: 북미 퍼블리셔가 iOS 1위 매출에서 압도적. 중화권은 iOS 상위권에서 강세, AOS에서는 한국·일본과 유사.",
     "인사이트: 북미 퍼블리셔(미국·캐나다)가 매출 규모 1위. 기타 그룹에도 Supercell 등 고매출 퍼블리셔 포함."),
]

CHART_FUNCS = {
    "s15": chart_s15, "s16": chart_s16, "s17": chart_s17, "s18": chart_s18,
    "s19": chart_s19, "s20": chart_s20, "s21": chart_s21, "s22": chart_s22,
}


def main():
    print("▶ 새 PPTX 생성 중 (슬라이드 15~22)...")
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    for key, title, subtitle, interpretation, insight in SLIDE_INFO:
        print(f"  차트 생성: {title[:35]}...")
        img_bytes = CHART_FUNCS[key]()
        slide = prs.slides.add_slide(blank)

        # 상단 헤더
        from pptx.enum.shapes import MSO_SHAPE
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.05))
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61); bg.line.fill.background()

        txb = slide.shapes.add_textbox(Inches(0.35), Inches(0.08), Inches(9.3), Inches(0.55))
        p = txb.text_frame.paragraphs[0]; p.text = title
        r = p.runs[0]; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

        txb2 = slide.shapes.add_textbox(Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.30))
        p2 = txb2.text_frame.paragraphs[0]; p2.text = subtitle
        r2 = p2.runs[0]; r2.font.size = Pt(10); r2.font.color.rgb = RGBColor(0xBB,0xBB,0xBB)

        # 차트 이미지
        chart_top = 1.10; chart_h = 3.50
        slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.15), Inches(chart_top), Inches(9.70), Inches(chart_h))

        # 해석 + 인사이트 (하나의 텍스트박스, 각 라벨에 형광펜)
        note_top = chart_top + chart_h + 0.03
        txb3 = slide.shapes.add_textbox(Inches(0.25), Inches(note_top), Inches(9.50), Inches(0.62))
        tf3 = txb3.text_frame; tf3.word_wrap = True

        # 해석 줄
        p_interp = tf3.paragraphs[0]
        run_label1 = p_interp.add_run()
        run_label1.text = "해석"
        run_label1.font.size = Pt(9); run_label1.font.bold = True
        run_label1.font.color.rgb = RGBColor(0x00,0x00,0x00)
        from pptx.oxml.ns import qn
        rPr1 = run_label1._r.get_or_add_rPr()
        hl1 = rPr1.makeelement(qn('a:highlight'), {})
        srgb1 = hl1.makeelement(qn('a:srgbClr'), {'val': 'FFFF00'})
        hl1.append(srgb1); rPr1.append(hl1)
        # 해석 내용 (": " 이후)
        interp_text = interpretation.replace("해석: ", "").replace("해석:", "")
        run_body1 = p_interp.add_run()
        run_body1.text = ": " + interp_text
        run_body1.font.size = Pt(9); run_body1.font.color.rgb = RGBColor(0x33,0x33,0x33)

        # 인사이트 줄
        from pptx.util import Pt as Pt2
        p_insight = tf3.add_paragraph()
        run_label2 = p_insight.add_run()
        run_label2.text = "인사이트"
        run_label2.font.size = Pt(9); run_label2.font.bold = True
        run_label2.font.color.rgb = RGBColor(0x00,0x00,0x00)
        rPr2 = run_label2._r.get_or_add_rPr()
        hl2 = rPr2.makeelement(qn('a:highlight'), {})
        srgb2 = hl2.makeelement(qn('a:srgbClr'), {'val': 'FFFF00'})
        hl2.append(srgb2); rPr2.append(hl2)
        insight_text = insight.replace("인사이트: ", "").replace("인사이트:", "")
        run_body2 = p_insight.add_run()
        run_body2.text = ": " + insight_text
        run_body2.font.size = Pt(9); run_body2.font.color.rgb = RGBColor(0x33,0x33,0x33)

        print(f"  ✔ {key} 완료")

    outpath = BASE / "MobileGame_Slides_15_22_v4_temp.pptx"
    prs.save(str(outpath))
    print(f"\n▶ 저장 완료: {outpath.name} ({len(prs.slides)}장)")


if __name__ == "__main__":
    main()
