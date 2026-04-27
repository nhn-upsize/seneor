import sys, io
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from lxml import etree

prs = Presentation("C:/Users/NHN/Documents/sensortower_api/MobileGame_Slides_15_22_v3.pptx")
prs4 = Presentation("C:/Users/NHN/Documents/sensortower_api/MobileGame_Slides_15_22_v4_temp.pptx")
for si in range(13, 21):
    sv3 = prs.slides[si]; sv4 = prs4.slides[si-13]
    rm = [sh for sh in sv3.shapes if sh.shape_type==13 or (sh.has_text_frame and sh.top/914400>4.5) or (sh.has_text_frame and 0.5<sh.top/914400<1.1)]
    for sh in rm: sh._element.getparent().remove(sh._element)
    for sh in sv4.shapes:
        t = sh.top/914400
        if sh.shape_type==13:
            sv3.shapes.add_picture(io.BytesIO(sh.image.blob), sh.left, sh.top, sh.width, sh.height)
        elif sh.has_text_frame and (t>4.5 or (0.5<t<1.1)):
            sv3.shapes._spTree.append(etree.fromstring(etree.tostring(sh._element)))
prs.save("C:/Users/NHN/Documents/sensortower_api/MobileGame_Slides_15_22_v3.pptx")
print("v3 done!")
