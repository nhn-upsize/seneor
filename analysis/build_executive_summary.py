"""Executive Summary PPT — v3 기준 데이터로 6장 요약본"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BASE = Path("C:/Users/NHN/Documents/sensortower_api")

NAVY  = RGBColor(0x1E, 0x27, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x66, 0x66, 0x66)
TEAL  = RGBColor(0x0D, 0x94, 0x88)
RED   = RGBColor(0xE7, 0x4C, 0x3C)
LGRAY = RGBColor(0xBB, 0xBB, 0xBB)
CARD_A = RGBColor(0xF0, 0xF4, 0xFF)
CARD_B = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank = prs.slide_layouts[6]


def add_header(slide, title, subtitle=""):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.05))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.35), Inches(0.08), Inches(9.3), Inches(0.55))
    p = txb.text_frame.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        txb2 = slide.shapes.add_textbox(Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.30))
        p2 = txb2.text_frame.paragraphs[0]; p2.text = subtitle
        r2 = p2.runs[0]; r2.font.size = Pt(10); r2.font.color.rgb = LGRAY


def add_footer(slide):
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(5.35), Inches(9.4), Inches(0.22))
    p = txb.text_frame.paragraphs[0]
    p.text = "Source: Sensor Tower API  |  iOS + Android TOP 100 기준  |  2022-2026.03"
    p.runs[0].font.size = Pt(8); p.runs[0].font.color.rgb = RGBColor(0x9B, 0xA8, 0xC0)


def add_text(slide, left, top, w, h, text, size=12, bold=False, color=NAVY, align=None):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    if align: p.alignment = align
    return txb


def add_card(slide, left, top, w, h, fill_rgb):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = fill_rgb; card.line.fill.background()
    return card


def add_insight(slide, left, top, w, text):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(0.45))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    rl = p.add_run(); rl.text = "인사이트"
    rl.font.size = Pt(10); rl.font.bold = True; rl.font.color.rgb = RGBColor(0, 0, 0)
    rPr = rl._r.get_or_add_rPr()
    hl = rPr.makeelement(qn("a:highlight"), {})
    srgb = hl.makeelement(qn("a:srgbClr"), {"val": "FFFF00"})
    hl.append(srgb); rPr.append(hl)
    rb = p.add_run(); rb.text = ": " + text
    rb.font.size = Pt(10); rb.font.color.rgb = GRAY


# ═══ Slide 1: 표지 ═══
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
add_text(s, 0.8, 1.5, 8.4, 1.0, "모바일 게임 시장 변화 분석", 36, True, WHITE)
add_text(s, 0.8, 2.5, 8.4, 0.5, "Executive Summary", 20, False, TEAL)
add_text(s, 0.8, 3.2, 8.4, 0.5, "한국 / 일본 / 미국  |  2022 - 2026.Q1  |  매출 TOP 100 기준", 14, False, LGRAY)
add_text(s, 0.8, 4.8, 8.4, 0.3, "iOS + Android 분리 분석  |  Sensor Tower API  |  2026.04", 10, False, RGBColor(0x88, 0x88, 0x88))
print("  S1 표지")


# ═══ Slide 2: 핵심 메시지 ═══
s = prs.slides.add_slide(blank)
add_header(s, "핵심 요약 - 3가지 메시지", "2025년 이후 모바일 게임 시장 구조 변화")

msgs = [
    ("1", "생존 난이도 급상승",
     "2025년 이후 전 시장 생존율 15~20%p 하락\niOS KR 53->35% / JP 59->37% / US 64->46%\nAndroid도 동일 추세 - 시장 성숙기 본격 진입\n* 생존 기준: 첫 진입월+3개월 시점 TOP100 존재 여부"),
    ("2", "일본 퍼블리셔의 잔존율 압도적 우위",
     "M+12 잔존율: 일본 4.8% vs 한국 2.1% vs 중화권 2.1%\nIP 기반 충성도와 라이브서비스가 핵심 차별 요소\n한국 퍼블리셔 장기 잔존 전략 강화 시급"),
    ("3", "미국 시장의 매출 집중도 심화",
     "AOS 1위 월매출: 미국 $28M vs 한국 $20M vs 일본 $12M\nTOP 10 진입이 수익성의 결정적 분기점\n50위 이하에서는 국가 간 격차 급감"),
]
for i, (num, title, body) in enumerate(msgs):
    y = 1.25 + i * 1.40
    add_card(s, 0.3, y, 0.5, 0.5, TEAL)
    add_text(s, 0.3, y, 0.5, 0.5, num, 20, True, WHITE, PP_ALIGN.CENTER)
    add_text(s, 0.95, y, 3.5, 0.45, title, 14, True, NAVY)
    add_text(s, 0.95, y + 0.45, 8.7, 0.90, body, 10, False, GRAY)
add_footer(s)
print("  S2 핵심메시지")


# ═══ Slide 3: 생존율 ═══
s = prs.slides.add_slide(blank)
add_header(s, "신규 게임 3개월 생존율 - 전 시장 하락", "매출 TOP 100 신규진입 기준 / pre2025 vs post2025 연평균")

data = [
    ("KR 한국", "53% -> 35%  ▼18p", "47% -> 37%  ▼10p"),
    ("JP 일본", "59% -> 37%  ▼22p", "60% -> 45%  ▼15p"),
    ("US 미국", "64% -> 46%  ▼18p", "68% -> 45%  ▼23p"),
]
cols = ["국가", "iOS (pre -> post)", "AOS (pre -> post)"]
col_w = [2.5, 3.2, 3.2]
col_x = [0.3, 2.9, 6.2]

for ci in range(3):
    add_card(s, col_x[ci], 1.25, col_w[ci], 0.40, NAVY)
    add_text(s, col_x[ci] + 0.1, 1.26, col_w[ci] - 0.2, 0.35, cols[ci], 11, True, WHITE)

for ri, (ctr, ios_v, aos_v) in enumerate(data):
    y = 1.75 + ri * 0.55
    bg_c = CARD_A if ri % 2 == 0 else CARD_B
    vals = [ctr, ios_v, aos_v]
    for ci in range(3):
        add_card(s, col_x[ci], y, col_w[ci], 0.50, bg_c)
        clr = NAVY if ci == 0 else RED
        add_text(s, col_x[ci] + 0.1, y + 0.05, col_w[ci] - 0.2, 0.40, vals[ci], 12, ci == 0, clr)

add_text(s, 0.3, 3.55, 9.4, 0.30,
    "퍼블리셔별: 중화권/서구권이 상대적으로 높은 생존율 유지, 한국은 플랫폼별 편차 큼", 10, False, GRAY)
add_insight(s, 0.3, 3.95, 9.4,
    "시장 성숙기 진입으로 신규 게임 정착 난이도가 급격히 상승. Paid UA 투자가 초기 생존의 필수 조건으로 확인됨.")
add_footer(s)
print("  S3 생존율")


# ═══ Slide 4: 잔존율 ═══
s = prs.slides.add_slide(blank)
add_header(s, "월별 잔존율 M+1~M+12 - 일본 압도적 우위", "매출 TOP 100 신규진입 / iOS+AOS / WW 기준")

ret_rows = [
    ("iOS", "KR 한국", "7.9%", "4.0%", "2.6%"),
    ("iOS", "JP 일본", "9.3%", "4.8%", "3.6%"),
    ("iOS", "US 미국", "7.6%", "3.9%", "2.5%"),
    ("AOS", "KR 한국", "6.6%", "3.1%", "2.0%"),
    ("AOS", "JP 일본", "10.8%", "5.8%", "3.6%"),
    ("AOS", "US 미국", "6.9%", "3.3%", "2.1%"),
]
hdr = ["플랫폼", "국가", "M+1", "M+3", "M+12"]
hw = [1.2, 2.0, 1.8, 1.8, 1.8]
hx = [0.3, 1.6, 3.7, 5.6, 7.5]

for ci in range(5):
    add_card(s, hx[ci], 1.25, hw[ci], 0.35, NAVY)
    add_text(s, hx[ci] + 0.05, 1.26, hw[ci] - 0.1, 0.30, hdr[ci], 10, True, WHITE)

for ri, (plat, ctr, m1, m3, m12) in enumerate(ret_rows):
    y = 1.68 + ri * 0.40
    bg_c = CARD_A if ri % 2 == 0 else CARD_B
    vals = [plat, ctr, m1, m3, m12]
    for ci in range(5):
        add_card(s, hx[ci], y, hw[ci], 0.36, bg_c)
        clr = NAVY if ci < 2 else TEAL
        add_text(s, hx[ci] + 0.05, y + 0.02, hw[ci] - 0.1, 0.32, vals[ci], 11, ci < 2, clr)

add_text(s, 0.3, 4.20, 9.4, 0.30,
    "퍼블리셔별: 일본 M+12=4.8% >> 서구권 2.6% > 중화권 2.1% = 한국 2.1%", 10, True, NAVY)
add_insight(s, 0.3, 4.55, 9.4,
    "일본 퍼블리셔의 IP 기반 라이브서비스가 장기 잔존의 핵심. 한국 퍼블리셔는 M+3 이후 급감하며 개선 시급.")
add_footer(s)
print("  S4 잔존율")


# ═══ Slide 5: 매출 ═══
s = prs.slides.add_slide(blank)
add_header(s, "순위별 월평균 매출 - 미국 시장 압도적", "iOS + AOS / 매출 TOP 100 / 2025 전/후 연평균")

rev_rows = [
    ("1위", "$20M", "$12M", "$28M"),
    ("10위", "$3.4M", "$3.7M", "$7.0M"),
    ("50위", "$0.8M", "$0.8M", "$1.8M"),
    ("100위", "$0.3M", "$0.4M", "$1.0M"),
]
rhdr = ["순위", "KR 한국", "JP 일본", "US 미국"]
rw = [1.5, 2.4, 2.4, 2.4]
rx = [0.3, 1.9, 4.4, 6.9]

for ci in range(4):
    add_card(s, rx[ci], 1.25, rw[ci], 0.40, NAVY)
    add_text(s, rx[ci] + 0.1, 1.26, rw[ci] - 0.2, 0.35, rhdr[ci], 11, True, WHITE)

for ri, (rank, kr, jp, us) in enumerate(rev_rows):
    y = 1.75 + ri * 0.55
    bg_c = CARD_A if ri % 2 == 0 else CARD_B
    vals = [rank, kr, jp, us]
    for ci in range(4):
        add_card(s, rx[ci], y, rw[ci], 0.50, bg_c)
        clr = NAVY if ci == 0 else TEAL
        add_text(s, rx[ci] + 0.1, y + 0.05, rw[ci] - 0.2, 0.40, vals[ci], 13, ci == 0, clr)

add_text(s, 0.3, 4.05, 9.4, 0.30,
    "AOS 기준 / 퍼블리셔별: 중화권이 상위권 매출 집중도 가장 높음", 10, True, NAVY)
add_insight(s, 0.3, 4.40, 9.4,
    "TOP 10 진입 여부가 수익성의 결정적 분기점. 미국 시장 상위권 매출이 압도적이며 글로벌 전략의 최우선 타겟.")
add_footer(s)
print("  S5 매출")


# ═══ Slide 6: NHN 시사점 ═══
s = prs.slides.add_slide(blank)
add_header(s, "NHN 시사점 & 대응 방향", "시장 분석 기반 전략적 시사점")

items = [
    ("위기", RGBColor(0xC0, 0x39, 0x2B),
     "3개월 생존율 35~46%대로 급락 + 한국 퍼블리셔 잔존율 최하위",
     "신규 론칭 시 초기 3개월 생존 전략 필수\nM+3 이후 잔존 개선 로드맵 선행 설계"),
    ("기회", RGBColor(0xE0, 0xA8, 0x00),
     "미국 시장 유일한 성장 + AOS 진입장벽 낮음",
     "글로벌 타겟 시 미국 Android 우선 진입 전략 검토\nAction/SLG 장르 중심 포트폴리오 재편"),
    ("과제", RGBColor(0x05, 0x96, 0x69),
     "일본 대비 장기 잔존율 2배 이상 격차",
     "IP 기반 라이브서비스 강화\n커뮤니티 운영 고도화 / Whale 유저 리텐션 집중"),
]

for i, (badge, badge_clr, diagnosis, action) in enumerate(items):
    y = 1.25 + i * 1.35
    add_text(s, 0.4, y, 1.5, 0.35, badge, 14, True, badge_clr)
    add_card(s, 0.3, y + 0.40, 4.5, 0.80, CARD_A)
    add_text(s, 0.4, y + 0.40, 4.3, 0.25, "현황 진단", 9, True, RGBColor(0x99, 0x99, 0x99))
    add_text(s, 0.4, y + 0.62, 4.3, 0.55, diagnosis, 10, False, NAVY)
    add_card(s, 5.0, y + 0.40, 4.7, 0.80, RGBColor(0xE8, 0xF5, 0xE9))
    add_text(s, 5.1, y + 0.40, 4.5, 0.25, "대응 방향", 9, True, RGBColor(0x99, 0x99, 0x99))
    add_text(s, 5.1, y + 0.62, 4.5, 0.55, action, 10, False, RGBColor(0x05, 0x96, 0x69))

add_footer(s)
print("  S6 시사점")

prs.save(str(BASE / "MobileGame_Executive_Summary_2026.pptx"))
print("\nDone: MobileGame_Executive_Summary_2026.pptx (6 slides)")
