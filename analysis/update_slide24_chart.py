"""
Slide 24 차트 업데이트:
- Paid Display + Paid Search 누적 막대(stacked bar)로 변경
- 합계(광고집행율 총량)가 막대 높이로 바로 보이도록
- 각 구간에 수치 레이블 + 상단에 합계 레이블
- v3.pptx 기준으로 수정 → v3.pptx 저장
"""
import sys, io, zipfile, re, shutil
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

rcParams["font.family"]     = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE    = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX_IN = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"

COUNTRY_COLORS = {"KR": "#E74C3C", "US": "#3498DB", "JP": "#2ECC71", "CN": "#F39C12"}


def fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def make_stacked_ad_chart(df: pd.DataFrame) -> bytes:
    """
    누적 막대:
    - 하단: Paid Display (진한색)
    - 상단: Paid Search (같은 색 연하게 + 빗금)
    - 상단 레이블: 합계 %
    - 우측: Organic(자연유입) % 텍스트 표시
    WW / US 나란히 2패널
    """
    countries  = ["KR", "JP", "US", "CN"]
    xlabels    = ["KR\n한국", "JP\n일본", "US\n미국", "CN\n중국"]
    colors     = [COUNTRY_COLORS[c] for c in countries]

    df_idx = df.set_index("country")

    fig, axes = plt.subplots(1, 2, figsize=(11, 5.0))
    fig.patch.set_alpha(0)

    spec = [
        (axes[0], "전세계(WW) 기준",
         "paid_display_pct_ww", "paid_search_pct_ww"),
        (axes[1], "미국 앱스토어(US) 기준",
         "paid_display_pct_us", "paid_search_pct_us"),
    ]

    for ax, title, disp_col, srch_col in spec:
        disp = [float(df_idx.loc[c, disp_col]) if c in df_idx.index else 0 for c in countries]
        srch = [float(df_idx.loc[c, srch_col]) if c in df_idx.index else 0 for c in countries]
        total = [d + s for d, s in zip(disp, srch)]
        organic = [100 - t for t in total]
        x = np.arange(len(countries))
        w = 0.5

        # 하단: Paid Display
        bars1 = ax.bar(x, disp, w, color=colors, alpha=0.92,
                       edgecolor="white", linewidth=1.2, label="Paid Display")
        # 상단: Paid Search (같은 색 연하게 + 빗금)
        bars2 = ax.bar(x, srch, w, bottom=disp, color=colors, alpha=0.45,
                       edgecolor="white", linewidth=1.2, hatch="///", label="Paid Search")

        # Paid Display 수치 (막대 중앙)
        for bar, v in zip(bars1, disp):
            if v >= 3:
                ax.text(bar.get_x() + bar.get_width()/2,
                        bar.get_height()/2,
                        f"{v:.1f}%", ha="center", va="center",
                        fontsize=8.5, fontweight="bold", color="white")

        # Paid Search 수치 (상단 구간 중앙)
        for bar1, bar2, v in zip(bars1, bars2, srch):
            if v >= 2:
                mid = bar1.get_height() + bar2.get_height()/2
                ax.text(bar2.get_x() + bar2.get_width()/2,
                        mid, f"{v:.1f}%",
                        ha="center", va="center",
                        fontsize=8, color="white", fontweight="bold")

        # 합계 레이블 (막대 위)
        for xi, t, org in zip(x, total, organic):
            ax.text(xi, t + 0.8, f"합계 {t:.1f}%",
                    ha="center", va="bottom", fontsize=8.5,
                    fontweight="bold", color="#333")
            ax.text(xi, t + 3.8, f"(Organic {org:.1f}%)",
                    ha="center", va="bottom", fontsize=7.5,
                    color="#888")

        ax.set_title(title, fontsize=12, fontweight="bold", color="#1E2761", pad=8)
        ax.set_xticks(x)
        ax.set_xticklabels(xlabels, fontsize=10, fontweight="bold")
        ax.set_ylabel("다운로드 중 비중 (%)", fontsize=9)
        ax.set_ylim(0, max(total) * 1.55)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35)
        ax.set_facecolor("#F8FAFF")
        ax.spines[["top", "right"]].set_visible(False)

        legend_handles = [
            mpatches.Patch(facecolor="#888", alpha=0.92, edgecolor="white",
                           label="Paid Display (배너·영상·플레이어블)"),
            mpatches.Patch(facecolor="#888", alpha=0.45, edgecolor="white",
                           hatch="///", label="Paid Search (앱스토어 검색광고)"),
        ]
        ax.legend(handles=legend_handles, loc="upper right",
                  fontsize=7.5, framealpha=0.88)

    fig.tight_layout()
    return fig_to_bytes(fig)


def replace_slide24_image(pptx_path: Path, new_img: bytes) -> bool:
    """Slide 24 (index 23)의 이미지를 교체."""
    prs = Presentation(str(pptx_path))
    slide = prs.slides[23]
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            shape._element.getparent().remove(shape._element)
            slide.shapes.add_picture(io.BytesIO(new_img), left, top, width, height)
            prs.save(str(pptx_path))
            return True
    return False


def main():
    df = pd.read_csv(BASE / "ad_rate_by_country.csv")
    print("데이터:\n", df.to_string(index=False))

    print("\n▶ 누적 막대 차트 생성...")
    img = make_stacked_ad_chart(df)

    print("▶ Slide 24 이미지 교체...")
    ok = replace_slide24_image(PPTX_IN, img)
    print(f"  {'성공' if ok else '실패 — 이미지 도형 없음'}")

    if ok:
        print(f"▶ 저장 완료: {PPTX_IN.name}")
    else:
        # 이미지가 없으면 새로 추가
        prs = Presentation(str(PPTX_IN))
        slide = prs.slides[23]
        slide.shapes.add_picture(io.BytesIO(img),
            Inches(0.25), Inches(1.1), Inches(9.2), Inches(5.2))
        prs.save(str(PPTX_IN))
        print(f"▶ 이미지 새로 추가 후 저장 완료: {PPTX_IN.name}")


if __name__ == "__main__":
    main()
