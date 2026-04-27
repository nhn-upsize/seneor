"""
v3.pptx 잔존율 슬라이드 수정 (권장안):
  Slide 18  — 국가별 × 2025 전/후 제대로 된 두 패널 차트
  Slide 19  — 퍼블리셔별 × 2025 전/후 두 패널 차트
  Slide 20  — 신규 추가: 퍼블리셔별 월단위 잔존율 (M+1 ~ M+12)
  (※ Slide 17 텍스트는 이미 수정되어 있으므로 그대로 유지)
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.lines as mlines
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
from lxml import etree

rcParams["font.family"]        = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"

CTR_COLORS = {"KR":"#E74C3C", "US":"#3498DB", "JP":"#2ECC71", "CN":"#F39C12"}
CTR_LABEL  = {"KR":"KR 한국", "JP":"JP 일본", "US":"US 미국", "CN":"CN 중국"}
PUB_COLORS = {"중화권":"#F39C12","서구권":"#3498DB","일본":"#2ECC71","한국":"#E74C3C","기타":"#9B59B6"}

# ── 데이터 로드 ───────────────────────────────────────────────────────────
df2  = pd.read_csv(BASE / "slide2_v2.csv")
tags = pd.read_csv(BASE / "app_tags.csv")

# 국가별 전/후 평균 (D1/D7/D14/D30)
ctr_stats = (
    df2.groupby(["country","period"])[["d1","d7","d14","d30"]]
       .agg(["mean","count"]).round(1)
)

# 퍼블리셔별 전/후 평균
pub_stats = (
    df2.groupby(["publisher_group","period"])[["d1","d7","d14","d30"]]
       .agg(["mean","count"]).round(1)
)

# 퍼블리셔별 월단위 잔존율 (app_tags)
pub_ret = (
    tags.groupby("publisher_group")
        .agg(n=("app_id","count"),
             m1=("d30_ret_ww","mean"), m2=("d60_ret_ww","mean"),
             m3=("d90_ret_ww","mean"), m6=("d180_ret_ww","mean"),
             m12=("d365_ret_ww","mean"))
        .round(1)
)


def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def replace_picture(slide, img_bytes: bytes) -> bool:
    for shape in slide.shapes:
        if shape.shape_type == 13:
            left, top, w, h = shape.left, shape.top, shape.width, shape.height
            shape._element.getparent().remove(shape._element)
            slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, w, h)
            return True
    return False


# ══════════════════════════════════════════════════════════════════════════
# 차트 1: 국가별 × 2025 전후 (Slide 18)
# ══════════════════════════════════════════════════════════════════════════
def make_ctr_before_after() -> bytes:
    ctrs    = ["KR","JP","US"]
    x_pts   = np.array([0, 1, 2, 3])
    xlabels = ["D1", "D7", "D14", "D30"]
    periods = [("pre2025","2025 이전  (2022–2024)"), ("post2025","2025 이후  (2025~)")]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.8), sharey=True)
    fig.patch.set_alpha(0)

    for ax, (period, panel_title) in zip(axes, periods):
        ax.set_facecolor("#F8FAFF")
        handles = []
        for ctr in ctrs:
            try:
                row = ctr_stats.loc[(ctr, period)]
            except KeyError:
                continue
            vals  = [row[("d1","mean")], row[("d7","mean")],
                     row[("d14","mean")], row[("d30","mean")]]
            n_val = int(row[("d1","count")])
            color = CTR_COLORS[ctr]
            ls    = "--" if ctr == "JP" else "-"
            mk    = {"KR":"o","JP":"^","US":"s"}.get(ctr,"o")
            ax.plot(x_pts, vals, color=color, lw=2.2, ls=ls,
                    marker=mk, markersize=7)
            # 마지막 포인트에 레이블
            ax.text(x_pts[-1]+0.06, vals[-1],
                    f"{CTR_LABEL[ctr]} (n={n_val})",
                    color=color, fontsize=8, va="center")
            handles.append(mlines.Line2D([], [], color=color, lw=2.2,
                           ls=ls, marker=mk, markersize=6,
                           label=f"{CTR_LABEL[ctr]}"))

        # 포인트 값 레이블
        for ctr in ctrs:
            try:
                row = ctr_stats.loc[(ctr, period)]
            except KeyError:
                continue
            vals = [row[("d1","mean")], row[("d7","mean")],
                    row[("d14","mean")], row[("d30","mean")]]
            for xi, v in zip(x_pts, vals):
                ax.text(xi, v+0.6, f"{v:.1f}", ha="center", va="bottom",
                        fontsize=7.5, color=CTR_COLORS[ctr])

        ax.set_title(panel_title, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_xticks(x_pts)
        ax.set_xticklabels(xlabels, fontsize=11, fontweight="bold")
        ax.set_ylabel("잔존율 (%)", fontsize=9)
        ax.set_ylim(0, 60)
        ax.set_xlim(-0.3, 3.8)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(handles=handles, loc="upper right", fontsize=8.5, framealpha=0.85)

    # 수직 구분선
    fig.add_artist(plt.Line2D([0.505, 0.505], [0.05, 0.95],
                              transform=fig.transFigure,
                              color="#CBD5E1", lw=1.5, ls="--"))
    fig.text(0.5, 0.01,
             "* post2025 샘플 수가 적어 대표성 제한 있음  |  출처: Sensor Tower slide2_v2.csv",
             ha="center", fontsize=7.5, color="#888")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 차트 2: 퍼블리셔별 × 2025 전후 (Slide 19)
# ══════════════════════════════════════════════════════════════════════════
def make_pub_before_after() -> bytes:
    pubs    = ["중화권","서구권","일본","한국"]
    x_pts   = np.array([0, 1, 2, 3])
    xlabels = ["D1", "D7", "D14", "D30"]
    periods = [("pre2025","2025 이전  (2022–2024)"), ("post2025","2025 이후  (2025~)")]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.8), sharey=True)
    fig.patch.set_alpha(0)

    for ax, (period, panel_title) in zip(axes, periods):
        ax.set_facecolor("#F8FAFF")
        handles = []
        for pub in pubs:
            try:
                row = pub_stats.loc[(pub, period)]
            except KeyError:
                continue
            vals  = [row[("d1","mean")], row[("d7","mean")],
                     row[("d14","mean")], row[("d30","mean")]]
            n_val = int(row[("d1","count")])
            color = PUB_COLORS.get(pub,"#888")
            ls    = "--" if pub == "한국" else "-"
            mk    = {"중화권":"o","서구권":"s","일본":"^","한국":"D"}.get(pub,"o")
            ax.plot(x_pts, vals, color=color, lw=2.2, ls=ls,
                    marker=mk, markersize=7)
            ax.text(x_pts[-1]+0.06, vals[-1],
                    f"{pub} (n={n_val})",
                    color=color, fontsize=8, va="center")
            handles.append(mlines.Line2D([], [], color=color, lw=2.2,
                           ls=ls, marker=mk, markersize=6, label=pub))

        for pub in pubs:
            try:
                row = pub_stats.loc[(pub, period)]
            except KeyError:
                continue
            vals = [row[("d1","mean")], row[("d7","mean")],
                    row[("d14","mean")], row[("d30","mean")]]
            for xi, v in zip(x_pts, vals):
                ax.text(xi, v+0.6, f"{v:.1f}", ha="center", va="bottom",
                        fontsize=7.5, color=PUB_COLORS.get(pub,"#888"))

        ax.set_title(panel_title, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_xticks(x_pts)
        ax.set_xticklabels(xlabels, fontsize=11, fontweight="bold")
        ax.set_ylabel("잔존율 (%)", fontsize=9)
        ax.set_ylim(0, 70)
        ax.set_xlim(-0.3, 3.8)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(handles=handles, loc="upper right", fontsize=8.5, framealpha=0.85)

    fig.add_artist(plt.Line2D([0.505, 0.505], [0.05, 0.95],
                              transform=fig.transFigure,
                              color="#CBD5E1", lw=1.5, ls="--"))
    fig.text(0.5, 0.01,
             "* post2025 샘플 수 극히 적음 (일부 퍼블리셔 1~2개) — 추세 참고용  |  출처: Sensor Tower slide2_v2.csv",
             ha="center", fontsize=7.5, color="#888")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 차트 3: 퍼블리셔별 월단위 잔존율 M+1~M+12 (신규 Slide 20)
# ══════════════════════════════════════════════════════════════════════════
def make_monthly_chart() -> bytes:
    pubs   = ["중화권","서구권","일본","한국"]
    m_cols = ["m1","m2","m3","m6","m12"]
    x_pts  = np.array([0, 1, 2, 3, 4])
    xlbls  = ["M+1\n(D30)","M+2\n(D60)","M+3\n(D90)","M+6\n(D180)","M+12\n(D365)"]

    fig, (ax_line, ax_tbl) = plt.subplots(
        1, 2, figsize=(11, 4.8), gridspec_kw={"width_ratios":[1.55,1]})
    fig.patch.set_alpha(0)

    # ── 좌: 선그래프 ──
    ax_line.set_facecolor("#F8FAFF")
    for pub in pubs:
        if pub not in pub_ret.index:
            continue
        row   = pub_ret.loc[pub]
        vals  = [row[c] for c in m_cols]
        color = PUB_COLORS.get(pub,"#888")
        ls    = "--" if pub == "한국" else "-"
        mk    = {"중화권":"o","서구권":"s","일본":"^","한국":"D"}.get(pub,"o")
        note  = "*" if pub == "한국" else ""
        ax_line.plot(x_pts, vals, color=color, lw=2.2, ls=ls,
                     marker=mk, markersize=7,
                     label=f"{pub}{note} (n={int(row['n'])})")
        ax_line.text(x_pts[-1]+0.08, vals[-1], f"{vals[-1]:.1f}%",
                     color=color, fontsize=8, va="center")

    ax_line.set_xticks(x_pts)
    ax_line.set_xticklabels(xlbls, fontsize=9.5)
    ax_line.set_ylabel("잔존율 (%)", fontsize=9)
    ax_line.set_ylim(0, 28)
    ax_line.set_xlim(-0.3, 4.8)
    ax_line.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax_line.spines[["top","right"]].set_visible(False)
    ax_line.set_title("월 단위 잔존율 곡선 (WW 평균)", fontsize=11,
                      fontweight="bold", color="#1E2761", pad=6)
    ax_line.legend(loc="upper right", fontsize=8.5, framealpha=0.85)

    # ── 우: 수치 테이블 ──
    ax_tbl.axis("off")
    ax_tbl.set_facecolor("#F8FAFF")
    col_labels = ["퍼블리셔", "M+1", "M+2", "M+3", "M+6", "M+12"]
    rows_data, cell_colors = [], []
    for pub in pubs:
        if pub not in pub_ret.index:
            continue
        row   = pub_ret.loc[pub]
        note  = "*" if pub == "한국" else ""
        vals  = [f"{row[c]:.1f}%" for c in m_cols]
        rows_data.append([f"{pub}{note}"] + vals)
        bg = "#FFF8E7" if pub == "한국" else "#FFFFFF"
        cell_colors.append([bg]*6)

    tbl = ax_tbl.table(cellText=rows_data, colLabels=col_labels,
                       cellLoc="center", loc="center",
                       cellColours=cell_colors)
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9)
    tbl.scale(1, 1.6)
    for j in range(len(col_labels)):
        tbl[(0,j)].set_facecolor("#1E2761")
        tbl[(0,j)].set_text_props(color="white", fontweight="bold")
    ax_tbl.set_title("퍼블리셔별 수치 요약", fontsize=10,
                     fontweight="bold", color="#1E2761", pad=6)

    fig.text(0.01, 0.01,
             "* 한국 퍼블리셔 n=1 — 참고 수준  |  M+1≈D30, M+2≈D60, M+3≈D90, M+6≈D180, M+12≈D365  |  출처: Sensor Tower app_tags",
             fontsize=7, color="#888")
    fig.tight_layout()
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 신규 슬라이드 삽입 헬퍼
# ══════════════════════════════════════════════════════════════════════════
def add_slide_at(prs, position: int, layout_index: int = 0):
    """빈 슬라이드를 추가하고 position 위치로 이동."""
    # v3.pptx는 레이아웃이 1개(index 0)뿐
    layout_index = min(layout_index, len(prs.slide_layouts) - 1)
    layout = prs.slide_layouts[layout_index]
    slide  = prs.slides.add_slide(layout)  # 맨 끝에 추가

    # _sldIdLst에서 맨 끝 요소를 position 위치로 이동
    sl_lst = prs.slides._sldIdLst
    el = sl_lst[-1]
    sl_lst.remove(el)
    sl_lst.insert(position, el)
    return slide


def build_new_slide(prs, img_bytes: bytes, title: str, subtitle: str):
    """Slide 19 다음(index 19)에 새 슬라이드 삽입."""
    slide = add_slide_at(prs, position=19, layout_index=6)

    # 배경
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFF)

    # 제목 박스
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.1), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = False
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20); run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)

    # 부제목
    sub = slide.shapes.add_textbox(Inches(0.3), Inches(0.78), Inches(9.1), Inches(0.28))
    r = sub.text_frame.paragraphs[0].add_run()
    r.text = subtitle
    r.font.size = Pt(8.5); r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # 차트 이미지
    slide.shapes.add_picture(io.BytesIO(img_bytes),
        Inches(0.25), Inches(1.1), Inches(9.2), Inches(5.1))

    return slide


# ══════════════════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════════════════
def main():
    print("▶ 차트 생성 중...")
    img_ctr     = make_ctr_before_after()
    print("  ✔ 국가별 전/후 차트")
    img_pub     = make_pub_before_after()
    print("  ✔ 퍼블리셔별 전/후 차트")
    img_monthly = make_monthly_chart()
    print("  ✔ 퍼블리셔별 월단위 차트")

    prs = Presentation(str(PPTX))

    # Slide 18 (index 17): 국가별 × 2025 전후
    print("\n▶ Slide 18 이미지 교체 (국가별 × 2025 전후)...")
    ok = replace_picture(prs.slides[17], img_ctr)
    print(f"  {'교체 완료' if ok else '이미지 도형 없음 — 추가'}")
    if not ok:
        prs.slides[17].shapes.add_picture(
            io.BytesIO(img_ctr), Inches(0.3), Inches(1.15), Inches(9.1), Inches(5.0))

    # Slide 19 (index 18): 퍼블리셔별 × 2025 전후
    print("▶ Slide 19 이미지 교체 (퍼블리셔별 × 2025 전후)...")
    ok = replace_picture(prs.slides[18], img_pub)
    print(f"  {'교체 완료' if ok else '이미지 도형 없음 — 추가'}")
    if not ok:
        prs.slides[18].shapes.add_picture(
            io.BytesIO(img_pub), Inches(0.3), Inches(1.15), Inches(9.1), Inches(5.0))

    # Slide 20 신규 삽입 (index 19 위치에)
    print("▶ Slide 20 신규 삽입 (퍼블리셔별 월단위 잔존율 M+1~M+12)...")
    build_new_slide(
        prs, img_monthly,
        title    = "퍼블리셔별 월단위 잔존율 — M+1 ~ M+12",
        subtitle = "Sensor Tower app_tags 기준  |  WW 평균  |  M+1=D30, M+2=D60, M+3=D90, M+6=D180, M+12=D365  |  2026.01"
    )
    print("  삽입 완료")

    prs.save(str(PPTX))
    print(f"\n▶ 저장 완료: {PPTX.name}")

    # 검증
    prs2 = Presentation(str(PPTX))
    print(f"  총 {len(prs2.slides)}장")
    for i in range(17, 22):
        s = prs2.slides[i]
        t = [sh.text_frame.text.strip()[:50] for sh in s.shapes
             if sh.has_text_frame and sh.text_frame.text.strip()]
        imgs = sum(1 for sh in s.shapes if sh.shape_type == 13)
        print(f"  Slide {i+1}: {t[0] if t else '(빈)'} | 이미지:{imgs}")


if __name__ == "__main__":
    main()
