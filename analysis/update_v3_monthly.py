"""
v3.pptx 업데이트:
  Slide 17 — 연도별 D1/D7/D14/D30 테이블 → 퍼블리셔별 월단위 잔존율 차트+테이블
  Slide 18 — 국가별 잔존율 차트 → 퍼블리셔별 월단위 잔존율 선그래프
  Slide 19 — 퍼블리셔별 잔존율 차트 → 동일하되 누적 % 기반 곡선으로 교체
  Slide 23 — 퍼블리셔별 광고집행율 → 퍼블리셔별+국가별 합산 차트
  Slide 24 — 국가별 광고집행율    → 국가별+퍼블리셔별 합산 차트
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

rcParams["font.family"]        = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"

PUB_COLORS = {
    "중화권": "#F39C12",
    "서구권": "#3498DB",
    "일본":   "#2ECC71",
    "한국":   "#E74C3C",
    "기타":   "#9B59B6",
}
CTR_COLORS = {"KR": "#E74C3C", "US": "#3498DB", "JP": "#2ECC71", "CN": "#F39C12"}

# ── 데이터 로드 ───────────────────────────────────────────────────────────
df         = pd.read_csv(BASE / "app_tags.csv")
country_df = pd.read_csv(BASE / "ad_rate_by_country.csv")

# 퍼블리셔별 월 잔존율 (D30=M1, D60=M2, D90=M3, D180=M6, D365=M12)
pub_ret = (
    df.groupby("publisher_group")
      .agg(n=("app_id","count"),
           m1=("d30_ret_ww","mean"), m2=("d60_ret_ww","mean"),
           m3=("d90_ret_ww","mean"), m6=("d180_ret_ww","mean"),
           m12=("d365_ret_ww","mean"))
      .round(1)
)

# 퍼블리셔별 광고집행율
pub_ad = (
    df.groupby("publisher_group")
      .agg(disp=("paid_display_pct_ww","mean"),
           srch=("paid_search_pct_ww","mean"))
      .round(1)
)
pub_ad["total"]   = (pub_ad["disp"] + pub_ad["srch"]).round(1)
pub_ad["organic"] = (100 - pub_ad["total"]).round(1)

PUB_ORDER = ["중화권", "서구권", "일본", "한국"]
CTR_ORDER = ["KR", "JP", "US"]
CTR_LABEL = {"KR":"KR\n한국","JP":"JP\n일본","US":"US\n미국"}


def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def replace_picture(slide, img_bytes: bytes):
    """슬라이드에서 첫 번째 PICTURE 도형을 교체."""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            left, top, w, h = shape.left, shape.top, shape.width, shape.height
            shape._element.getparent().remove(shape._element)
            slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, w, h)
            return True
    return False


def add_picture_fullwidth(slide, img_bytes: bytes, top_in: float, h_in: float):
    """이미지가 없는 슬라이드에 새 이미지 추가."""
    slide.shapes.add_picture(
        io.BytesIO(img_bytes),
        Inches(0.3), Inches(top_in), Inches(9.1), Inches(h_in)
    )


# ══════════════════════════════════════════════════════════════════════════
# 1) 월단위 잔존율 차트 (퍼블리셔별 선그래프)
# ══════════════════════════════════════════════════════════════════════════
def make_monthly_retention_chart() -> bytes:
    months  = ["M+1", "M+2", "M+3", "M+6", "M+12"]
    m_cols  = ["m1", "m2", "m3", "m6", "m12"]
    x       = np.arange(len(months))

    fig, ax = plt.subplots(figsize=(10, 4.5))
    fig.patch.set_alpha(0)
    ax.set_facecolor("#F8FAFF")

    for pub in PUB_ORDER:
        if pub not in pub_ret.index:
            continue
        row    = pub_ret.loc[pub]
        vals   = [row[c] for c in m_cols]
        color  = PUB_COLORS.get(pub, "#888")
        note   = " *" if pub == "한국" else ""
        label  = f"{pub}{note} (n={int(row['n'])})"
        ls     = "--" if pub == "한국" else "-"
        marker = "o" if pub in ("중화권","서구권") else "^" if pub=="일본" else "s"
        ax.plot(x, vals, color=color, linewidth=2.2, linestyle=ls,
                marker=marker, markersize=7, label=label)
        ax.text(x[-1] + 0.08, vals[-1], f"{vals[-1]}%",
                color=color, fontsize=8.5, va="center")

    ax.set_xticks(x)
    ax.set_xticklabels(months, fontsize=11, fontweight="bold")
    ax.set_ylabel("잔존율 (%)", fontsize=10)
    ax.set_ylim(0, 28)
    ax.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_title("퍼블리셔별 월 단위 평균 잔존율 (WW 기준)", fontsize=13,
                 fontweight="bold", color="#1E2761", pad=8)
    ax.legend(loc="upper right", fontsize=9, framealpha=0.85)
    ax.text(0, -3.8, "* 한국 퍼블리셔 n=1 — 참고 수준  |  M+1≈D30, M+2≈D60, M+3≈D90, M+6≈D180, M+12≈D365",
            fontsize=7.5, color="#888", transform=ax.transData)
    fig.tight_layout()
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 2) 퍼블리셔별 월단위 잔존율 + 테이블 복합 차트 (Slide 17·18용)
# ══════════════════════════════════════════════════════════════════════════
def make_retention_with_table() -> bytes:
    months = ["M+1", "M+2", "M+3", "M+6", "M+12"]
    m_cols = ["m1", "m2", "m3", "m6", "m12"]
    x      = np.arange(len(months))

    fig, (ax_line, ax_tbl) = plt.subplots(1, 2, figsize=(11, 4.8),
                                           gridspec_kw={"width_ratios": [1.6, 1]})
    fig.patch.set_alpha(0)

    # ── 좌측: 선그래프 ──
    ax_line.set_facecolor("#F8FAFF")
    for pub in PUB_ORDER:
        if pub not in pub_ret.index:
            continue
        row   = pub_ret.loc[pub]
        vals  = [row[c] for c in m_cols]
        color = PUB_COLORS.get(pub, "#888")
        note  = " *" if pub == "한국" else ""
        label = f"{pub}{note}"
        ls    = "--" if pub == "한국" else "-"
        mk    = {"중화권":"o","서구권":"s","일본":"^","한국":"D"}.get(pub,"o")
        ax_line.plot(x, vals, color=color, lw=2.2, ls=ls,
                     marker=mk, markersize=7, label=label)
        ax_line.text(x[-1]+0.08, vals[-1], f"{vals[-1]}%",
                     color=color, fontsize=8, va="center")

    ax_line.set_xticks(x); ax_line.set_xticklabels(months, fontsize=10.5, fontweight="bold")
    ax_line.set_ylabel("잔존율 (%)", fontsize=9)
    ax_line.set_ylim(0, 28)
    ax_line.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax_line.spines[["top","right"]].set_visible(False)
    ax_line.set_title("월 단위 잔존율 곡선 (WW 평균)", fontsize=11,
                      fontweight="bold", color="#1E2761", pad=6)
    ax_line.legend(loc="upper right", fontsize=8.5, framealpha=0.85)

    # ── 우측: 수치 테이블 ──
    ax_tbl.set_facecolor("#F8FAFF")
    ax_tbl.axis("off")
    col_labels = ["퍼블리셔", "M+1", "M+2", "M+3", "M+6", "M+12"]
    rows_data  = []
    cell_colors = []
    for pub in PUB_ORDER:
        if pub not in pub_ret.index:
            continue
        row   = pub_ret.loc[pub]
        note  = "*" if pub == "한국" else ""
        vals  = [f"{row[c]:.1f}%" for c in m_cols]
        rows_data.append([f"{pub}{note}"] + vals)
        bg    = "#FFF8E7" if pub == "한국" else "#FFFFFF"
        cell_colors.append([bg] * 6)

    tbl = ax_tbl.table(
        cellText    = rows_data,
        colLabels   = col_labels,
        cellLoc     = "center",
        loc         = "center",
        cellColours = cell_colors,
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9)
    tbl.scale(1, 1.55)
    # 헤더 스타일
    for j in range(len(col_labels)):
        tbl[(0, j)].set_facecolor("#1E2761")
        tbl[(0, j)].set_text_props(color="white", fontweight="bold")
    ax_tbl.set_title("퍼블리셔별 수치 요약", fontsize=10,
                     fontweight="bold", color="#1E2761", pad=6)

    fig.text(0.01, 0.01,
             "* 한국 퍼블리셔 n=1 — 참고 수준  |  M+1=D30, M+2=D60, M+3=D90, M+6=D180, M+12=D365",
             fontsize=7, color="#888")
    fig.tight_layout()
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 3) 광고집행율 — 국가별 + 퍼블리셔별 통합 차트
# ══════════════════════════════════════════════════════════════════════════
def make_combined_ad_chart() -> bytes:
    fig, (ax_ctr, ax_pub) = plt.subplots(1, 2, figsize=(11, 4.8))
    fig.patch.set_alpha(0)

    # ── 좌: 국가별 시장 기준 ──
    ax_ctr.set_facecolor("#F8FAFF")
    ctr_df = country_df.set_index("country")
    x_ctr  = np.arange(len(CTR_ORDER))
    w      = 0.5

    disp_ctr = [float(ctr_df.loc[c, "paid_display_pct_ww"]) for c in CTR_ORDER]
    srch_ctr = [float(ctr_df.loc[c, "paid_search_pct_ww"]) for c in CTR_ORDER]
    tot_ctr  = [d+s for d,s in zip(disp_ctr, srch_ctr)]
    org_ctr  = [100-t for t in tot_ctr]
    colors_c = [CTR_COLORS[c] for c in CTR_ORDER]

    b1 = ax_ctr.bar(x_ctr, disp_ctr, w, color=colors_c, alpha=0.9,
                    edgecolor="white", lw=1.1, label="Paid Display")
    b2 = ax_ctr.bar(x_ctr, srch_ctr, w, bottom=disp_ctr, color=colors_c,
                    alpha=0.45, edgecolor="white", lw=1.1, hatch="///",
                    label="Paid Search")
    for bar, v in zip(b1, disp_ctr):
        if v >= 3:
            ax_ctr.text(bar.get_x()+bar.get_width()/2, bar.get_height()/2,
                        f"{v:.1f}%", ha="center", va="center",
                        fontsize=8, fontweight="bold", color="white")
    for bar1, bar2, v in zip(b1, b2, srch_ctr):
        if v >= 2:
            ax_ctr.text(bar2.get_x()+bar2.get_width()/2,
                        bar1.get_height()+bar2.get_height()/2,
                        f"{v:.1f}%", ha="center", va="center",
                        fontsize=7.5, color="white", fontweight="bold")
    for xi, t, org in zip(x_ctr, tot_ctr, org_ctr):
        ax_ctr.text(xi, t+0.8, f"합계 {t:.1f}%", ha="center", va="bottom",
                    fontsize=8, fontweight="bold", color="#333")
        ax_ctr.text(xi, t+3.4, f"(Organic {org:.1f}%)", ha="center", va="bottom",
                    fontsize=7, color="#888")

    ax_ctr.set_title("국가별 시장 기준 (앱스토어 TOP 25)", fontsize=11,
                     fontweight="bold", color="#1E2761", pad=6)
    ax_ctr.set_xticks(x_ctr)
    ax_ctr.set_xticklabels([CTR_LABEL[c] for c in CTR_ORDER], fontsize=9.5, fontweight="bold")
    ax_ctr.set_ylabel("다운로드 중 비중 (%)", fontsize=9)
    ax_ctr.set_ylim(0, max(tot_ctr)*1.55)
    ax_ctr.yaxis.grid(True, linestyle="--", alpha=0.35)
    ax_ctr.spines[["top","right"]].set_visible(False)
    legend_h = [
        mpatches.Patch(facecolor="#888", alpha=0.9, edgecolor="white", label="Paid Display"),
        mpatches.Patch(facecolor="#888", alpha=0.45, edgecolor="white", hatch="///", label="Paid Search"),
    ]
    ax_ctr.legend(handles=legend_h, loc="upper right", fontsize=7.5, framealpha=0.85)

    # ── 우: 퍼블리셔별 기준 ──
    ax_pub.set_facecolor("#F8FAFF")
    pubs_to_show = [p for p in PUB_ORDER if p in pub_ad.index]
    x_pub = np.arange(len(pubs_to_show))
    colors_p = [PUB_COLORS[p] for p in pubs_to_show]

    disp_p = [float(pub_ad.loc[p,"disp"]) for p in pubs_to_show]
    srch_p = [float(pub_ad.loc[p,"srch"]) for p in pubs_to_show]
    tot_p  = [float(pub_ad.loc[p,"total"]) for p in pubs_to_show]
    org_p  = [float(pub_ad.loc[p,"organic"]) for p in pubs_to_show]

    p1 = ax_pub.bar(x_pub, disp_p, w, color=colors_p, alpha=0.9,
                    edgecolor="white", lw=1.1, label="Paid Display")
    p2 = ax_pub.bar(x_pub, srch_p, w, bottom=disp_p, color=colors_p,
                    alpha=0.45, edgecolor="white", lw=1.1, hatch="///",
                    label="Paid Search")
    for bar, v in zip(p1, disp_p):
        if v >= 2:
            ax_pub.text(bar.get_x()+bar.get_width()/2, bar.get_height()/2,
                        f"{v:.1f}%", ha="center", va="center",
                        fontsize=8, fontweight="bold", color="white")
    for bar1, bar2, v in zip(p1, p2, srch_p):
        if v >= 2:
            ax_pub.text(bar2.get_x()+bar2.get_width()/2,
                        bar1.get_height()+bar2.get_height()/2,
                        f"{v:.1f}%", ha="center", va="center",
                        fontsize=7.5, color="white", fontweight="bold")
    for xi, t, org in zip(x_pub, tot_p, org_p):
        ax_pub.text(xi, t+0.8, f"합계 {t:.1f}%", ha="center", va="bottom",
                    fontsize=8, fontweight="bold", color="#333")
        ax_pub.text(xi, t+3.4, f"(Organic {org:.1f}%)", ha="center", va="bottom",
                    fontsize=7, color="#888")

    xlabels_p = [f"{p}\n(n={int(pub_ret.loc[p,'n'])})" if p in pub_ret.index else p
                 for p in pubs_to_show]
    ax_pub.set_title("퍼블리셔별 기준 (WW 평균)", fontsize=11,
                     fontweight="bold", color="#1E2761", pad=6)
    ax_pub.set_xticks(x_pub)
    ax_pub.set_xticklabels(xlabels_p, fontsize=9.5, fontweight="bold")
    ax_pub.set_ylabel("다운로드 중 비중 (%)", fontsize=9)
    ax_pub.set_ylim(0, max(tot_p)*1.55)
    ax_pub.yaxis.grid(True, linestyle="--", alpha=0.35)
    ax_pub.spines[["top","right"]].set_visible(False)
    ax_pub.legend(handles=legend_h, loc="upper right", fontsize=7.5, framealpha=0.85)

    fig.text(0.01, 0.0,
             "* 한국 퍼블리셔 n=1 — 참고 수준  |  국가별: 앱스토어 TOP 25 기준 (2026.01)  "
             "|  퍼블리셔별: app_tags WW 평균",
             fontsize=7, color="#888")
    fig.tight_layout()
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# Slide 17 — 텍스트 테이블을 월단위로 교체 + 차트 이미지 추가
# ══════════════════════════════════════════════════════════════════════════
def update_slide17(slide):
    # (A) 기존 텍스트 도형 내용 교체
    group_map = {
        "US_TH": "중화권 퍼블리셔 — 월 단위 잔존율 (WW, %)",
        "US_CH": "   M+1       M+2       M+3       M+6       M+12",
        "US_R1": f"  {pub_ret.loc['중화권','m1']:.1f}       {pub_ret.loc['중화권','m2']:.1f}       {pub_ret.loc['중화권','m3']:.1f}       {pub_ret.loc['중화권','m6']:.1f}       {pub_ret.loc['중화권','m12']:.1f}",
        "US_R2": "  (n=39 앱 평균 — 2026.01 기준)",
        "US_R3": "", "US_R4": "", "US_R5": "",

        "KR_TH": "서구권 퍼블리셔 — 월 단위 잔존율 (WW, %)",
        "KR_CH": "   M+1       M+2       M+3       M+6       M+12",
        "KR_R1": f"  {pub_ret.loc['서구권','m1']:.1f}       {pub_ret.loc['서구권','m2']:.1f}       {pub_ret.loc['서구권','m3']:.1f}       {pub_ret.loc['서구권','m6']:.1f}       {pub_ret.loc['서구권','m12']:.1f}",
        "KR_R2": "  (n=17 앱 평균 — 2026.01 기준)",
        "KR_R3": "", "KR_R4": "",

        "JP_TH": "일본 퍼블리셔 — 월 단위 잔존율 (WW, %)",
        "JP_CH": "   M+1       M+2       M+3       M+6       M+12",
        "JP_R1": f"  {pub_ret.loc['일본','m1']:.1f}       {pub_ret.loc['일본','m2']:.1f}       {pub_ret.loc['일본','m3']:.1f}       {pub_ret.loc['일본','m6']:.1f}       {pub_ret.loc['일본','m12']:.1f}",
        "JP_R2": "  (n=11 앱 평균 — 2026.01 기준)",
        "JP_R3": "", "JP_R4": "",

        "CN_TH": "한국 퍼블리셔 — 월 단위 잔존율 (WW, %) *",
        "CN_CH": "   M+1       M+2       M+3       M+6       M+12",
        "CN_R1": f"  {pub_ret.loc['한국','m1']:.1f}         {pub_ret.loc['한국','m2']:.1f}         {pub_ret.loc['한국','m3']:.1f}         {pub_ret.loc['한국','m6']:.1f}         {pub_ret.loc['한국','m12']:.1f}",
        "CN_R2": "  (n=1 — 샘플 매우 적음, 참고 수준)",
        "CN_R3": "",
    }
    insight_text = ("일본 퍼블리셔가 전 구간 최고 잔존율 (M+1 20.8% → M+12 9.2%)  |  "
                    "서구권 2위 (M+1 14.3%)  |  중화권·한국은 상대적으로 낮음  |  "
                    "* 한국 n=1 참고")

    for shape in slide.shapes:
        n = shape.name
        if n in group_map and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            if group_map[n]:
                shape.text_frame.paragraphs[0].runs[0].text = group_map[n] if shape.text_frame.paragraphs[0].runs else ""
                if not shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].add_run().text = group_map[n]
        if n == "Insight_BG" and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            para0 = tf.paragraphs[0]
            if para0.runs:
                para0.runs[0].text = insight_text
            else:
                para0.add_run().text = insight_text

        # 제목 업데이트
        if n == "Title" and shape.has_text_frame:
            tf = shape.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = "퍼블리셔별 월 단위 잔존율 (M+1 ~ M+12)"

    print("  Slide 17 텍스트 수정 완료")


# ══════════════════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════════════════
def main():
    print("▶ 차트 생성 중...")
    img_monthly    = make_monthly_retention_chart()
    img_ret_table  = make_retention_with_table()
    img_combined   = make_combined_ad_chart()
    print("  차트 3종 생성 완료")

    prs = Presentation(str(PPTX))

    # Slide 17 (index 16): 텍스트 수정 + 이미지 추가
    print("▶ Slide 17 수정...")
    update_slide17(prs.slides[16])

    # Slide 18 (index 17): 이미지 교체
    print("▶ Slide 18 이미지 교체...")
    ok = replace_picture(prs.slides[17], img_ret_table)
    if not ok:
        add_picture_fullwidth(prs.slides[17], img_ret_table, 1.15, 5.0)
    print(f"  {'교체' if ok else '추가'} 완료")

    # Slide 19 (index 18): 이미지 교체
    print("▶ Slide 19 이미지 교체...")
    ok = replace_picture(prs.slides[18], img_monthly)
    if not ok:
        add_picture_fullwidth(prs.slides[18], img_monthly, 1.15, 5.0)
    print(f"  {'교체' if ok else '추가'} 완료")

    # Slide 23 (index 22): 이미지 교체
    print("▶ Slide 23 이미지 교체...")
    ok = replace_picture(prs.slides[22], img_combined)
    if not ok:
        add_picture_fullwidth(prs.slides[22], img_combined, 1.15, 5.0)
    print(f"  {'교체' if ok else '추가'} 완료")

    # Slide 24 (index 23): 이미지 교체
    print("▶ Slide 24 이미지 교체...")
    ok = replace_picture(prs.slides[23], img_combined)
    if not ok:
        add_picture_fullwidth(prs.slides[23], img_combined, 1.15, 5.0)
    print(f"  {'교체' if ok else '추가'} 완료")

    prs.save(str(PPTX))
    print(f"\n▶ 저장 완료: {PPTX.name}")


if __name__ == "__main__":
    main()
