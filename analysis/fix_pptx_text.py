"""
저장된 PPTX 파일 내 텍스트 주석 일괄 수정
  - 서부권 → 서구권  (v3.pptx 슬라이드 본문 텍스트)
  - 서구권 = 미국 · 캐나다 · 영국 ... 터키 포함 명시
"""
import sys, io, copy
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn

BASE = Path("C:/Users/NHN/Documents/sensortower_api")

REPLACEMENTS = [
    # (검색문자열, 교체문자열)
    ("서부권", "서구권"),
]

def fix_run_text(run_el, old: str, new: str) -> int:
    """단일 <a:r> 요소 안의 텍스트를 교체. 교체 횟수 반환."""
    t_el = run_el.find(qn("a:t"))
    if t_el is None or t_el.text is None:
        return 0
    if old in t_el.text:
        t_el.text = t_el.text.replace(old, new)
        return 1
    return 0


def fix_para_text(para_el, old: str, new: str) -> int:
    """
    단락 내 런(run)들을 합쳐서 old가 포함되는지 확인 후 교체.
    단순 케이스(old가 단일 런 안에 있을 때)와
    복잡 케이스(old가 여러 런에 분산) 모두 처리.
    """
    count = 0
    runs = para_el.findall(qn("a:r"))

    # ① 각 런 단독으로 먼저 교체
    for r in runs:
        count += fix_run_text(r, old, new)

    # ② 여러 런에 걸친 경우 — 전체 텍스트 합산 후 재조합
    if count == 0:
        full = "".join(
            (r.find(qn("a:t")).text or "")
            for r in runs if r.find(qn("a:t")) is not None
        )
        if old in full:
            new_full = full.replace(old, new)
            # 첫 번째 런에 전체 텍스트 몰아넣고 나머지 런 제거
            if runs:
                t0 = runs[0].find(qn("a:t"))
                if t0 is not None:
                    t0.text = new_full
                for r in runs[1:]:
                    para_el.remove(r)
                count += 1
    return count


def fix_pptx(fname: str):
    fpath = BASE / fname
    prs   = Presentation(str(fpath))
    total = 0

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                para_el = para._p
                for old, new in REPLACEMENTS:
                    n = fix_para_text(para_el, old, new)
                    if n:
                        # 교체된 텍스트 미리보기
                        preview = "".join(
                            (r.find(qn("a:t")).text or "")
                            for r in para_el.findall(qn("a:r"))
                            if r.find(qn("a:t")) is not None
                        )
                        print(f"  Slide {slide_idx+1} [{old[:6]}→{new[:6]}]: {repr(preview[:80])}")
                        total += n

    prs.save(str(fpath))
    print(f"  → {total}곳 수정 후 저장: {fname}\n")
    return total


if __name__ == "__main__":
    print("▶ v3.pptx 텍스트 수정 중...")
    fix_pptx("MobileGame_Market_Analysis_2022-2026_v3.pptx")

    print("▶ Executive Report 텍스트 확인 중...")
    fix_pptx("MobileGame_Executive_Report_2026.pptx")

    print("완료!")
