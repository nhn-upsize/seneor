"""
Slide 21 (퍼블리셔별 매출 × 2025 전후) 차트만 재생성 — 서구권 레이블 반영
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches

rcParams["font.family"]        = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"

PUB_COLORS = {
    "중화권": "#F39C12", "서구권": "#3498DB",
    "일본":   "#2ECC71", "한국":   "#E74C3C",
}
PUB_ORDER = ["중화권", "서구권", "일본", "한국"]


def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def replace_or_add(slide, img_bytes, top_in=1.15, h_in=5.0):
    for sh in slide.shapes:
        if sh.shape_type == 13:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
            sh._element.getparent().remove(sh._element)
            slide.shapes.add_picture(io.BytesIO(img_bytes), l, t, w, h)
            return True
    slide.shapes.add_picture(io.BytesIO(img_bytes),
        Inches(0.3), Inches(top_in), Inches(9.1), Inches(h_in))
    return False


def find_slide(prs, keyword):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and keyword in sh.text_frame.text:
                return i, s
    return None, None


def make_revenue_by_publisher() -> bytes:
    df3 = pd.read_csv(BASE / "slide3_v2.csv")
    df  = df3[df3["publisher_group"].isin(PUB_ORDER)].copy()
    df["revenue_m_usd"] = pd.to_numeric(df["revenue_m_usd"], errors="coerce")

    rank_bins = [1, 10, 25, 50, 100]
    rank_lbls = ["TOP 1-10", "11-25", "26-50", "51-100"]
    df["rank_group"] = pd.cut(df["rank"], bins=rank_bins,
                              labels=rank_lbls, include_lowest=True)
    agg = (df.groupby(["publisher_group", "period", "rank_group"])["revenue_m_usd"]
             .mean().round(3).reset_index())

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.8), sharey=False)
    fig.patch.set_alpha(0)
    periods = [("pre2025",  "2025 이전  (2022–2024)"),
               ("post2025", "2025 이후  (2025~)")]

    for ax, (period, title) in zip(axes, periods):
        ax.set_facecolor("#F8FAFF")
        sub = agg[agg["period"] == period]
        x   = np.arange(len(rank_lbls))
        w   = 0.2
        for k, pub in enumerate(PUB_ORDER):
            vals = []
            for rg in rank_lbls:
                row = sub[(sub["publisher_group"] == pub) & (sub["rank_group"] == rg)]
                vals.append(float(row["revenue_m_usd"].values[0]) if len(row) else 0)
            bars = ax.bar(x + k*w, vals, w, label=pub,
                          color=PUB_COLORS[pub], alpha=0.85, edgecolor="white")
            for bar, v in zip(bars, vals):
                if v > 0.001:
                    ax.text(bar.get_x() + bar.get_width()/2,
                            bar.get_height() + 0.005,
                            f"{v:.2f}", ha="center", va="bottom",
                            fontsize=6.5, color=PUB_COLORS[pub])

        ax.set_title(title, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_xticks(x + w * 1.5)
        ax.set_xticklabels(rank_lbls, fontsize=9.5)
        ax.set_ylabel("월 평균 매출 (백만 USD)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top", "right"]].set_visible(False)
        ax.legend(fontsize=8.5, framealpha=0.85, ncol=2)

    fig.add_artist(plt.Line2D([0.505, 0.505], [0.05, 0.95],
                              transform=fig.transFigure,
                              color="#CBD5E1", lw=1.5, ls="--"))
    fig.text(0.5, 0.01,
             "* 중화권·서구권·일본·한국 퍼블리셔 기준  |  서구권 = 미국·캐나다·영국·독일·프랑스·터키 등  |  출처: Sensor Tower",
             ha="center", fontsize=7.5, color="#888")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


def main():
    print("▶ Slide 21 퍼블리셔별 매출 차트 재생성 중 (서구권 레이블)...")
    img = make_revenue_by_publisher()
    print("  ✔ 차트 생성 완료")

    prs = Presentation(str(PPTX))
    idx, slide = find_slide(prs, "순위별 월 평균 매출 — 퍼블리셔")
    if slide is None:
        print("  ⚠ Slide 21 찾기 실패")
        return

    replaced = replace_or_add(slide, img)
    print(f"  ✔ Slide {idx+1} — {'교체' if replaced else '추가'} 완료")

    prs.save(str(PPTX))
    print(f"▶ 저장 완료: {PPTX.name}")


if __name__ == "__main__":
    main()
