"""
최종 PPTX 복구:
- 백업 기준 (slide25 = 결론 × 2)
- 광고집행율 국가별 슬라이드 XML을 python-pptx로 새로 생성 → slide25.xml로 주입
- 결론 (slide14.xml) → slide26.xml로 추가
- rId32 → slide25.xml (광고집행율 국가별)
- rId33 → slide26.xml (결론)
"""
import sys, io, zipfile, re, shutil
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

rcParams["font.family"]     = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE     = Path("C:/Users/NHN/Documents/sensortower_api")
BACKUP   = BASE / "MobileGame_Market_Analysis_2022-2026_v2_backup.pptx"
UNPACKED = BASE / "pptx_unpacked3"
PPTX_OUT = BASE / "MobileGame_Market_Analysis_2022-2026_v2.pptx"
TMP_NEW  = BASE / "_new_slide_tmp.pptx"
TMP_OUT  = BASE / "_rebuild_final_tmp.pptx"

COUNTRY_COLORS = {"KR": "#E74C3C", "US": "#3498DB", "JP": "#2ECC71", "CN": "#F39C12"}


def fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def make_ad_rate_country_chart(df: pd.DataFrame) -> bytes:
    countries  = ["KR", "JP", "US", "CN"]
    xlabels    = ["KR\n한국 앱스토어", "JP\n일본 앱스토어",
                  "US\n미국 앱스토어", "CN\n중국 앱스토어"]
    colors_bar = [COUNTRY_COLORS[c] for c in countries]
    df_idx = df.set_index("country")

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.8))
    fig.patch.set_alpha(0)

    spec = [
        (axes[0], "전세계(WW) 기준 광고 설치 비중",
         "paid_display_pct_ww", "paid_search_pct_ww"),
        (axes[1], "미국 앱스토어(US) 기준 광고 설치 비중",
         "paid_display_pct_us", "paid_search_pct_us"),
    ]
    for ax, title, disp_col, srch_col in spec:
        disp_vals = [float(df_idx.loc[c, disp_col]) if c in df_idx.index else 0 for c in countries]
        srch_vals = [float(df_idx.loc[c, srch_col]) if c in df_idx.index else 0 for c in countries]
        x = np.arange(len(countries)); w = 0.38
        bars1 = ax.bar(x - w/2, disp_vals, w, color=colors_bar,
                       edgecolor="white", linewidth=1.1, alpha=0.92)
        bars2 = ax.bar(x + w/2, srch_vals, w, color=colors_bar,
                       edgecolor="white", linewidth=1.1, alpha=0.5, hatch="///")
        for bar, v in zip(bars1, disp_vals):
            if v > 0:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                        f"{v:.1f}%", ha="center", va="bottom", fontsize=9, fontweight="bold")
        for bar, v in zip(bars2, srch_vals):
            if v > 0:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                        f"{v:.1f}%", ha="center", va="bottom", fontsize=9, fontweight="bold")
        ax.set_title(title, fontsize=11, fontweight="bold", color="#1E2761", pad=8)
        ax.set_xticks(x); ax.set_xticklabels(xlabels, fontsize=9)
        ax.set_ylabel("다운로드 중 광고 설치 비중 (%)", fontsize=9)
        ymax = max([v for v in disp_vals + srch_vals if v] or [10])
        ax.set_ylim(0, ymax * 1.35)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.set_facecolor("#F8FAFF")
        ax.spines[["top", "right"]].set_visible(False)
        legend_handles = [
            mpatches.Patch(facecolor="#888", edgecolor="white",
                           label="Paid Display — 배너·영상·플레이어블 광고"),
            mpatches.Patch(facecolor="#888", edgecolor="white", hatch="///",
                           alpha=0.5, label="Paid Search — 앱스토어 검색광고"),
        ]
        ax.legend(handles=legend_handles, loc="upper left", fontsize=7.5, framealpha=0.85)
    fig.tight_layout()
    return fig_to_bytes(fig)


def build_ad_slide_xml() -> tuple[bytes, bytes]:
    """광고집행율 국가별 슬라이드를 임시 PPTX로 만들어 XML 추출."""
    df_country = pd.read_csv(BASE / "ad_rate_by_country.csv")
    img_bytes  = make_ad_rate_country_chart(df_country)

    # 임시 PPTX에 슬라이드 1개 생성
    tmp_prs = Presentation()
    blank   = tmp_prs.slide_layouts[6]  # blank layout
    slide   = tmp_prs.slides.add_slide(blank)

    # 배경
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFF)

    # 제목 박스
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.1), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "광고 집행율 & Paid Install 비중 — 국가별 (시장 기준)"
    run.font.size = Pt(20); run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)

    # 서브타이틀
    sub = slide.shapes.add_textbox(Inches(0.3), Inches(0.78), Inches(9.1), Inches(0.28))
    r = sub.text_frame.paragraphs[0].add_run()
    r.text = ("각 나라 앱스토어 TOP 25 게임 기준 (2026년 1월)  |  "
              "해당 시장 인기 게임들의 광고 의존도")
    r.font.size = Pt(8.5); r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # 차트
    slide.shapes.add_picture(io.BytesIO(img_bytes),
        Inches(0.25), Inches(1.1), Inches(9.2), Inches(5.2))

    # 하단 노트
    note = slide.shapes.add_textbox(Inches(0.25), Inches(6.4), Inches(9.2), Inches(0.55))
    note.text_frame.word_wrap = True
    nr = note.text_frame.paragraphs[0].add_run()
    nr.text = ("WW = 해당 게임의 전세계 다운로드 중 광고 설치 비중  |  "
               "US = 미국 앱스토어 기준 광고 설치 비중\n"
               "Paid Display: 배너·영상·플레이어블 광고  |  "
               "Paid Search: Apple Search Ads 등 앱스토어 검색광고  |  "
               "Source: Sensor Tower API (2026-01)")
    nr.font.size = Pt(7.8); nr.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # 임시 파일로 저장
    tmp_prs.save(str(TMP_NEW))

    # 임시 PPTX에서 slide1.xml과 slide1.xml.rels 추출
    with zipfile.ZipFile(TMP_NEW, "r") as zf:
        slide_xml  = zf.read("ppt/slides/slide1.xml")
        slide_rels = zf.read("ppt/slides/_rels/slide1.xml.rels")
        # rels의 이미지 참조를 확인
        print(f"  새 슬라이드 rels 참조:\n{slide_rels.decode('utf-8')[:600]}")

        # 이미지 파일들도 추출 (slide XML에서 이미지 참조 번호 찾기)
        img_files = [(n, zf.read(n)) for n in zf.namelist()
                     if n.startswith("ppt/media/")]
        print(f"  미디어 파일 수: {len(img_files)}")

    TMP_NEW.unlink(missing_ok=True)
    return slide_xml, slide_rels, img_files


def main():
    print("▶ 광고집행율 국가별 슬라이드 XML 생성 중...")
    slide_xml, slide_rels, img_files = build_ad_slide_xml()

    # 결론 슬라이드 (pptx_unpacked3/slide14.xml)
    conclusion_xml  = (UNPACKED / "ppt/slides/slide14.xml").read_bytes()
    conclusion_rels = (UNPACKED / "ppt/slides/_rels/slide14.xml.rels").read_bytes()
    print(f"  결론 슬라이드: {len(conclusion_xml)} bytes")

    # 기존 미디어 파일 번호 확인 (겹치지 않게)
    with zipfile.ZipFile(BACKUP, "r") as zf:
        existing_media = [n for n in zf.namelist() if n.startswith("ppt/media/")]
        print(f"  기존 미디어 파일: {len(existing_media)}개")
    # 새 이미지에 새 번호 부여
    media_offset = len(existing_media) + 1

    # rels XML에서 이미지 참조 경로 업데이트
    rels_text = slide_rels.decode("utf-8")
    new_media_map = {}  # 원래 이름 → 새 이름
    for i, (orig_path, img_data) in enumerate(img_files):
        orig_name = orig_path.split("/")[-1]          # image1.png
        ext = orig_name.rsplit(".", 1)[-1]
        new_name = f"image{media_offset + i}.{ext}"
        new_path = f"ppt/media/{new_name}"
        new_media_map[orig_path] = (new_path, img_data)
        # rels 내 ../media/image1.png → ../media/imageN.png
        rels_text = rels_text.replace(f"../media/{orig_name}",
                                       f"../media/{new_name}")
    slide_rels_updated = rels_text.encode("utf-8")

    print(f"\n▶ PPTX 재빌드: 백업 기반")
    with zipfile.ZipFile(BACKUP, "r") as src, \
         zipfile.ZipFile(TMP_OUT, "w", zipfile.ZIP_DEFLATED) as dst:

        written = set()
        for item in src.infolist():
            name = item.filename

            if name in written:
                print(f"  [스킵 중복] {name}")
                continue

            # slide25.xml → 광고집행율 국가별 내용으로 교체
            if name == "ppt/slides/slide25.xml":
                print(f"  [교체] {name} → 광고집행율 국가별")
                dst.writestr(item, slide_xml)
                written.add(name)
                continue

            if name == "ppt/slides/_rels/slide25.xml.rels":
                print(f"  [교체] {name} → 광고집행율 국가별 rels")
                dst.writestr(item, slide_rels_updated)
                written.add(name)
                continue

            # presentation.xml.rels: rId33 → slide26.xml (결론)
            if name == "ppt/_rels/presentation.xml.rels":
                text = src.read(name).decode("utf-8")
                text = re.sub(
                    r'(Id="rId33"[^>]+Target=")slides/slide25\.xml(")',
                    r'\1slides/slide26.xml\2', text)
                print(f"  [수정] {name}: rId33 → slide26.xml")
                dst.writestr(item, text.encode("utf-8"))
                written.add(name)
                continue

            # [Content_Types].xml: slide26.xml + 새 이미지 타입 추가
            if name == "[Content_Types].xml":
                text = src.read(name).decode("utf-8")
                insert_before = "</Types>"
                additions = ""
                if "slide26.xml" not in text:
                    additions += ('  <Override PartName="/ppt/slides/slide26.xml" '
                                  'ContentType="application/vnd.openxmlformats-officedocument'
                                  '.presentationml.slide+xml"/>\n')
                for new_path, _ in new_media_map.values():
                    fname = new_path.split("/")[-1]
                    if fname not in text:
                        ext = fname.rsplit(".", 1)[-1].lower()
                        ctype = "image/png" if ext == "png" else f"image/{ext}"
                        additions += (f'  <Override PartName="/{new_path}" '
                                      f'ContentType="{ctype}"/>\n')
                if additions:
                    text = text.replace(insert_before, additions + insert_before)
                    print(f"  [수정] {name}")
                dst.writestr(item, text.encode("utf-8"))
                written.add(name)
                continue

            dst.writestr(item, src.read(name))
            written.add(name)

        # slide26.xml (결론) 추가
        dst.writestr("ppt/slides/slide26.xml", conclusion_xml)
        dst.writestr("ppt/slides/_rels/slide26.xml.rels", conclusion_rels)
        print("  [추가] ppt/slides/slide26.xml (결론)")

        # 새 미디어 파일 추가
        for orig_path, (new_path, img_data) in new_media_map.items():
            dst.writestr(new_path, img_data)
            print(f"  [추가] {new_path}")

    shutil.move(str(TMP_OUT), str(PPTX_OUT))
    print(f"\n▶ 저장 완료: {PPTX_OUT.name}")

    # 검증
    prs = Presentation(str(PPTX_OUT))
    print(f"  총 {len(prs.slides)}개 슬라이드")
    for i, s in enumerate(prs.slides):
        texts = [sh.text_frame.text.strip()[:55] for sh in s.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        marker = " ◀" if i >= 22 else ""
        print(f"  Slide {i+1}: {texts[0] if texts else '(빈 슬라이드)'}{marker}")


if __name__ == "__main__":
    main()
