"""
iOS vs Android(AOS) 비교 차트 생성 및 v3.pptx 업데이트

업데이트 대상 슬라이드 (제목으로 찾기):
  - 퍼블리셔별 월단위 잔존율 M+1~M+12  → iOS/AOS 나란히 비교
  - 광고 집행율 & Paid Install 비중 — 퍼블리셔 출신별  → iOS/AOS 비교
  - 광고 집행율 & Paid Install — 국가별(KR·JP·US)     → iOS/AOS 비교
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.lines as mlines
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

rcParams["font.family"]        = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"

PUB_COLORS  = {"중화권":"#F39C12","서구권":"#3498DB","일본":"#2ECC71","한국":"#E74C3C","기타":"#9B59B6"}
CTR_COLORS  = {"KR":"#E74C3C","JP":"#2ECC71","US":"#3498DB"}
CTR_LABEL   = {"KR":"KR 한국","JP":"JP 일본","US":"US 미국"}
PUB_ORDER   = ["중화권","서구권","일본","한국"]
CTR_ORDER   = ["KR","JP","US"]

# ── 데이터 로드 ───────────────────────────────────────────────────────────
ios_tags = pd.read_csv(BASE / "app_tags.csv")
aod_tags = pd.read_csv(BASE / "app_tags_android.csv")
ios_ctr  = pd.read_csv(BASE / "ad_rate_by_country.csv")
aod_ctr  = pd.read_csv(BASE / "ad_rate_by_country_android.csv")

# 퍼블리셔별 월잔존율 집계
def pub_ret(df):
    return df.groupby("publisher_group").agg(
        n=("app_id","count"),
        m1=("d30_ret_ww","mean"), m2=("d60_ret_ww","mean"),
        m3=("d90_ret_ww","mean"), m6=("d180_ret_ww","mean"),
        m12=("d365_ret_ww","mean")
    ).round(1)

ios_pub_ret = pub_ret(ios_tags)
aod_pub_ret = pub_ret(aod_tags)

# 퍼블리셔별 광고집행율 집계
def pub_ad(df):
    return df.groupby("publisher_group").agg(
        n=("app_id","count"),
        disp=("paid_display_pct_ww","mean"),
        srch=("paid_search_pct_ww","mean")
    ).round(1)

ios_pub_ad = pub_ad(ios_tags)
aod_pub_ad = pub_ad(aod_tags)


def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def replace_or_add(slide, img_bytes, top_in=1.15, h_in=5.0):
    for sh in slide.shapes:
        if sh.shape_type == 13:
            l,t,w,h = sh.left, sh.top, sh.width, sh.height
            sh._element.getparent().remove(sh._element)
            slide.shapes.add_picture(io.BytesIO(img_bytes), l, t, w, h)
            return True
    slide.shapes.add_picture(io.BytesIO(img_bytes),
        Inches(0.3), Inches(top_in), Inches(9.1), Inches(h_in))
    return False


def find_slide(prs, keyword):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and keyword in sh.text_frame.text:
                return i, s
    return None, None


# ══════════════════════════════════════════════════════════════════════════
# 차트 1: 퍼블리셔별 월단위 잔존율 iOS vs AOS
# ══════════════════════════════════════════════════════════════════════════
def make_retention_ios_aod() -> bytes:
    months = ["M+1","M+2","M+3","M+6","M+12"]
    mcols  = ["m1","m2","m3","m6","m12"]
    x      = np.arange(len(months))

    fig, axes = plt.subplots(1, 2, figsize=(12, 5.0),
                              gridspec_kw={"width_ratios":[1,1]})
    fig.patch.set_alpha(0)

    for ax, (platform, pret, marker_style) in zip(
            axes,
            [("iOS", ios_pub_ret, "-"), ("Android", aod_pub_ret, "--")]):
        ax.set_facecolor("#F8FAFF")
        ax.set_title(f"📱 {platform} — 퍼블리셔별 월단위 잔존율",
                     fontsize=11, fontweight="bold", color="#1E2761", pad=7)
        handles = []
        for pub in PUB_ORDER:
            if pub not in pret.index:
                continue
            row   = pret.loc[pub]
            vals  = [row[c] if not pd.isna(row[c]) else 0 for c in mcols]
            color = PUB_COLORS[pub]
            mk    = {"중화권":"o","서구권":"s","일본":"^","한국":"D"}.get(pub,"o")
            n     = int(row["n"])
            ax.plot(x, vals, color=color, lw=2.2, ls=marker_style,
                    marker=mk, markersize=7, label=f"{pub} (n={n})")
            ax.text(x[-1]+0.08, vals[-1], f"{vals[-1]:.1f}%",
                    color=color, fontsize=8, va="center")
            handles.append(mlines.Line2D([],[],color=color,lw=2.2,
                           ls=marker_style,marker=mk,markersize=6,
                           label=f"{pub} (n={n})"))
            for xi, v in zip(x, vals):
                if v > 0:
                    ax.text(xi, v+0.4, f"{v:.1f}", ha="center", va="bottom",
                            fontsize=6.5, color=color)

        ax.set_xticks(x)
        ax.set_xticklabels(months, fontsize=10.5, fontweight="bold")
        ax.set_ylabel("잔존율 (%)", fontsize=9)
        ax.set_ylim(0, 28)
        ax.set_xlim(-0.3, 4.9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(handles=handles, loc="upper right", fontsize=8, framealpha=0.85)

    # 구분선
    fig.add_artist(plt.Line2D([0.5,0.5],[0.05,0.95],
                              transform=fig.transFigure,
                              color="#CBD5E1",lw=1.5,ls="--"))
    fig.text(0.5, 0.01,
             "iOS: App Store TOP 25 기준 (n=74)  |  Android: Google Play TOP 25 기준 (n=74)  |"
             "  출처: Sensor Tower comparison_attributes WW 기준",
             ha="center", fontsize=7.5, color="#888")
    fig.tight_layout(rect=[0,0.04,1,1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 차트 2: 광고집행율 통합 — 국가별 + 퍼블리셔별, iOS vs AOS 2패널
# ══════════════════════════════════════════════════════════════════════════
def make_ad_ios_aod() -> bytes:
    """iOS / AOS 각각 국가별+퍼블리셔별 광고집행율 — 2×2 레이아웃"""
    fig, axes = plt.subplots(2, 2, figsize=(12, 9))
    fig.patch.set_alpha(0)

    legend_h = [
        mpatches.Patch(facecolor="#888",alpha=0.9,edgecolor="white",label="Paid Display"),
        mpatches.Patch(facecolor="#888",alpha=0.45,edgecolor="white",hatch="///",label="Paid Search"),
    ]

    configs = [
        (axes[0,0], "iOS — 국가별 시장 기준",    ios_ctr,    "country",   CTR_ORDER, CTR_LABEL,  CTR_COLORS),
        (axes[0,1], "iOS — 퍼블리셔별 기준",      ios_pub_ad, "publisher", PUB_ORDER, None,        PUB_COLORS),
        (axes[1,0], "Android — 국가별 시장 기준", aod_ctr,    "country",   CTR_ORDER, CTR_LABEL,  CTR_COLORS),
        (axes[1,1], "Android — 퍼블리셔별 기준",  aod_pub_ad, "publisher", PUB_ORDER, None,        PUB_COLORS),
    ]

    for ax, title, data, dtype, order, label_map, color_map in configs:
        ax.set_facecolor("#F8FAFF")
        ax.set_title(title, fontsize=10.5, fontweight="bold", color="#1E2761", pad=5)

        if dtype == "country":
            df = data[data["country"].isin(order)].set_index("country")
            keys  = [k for k in order if k in df.index]
            disps = [float(df.loc[k,"paid_display_pct_ww"]) for k in keys]
            srchs = [float(df.loc[k,"paid_search_pct_ww"])  for k in keys]
            xlbls = [label_map[k] for k in keys]
            clrs  = [color_map[k] for k in keys]
        else:
            keys  = [p for p in order if p in data.index]
            disps = [float(data.loc[p,"disp"]) for p in keys]
            srchs = [float(data.loc[p,"srch"])  for p in keys]
            ns    = [int(data.loc[p,"n"])       for p in keys]
            xlbls = [f"{p}\n(n={n})" for p, n in zip(keys, ns)]
            clrs  = [color_map[p] for p in keys]

        tots = [d+s for d,s in zip(disps,srchs)]
        orgs = [100-t for t in tots]
        xp   = np.arange(len(keys))
        w    = 0.5

        ax.bar(xp, disps, w, color=clrs, alpha=0.9, edgecolor="white", lw=1.1)
        ax.bar(xp, srchs, w, bottom=disps, color=clrs, alpha=0.45,
               edgecolor="white", lw=1.1, hatch="///")

        for xi, d, s, t, org in zip(xp, disps, srchs, tots, orgs):
            if d >= 3:
                ax.text(xi, d/2, f"{d:.1f}%", ha="center", va="center",
                        fontsize=8, fontweight="bold", color="white")
            if s >= 2:
                ax.text(xi, d+s/2, f"{s:.1f}%", ha="center", va="center",
                        fontsize=7.5, fontweight="bold", color="white")
            ax.text(xi, t+0.8, f"{t:.1f}%", ha="center", va="bottom",
                    fontsize=8.5, fontweight="bold", color="#333")
            ax.text(xi, t+3.5, f"Organic\n{org:.1f}%", ha="center", va="bottom",
                    fontsize=7, color="#888")

        ax.set_xticks(xp)
        ax.set_xticklabels(xlbls, fontsize=9.5, fontweight="bold")
        ax.set_ylabel("다운로드 중 비중 (%)", fontsize=8)
        ax.set_ylim(0, max(tots)*1.7 if tots else 80)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35)
        ax.spines[["top","right"]].set_visible(False)
        ax.legend(handles=legend_h, loc="upper right", fontsize=7, framealpha=0.85)

    fig.text(0.5, 0.01,
             "iOS: App Store TOP 25 (2026.01, CN 제외)  |  Android: Google Play TOP 25 (2026.01)  |"
             "  중화권=중국·대만·홍콩  |  서구권=미국·유럽  |  출처: Sensor Tower",
             ha="center", fontsize=7.5, color="#888")
    fig.suptitle("광고집행율 — iOS vs Android 비교", fontsize=13,
                 fontweight="bold", color="#1E2761", y=1.0)
    fig.tight_layout(rect=[0,0.03,1,0.99])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════════════════
def main():
    print("▶ iOS vs AOS 비교 차트 생성 중...")
    img_ret = make_retention_ios_aod()
    print("  ✔ 퍼블리셔별 월단위 잔존율 iOS vs AOS")
    img_ad  = make_ad_ios_aod()
    print("  ✔ 광고집행율 iOS vs AOS (2×2)")

    prs = Presentation(str(PPTX))
    print(f"\n현재 슬라이드 수: {len(prs.slides)}장")

    updates = [
        ("퍼블리셔별 월단위 잔존율",   img_ret, "Slide 20 (퍼블리셔별 월단위 잔존율 iOS vs AOS)"),
        ("Paid Install 비중 — 퍼블리셔", img_ad,  "Slide 24 (광고집행율 퍼블리셔별 iOS vs AOS)"),
        ("국가별(KR",                    img_ad,  "Slide 25 (광고집행율 국가별 iOS vs AOS)"),
    ]

    for keyword, img, label in updates:
        idx, slide = find_slide(prs, keyword)
        if slide is None:
            print(f"  ⚠ [{label}] 슬라이드 찾기 실패 (keyword='{keyword}')")
            continue
        replaced = replace_or_add(slide, img)
        print(f"  ✔ [{label}] Slide {idx+1} — {'교체' if replaced else '추가'} 완료")

    prs.save(str(PPTX))
    print(f"\n▶ 저장 완료: {PPTX.name}")


if __name__ == "__main__":
    main()
