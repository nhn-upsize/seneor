"""
v3.pptx 종합 수정:
  1. 중복 슬라이드 20/21 제거 (둘 중 하나 삭제)
  2. 슬라이드 23 (국가별 매출 × 2025 전후) 이미지 복구 — KR/JP/US 전용
  3. 슬라이드 24 (퍼블리셔별 매출 × 2025 전후) 이미지 복구
  4. 슬라이드 25 (퍼블리셔별 광고집행율) — 통합 차트 (CN 제외)
  5. 슬라이드 26 (국가별 광고집행율) — 통합 차트 (CN 제외)

  ※ 슬라이드는 제목(title text)으로 식별 (index 하드코딩 금지)
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
from copy import deepcopy
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

rcParams["font.family"]        = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"

# ── 색상 설정 ─────────────────────────────────────────────────────────────
CTR_COLORS = {"KR": "#E74C3C", "JP": "#2ECC71", "US": "#3498DB"}
CTR_LABEL  = {"KR": "KR 한국", "JP": "JP 일본", "US": "US 미국"}
PUB_COLORS = {
    "중화권": "#F39C12", "서구권": "#3498DB",
    "일본":   "#2ECC71", "한국":   "#E74C3C", "기타": "#9B59B6",
}
PUB_ORDER  = ["중화권", "서구권", "일본", "한국"]
CTR_ORDER  = ["KR", "JP", "US"]

# ── 데이터 로드 ───────────────────────────────────────────────────────────
df3       = pd.read_csv(BASE / "slide3_v2.csv")
df_tags   = pd.read_csv(BASE / "app_tags.csv")
df_ctr    = pd.read_csv(BASE / "ad_rate_by_country.csv")

# 광고집행율 (퍼블리셔별)
pub_ad = (
    df_tags.groupby("publisher_group")
    .agg(disp=("paid_display_pct_ww", "mean"),
         srch=("paid_search_pct_ww", "mean"),
         n=("app_id", "count"))
    .round(1)
)
pub_ad["total"]   = (pub_ad["disp"] + pub_ad["srch"]).round(1)
pub_ad["organic"] = (100 - pub_ad["total"]).round(1)

pub_ret = (
    df_tags.groupby("publisher_group")
    .agg(n=("app_id", "count"))
)


def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def replace_or_add_picture(slide, img_bytes: bytes, top_in=1.15, h_in=5.0):
    for shape in slide.shapes:
        if shape.shape_type == 13:
            left, top, w, h = shape.left, shape.top, shape.width, shape.height
            shape._element.getparent().remove(shape._element)
            slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, w, h)
            return True
    slide.shapes.add_picture(
        io.BytesIO(img_bytes),
        Inches(0.3), Inches(top_in), Inches(9.1), Inches(h_in)
    )
    return False


def find_slide_by_title(prs, keyword: str):
    """제목에 keyword가 포함된 슬라이드와 index 반환."""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if keyword in txt:
                    return i, slide
    return None, None


# ══════════════════════════════════════════════════════════════════════════
# 차트 1: 국가별 매출 × 2025 전후 (Slide 23 복구)
# ══════════════════════════════════════════════════════════════════════════
def make_revenue_by_country() -> bytes:
    """slide3_v2.csv → KR/JP/US × pre2025/post2025 월 평균 매출 비교"""
    df = df3[df3["country"].isin(CTR_ORDER)].copy()
    df["revenue_m_usd"] = pd.to_numeric(df["revenue_m_usd"], errors="coerce")

    # 순위구간 × period × country 로 평균 매출
    rank_bins  = [1, 10, 25, 50, 100]
    rank_lbls  = ["TOP 1-10", "11-25", "26-50", "51-100"]
    df["rank_group"] = pd.cut(df["rank"], bins=rank_bins,
                              labels=rank_lbls, include_lowest=True)
    agg = (df.groupby(["country", "period", "rank_group"])["revenue_m_usd"]
             .mean().round(3).reset_index())

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.8), sharey=False)
    fig.patch.set_alpha(0)
    periods = [("pre2025", "2025 이전  (2022–2024)"),
               ("post2025", "2025 이후  (2025~)")]

    for ax, (period, title) in zip(axes, periods):
        ax.set_facecolor("#F8FAFF")
        sub = agg[agg["period"] == period]
        x   = np.arange(len(rank_lbls))
        w   = 0.25
        for k, ctr in enumerate(CTR_ORDER):
            vals = []
            for rg in rank_lbls:
                row = sub[(sub["country"] == ctr) & (sub["rank_group"] == rg)]
                vals.append(float(row["revenue_m_usd"].values[0]) if len(row) else 0)
            bars = ax.bar(x + k*w, vals, w, label=CTR_LABEL[ctr],
                          color=CTR_COLORS[ctr], alpha=0.85, edgecolor="white")
            for bar, v in zip(bars, vals):
                if v > 0.001:
                    ax.text(bar.get_x() + bar.get_width()/2,
                            bar.get_height() + 0.005,
                            f"{v:.2f}", ha="center", va="bottom",
                            fontsize=7, color=CTR_COLORS[ctr])

        ax.set_title(title, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_xticks(x + w)
        ax.set_xticklabels(rank_lbls, fontsize=9.5)
        ax.set_ylabel("월 평균 매출 (백만 USD)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top", "right"]].set_visible(False)
        ax.legend(fontsize=8.5, framealpha=0.85)

    fig.add_artist(plt.Line2D([0.505, 0.505], [0.05, 0.95],
                              transform=fig.transFigure,
                              color="#CBD5E1", lw=1.5, ls="--"))
    fig.text(0.5, 0.01,
             "* KR·JP·US 3개 시장 기준  |  CN 제외  |  출처: Sensor Tower slide3_v2.csv",
             ha="center", fontsize=7.5, color="#888")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 차트 2: 퍼블리셔별 매출 × 2025 전후 (Slide 24 복구)
# ══════════════════════════════════════════════════════════════════════════
def make_revenue_by_publisher() -> bytes:
    """slide3_v2.csv → 퍼블리셔별 × pre2025/post2025 매출 비교"""
    df = df3[df3["publisher_group"].isin(PUB_ORDER)].copy()
    df["revenue_m_usd"] = pd.to_numeric(df["revenue_m_usd"], errors="coerce")

    rank_bins = [1, 10, 25, 50, 100]
    rank_lbls = ["TOP 1-10", "11-25", "26-50", "51-100"]
    df["rank_group"] = pd.cut(df["rank"], bins=rank_bins,
                              labels=rank_lbls, include_lowest=True)
    agg = (df.groupby(["publisher_group", "period", "rank_group"])["revenue_m_usd"]
             .mean().round(3).reset_index())

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.8), sharey=False)
    fig.patch.set_alpha(0)
    periods = [("pre2025", "2025 이전  (2022–2024)"),
               ("post2025", "2025 이후  (2025~)")]

    for ax, (period, title) in zip(axes, periods):
        ax.set_facecolor("#F8FAFF")
        sub = agg[agg["period"] == period]
        x   = np.arange(len(rank_lbls))
        w   = 0.2
        for k, pub in enumerate(PUB_ORDER):
            vals = []
            for rg in rank_lbls:
                row = sub[(sub["publisher_group"] == pub) & (sub["rank_group"] == rg)]
                vals.append(float(row["revenue_m_usd"].values[0]) if len(row) else 0)
            bars = ax.bar(x + k*w, vals, w, label=pub,
                          color=PUB_COLORS[pub], alpha=0.85, edgecolor="white")
            for bar, v in zip(bars, vals):
                if v > 0.001:
                    ax.text(bar.get_x() + bar.get_width()/2,
                            bar.get_height() + 0.005,
                            f"{v:.2f}", ha="center", va="bottom",
                            fontsize=6.5, color=PUB_COLORS[pub])

        ax.set_title(title, fontsize=12, fontweight="bold",
                     color="#1E2761", pad=8)
        ax.set_xticks(x + w * 1.5)
        ax.set_xticklabels(rank_lbls, fontsize=9.5)
        ax.set_ylabel("월 평균 매출 (백만 USD)", fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.spines[["top", "right"]].set_visible(False)
        ax.legend(fontsize=8, framealpha=0.85, ncol=2)

    fig.add_artist(plt.Line2D([0.505, 0.505], [0.05, 0.95],
                              transform=fig.transFigure,
                              color="#CBD5E1", lw=1.5, ls="--"))
    fig.text(0.5, 0.01,
             "* 중화권·서구권·일본·한국 퍼블리셔 기준  |  출처: Sensor Tower slide3_v2.csv",
             ha="center", fontsize=7.5, color="#888")
    fig.tight_layout(rect=[0, 0.04, 1, 1])
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 차트 3: 광고집행율 통합 (국가별 + 퍼블리셔별) — CN 제외
# ══════════════════════════════════════════════════════════════════════════
def make_combined_ad_chart() -> bytes:
    fig, (ax_ctr, ax_pub) = plt.subplots(1, 2, figsize=(11, 4.8))
    fig.patch.set_alpha(0)
    legend_h = [
        mpatches.Patch(facecolor="#888", alpha=0.9, edgecolor="white", label="Paid Display"),
        mpatches.Patch(facecolor="#888", alpha=0.45, edgecolor="white",
                       hatch="///", label="Paid Search"),
    ]

    # ── 좌: 국가별 (KR/JP/US) ──
    ax_ctr.set_facecolor("#F8FAFF")
    ctr_df  = df_ctr[df_ctr["country"].isin(CTR_ORDER)].set_index("country")
    x_ctr   = np.arange(len(CTR_ORDER))
    w       = 0.55

    disp_c = [float(ctr_df.loc[c, "paid_display_pct_ww"]) for c in CTR_ORDER]
    srch_c = [float(ctr_df.loc[c, "paid_search_pct_ww"])  for c in CTR_ORDER]
    tot_c  = [d+s for d, s in zip(disp_c, srch_c)]
    org_c  = [100-t for t in tot_c]
    clrs_c = [CTR_COLORS[c] for c in CTR_ORDER]

    ax_ctr.bar(x_ctr, disp_c, w, color=clrs_c, alpha=0.9,
               edgecolor="white", lw=1.1)
    ax_ctr.bar(x_ctr, srch_c, w, bottom=disp_c, color=clrs_c,
               alpha=0.45, edgecolor="white", lw=1.1, hatch="///")
    for xi, d, s, t, org in zip(x_ctr, disp_c, srch_c, tot_c, org_c):
        if d >= 3:
            ax_ctr.text(xi, d/2, f"{d:.1f}%", ha="center", va="center",
                        fontsize=9, fontweight="bold", color="white")
        if s >= 2:
            ax_ctr.text(xi, d+s/2, f"{s:.1f}%", ha="center", va="center",
                        fontsize=8, fontweight="bold", color="white")
        ax_ctr.text(xi, t+0.8, f"합계 {t:.1f}%", ha="center", va="bottom",
                    fontsize=8.5, fontweight="bold", color="#333")
        ax_ctr.text(xi, t+3.8, f"(Organic {org:.1f}%)", ha="center", va="bottom",
                    fontsize=7.5, color="#888")

    ax_ctr.set_title("국가별 시장 기준 (앱스토어 TOP 25)", fontsize=11,
                     fontweight="bold", color="#1E2761", pad=6)
    xtkls = [f"{CTR_LABEL[c]}" for c in CTR_ORDER]
    ax_ctr.set_xticks(x_ctr)
    ax_ctr.set_xticklabels(xtkls, fontsize=10.5, fontweight="bold")
    ax_ctr.set_ylabel("다운로드 중 비중 (%)", fontsize=9)
    ax_ctr.set_ylim(0, max(tot_c)*1.6)
    ax_ctr.yaxis.grid(True, linestyle="--", alpha=0.35)
    ax_ctr.spines[["top", "right"]].set_visible(False)
    ax_ctr.legend(handles=legend_h, loc="upper right", fontsize=7.5, framealpha=0.85)

    # ── 우: 퍼블리셔별 ──
    ax_pub.set_facecolor("#F8FAFF")
    pubs   = [p for p in PUB_ORDER if p in pub_ad.index]
    x_pub  = np.arange(len(pubs))
    clrs_p = [PUB_COLORS[p] for p in pubs]

    disp_p = [float(pub_ad.loc[p, "disp"])    for p in pubs]
    srch_p = [float(pub_ad.loc[p, "srch"])    for p in pubs]
    tot_p  = [float(pub_ad.loc[p, "total"])   for p in pubs]
    org_p  = [float(pub_ad.loc[p, "organic"]) for p in pubs]

    ax_pub.bar(x_pub, disp_p, w, color=clrs_p, alpha=0.9,
               edgecolor="white", lw=1.1)
    ax_pub.bar(x_pub, srch_p, w, bottom=disp_p, color=clrs_p,
               alpha=0.45, edgecolor="white", lw=1.1, hatch="///")
    for xi, d, s, t, org in zip(x_pub, disp_p, srch_p, tot_p, org_p):
        if d >= 2:
            ax_pub.text(xi, d/2, f"{d:.1f}%", ha="center", va="center",
                        fontsize=9, fontweight="bold", color="white")
        if s >= 2:
            ax_pub.text(xi, d+s/2, f"{s:.1f}%", ha="center", va="center",
                        fontsize=8, fontweight="bold", color="white")
        ax_pub.text(xi, t+0.8, f"합계 {t:.1f}%", ha="center", va="bottom",
                    fontsize=8.5, fontweight="bold", color="#333")
        ax_pub.text(xi, t+3.8, f"(Organic {org:.1f}%)", ha="center", va="bottom",
                    fontsize=7.5, color="#888")

    n_lbls = [f"{p}\n(n={int(pub_ad.loc[p,'n'])})" for p in pubs]
    ax_pub.set_title("퍼블리셔별 기준 (WW 평균)", fontsize=11,
                     fontweight="bold", color="#1E2761", pad=6)
    ax_pub.set_xticks(x_pub)
    ax_pub.set_xticklabels(n_lbls, fontsize=10, fontweight="bold")
    ax_pub.set_ylabel("다운로드 중 비중 (%)", fontsize=9)
    ax_pub.set_ylim(0, max(tot_p)*1.6)
    ax_pub.yaxis.grid(True, linestyle="--", alpha=0.35)
    ax_pub.spines[["top", "right"]].set_visible(False)
    ax_pub.legend(handles=legend_h, loc="upper right", fontsize=7.5, framealpha=0.85)

    fig.text(0.01, 0.0,
             "* 한국 n=1 — 참고 수준  |  국가별: TOP 25 기준 (2026.01, CN 제외)  "
             "|  퍼블리셔별: app_tags WW 평균  "
             "|  중화권=중국·대만·홍콩  |  서구권=미국·유럽",
             fontsize=7, color="#888")
    fig.tight_layout()
    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 중복 슬라이드 삭제 헬퍼
# ══════════════════════════════════════════════════════════════════════════
def remove_slide(prs, idx: int):
    """index 위치 슬라이드 삭제."""
    xml_slides = prs.slides._sldIdLst
    slide_el   = xml_slides[idx]
    xml_slides.remove(slide_el)
    # 실제 파트 삭제
    slide = prs.slides[idx] if idx < len(prs.slides) else None
    # prs.slides._sldIdLst 기준으로 제거된 것으로 충분 (python-pptx 저장 시 정리됨)


# ══════════════════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════════════════
def main():
    print("▶ 차트 생성 중...")
    img_rev_ctr = make_revenue_by_country()
    print("  ✔ 국가별 매출 × 2025 전후 차트")
    img_rev_pub = make_revenue_by_publisher()
    print("  ✔ 퍼블리셔별 매출 × 2025 전후 차트")
    img_ad      = make_combined_ad_chart()
    print("  ✔ 통합 광고집행율 차트 (CN 제외)")

    prs = Presentation(str(PPTX))

    print(f"\n현재 슬라이드 수: {len(prs.slides)}장")
    for i, s in enumerate(prs.slides):
        txts = [sh.text_frame.text.strip()[:45] for sh in s.shapes
                if sh.has_text_frame and sh.text_frame.text.strip()]
        imgs = sum(1 for sh in s.shapes if sh.shape_type == 13)
        print(f"  [{i+1:2d}] {txts[0][:40] if txts else '(빈)'} | 이미지:{imgs}")

    # ── 1) 중복 슬라이드 제거 (두 번째 "퍼블리셔별 월단위 잔존율") ──
    DUP_KEYWORD = "퍼블리셔별 월단위 잔존율"
    dup_indices = []
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and DUP_KEYWORD in sh.text_frame.text:
                dup_indices.append(i)
                break
    print(f"\n▶ 중복 슬라이드 인덱스: {[x+1 for x in dup_indices]}")
    if len(dup_indices) >= 2:
        del_idx = dup_indices[1]   # 두 번째 삭제
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(xml_slides[del_idx])
        print(f"  슬라이드 {del_idx+1} 삭제 완료")

    # ── 슬라이드 목록 재확인 ──
    print(f"\n중복 제거 후: {len(prs.slides)}장")

    # ── 2) 슬라이드를 제목으로 찾아 이미지 교체 ──
    targets = [
        ("순위별 월 평균 매출 — 국가별",         img_rev_ctr, "Slide 22 (국가별 매출 × 2025 전후)"),
        ("순위별 월 평균 매출 — 퍼블리셔",       img_rev_pub, "Slide 23 (퍼블리셔별 매출 × 2025 전후)"),
        ("Paid Install 비중 — 퍼블리셔 출신별",  img_ad,      "Slide 24 (퍼블리셔별 광고집행율)"),
        ("국가별(KR",                           img_ad,      "Slide 25 (국가별·퍼블리셔별 광고집행율)"),
    ]

    for keyword, img, label in targets:
        idx, slide = find_slide_by_title(prs, keyword)
        if slide is None:
            print(f"  ⚠  [{label}] 슬라이드를 찾지 못함 (keyword='{keyword}')")
            continue
        replaced = replace_or_add_picture(slide, img)
        print(f"  ✔  [{label}] 슬라이드 {idx+1} — {'교체' if replaced else '추가'} 완료")

    prs.save(str(PPTX))
    print(f"\n▶ 저장 완료: {PPTX.name}")

    # ── 검증 ──
    prs2 = Presentation(str(PPTX))
    print(f"\n최종 슬라이드 수: {len(prs2.slides)}장")
    for i, s in enumerate(prs2.slides):
        txts = [sh.text_frame.text.strip()[:45] for sh in s.shapes
                if sh.has_text_frame and sh.text_frame.text.strip()]
        imgs = sum(1 for sh in s.shapes if sh.shape_type == 13)
        print(f"  [{i+1:2d}] {txts[0][:40] if txts else '(빈)'} | 이미지:{imgs}")


if __name__ == "__main__":
    main()
