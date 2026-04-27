"""
Slide 24 수정: 진짜 국가별(시장 기준) 광고집행율로 교체
- ad_rate_by_country.csv (각 나라 앱스토어 TOP 25 기준) 사용
- 기존 Slide 24 삭제 후 올바른 데이터로 재삽입
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

rcParams["font.family"]     = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE    = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX_IN = BASE / "MobileGame_Market_Analysis_2022-2026_v2.pptx"

COUNTRY_COLORS = {"KR": "#E74C3C", "US": "#3498DB", "JP": "#2ECC71", "CN": "#F39C12"}
GRAY = RGBColor(0x6B, 0x72, 0x80)


def fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def add_title_box(slide, text: str):
    txBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.1), Inches(9.1), Inches(0.65))
    tf = txBox.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(20); run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)


def insert_slide_at(prs, slide, position: int):
    xml_slides = prs.slides._sldIdLst
    last = xml_slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(position, last)


def delete_slide(prs, index: int):
    xml_slides = prs.slides._sldIdLst
    slide_elem = xml_slides[index]
    xml_slides.remove(slide_elem)


def make_ad_rate_country_chart(df: pd.DataFrame) -> bytes:
    """
    국가별(시장 기준) 광고집행율.
    x축: KR / JP / US / CN (앱스토어 시장 기준)
    WW vs US 나란히 비교.
    """
    countries  = ["KR", "JP", "US", "CN"]
    xlabels    = ["KR\n한국 앱스토어", "JP\n일본 앱스토어",
                  "US\n미국 앱스토어", "CN\n중국 앱스토어"]
    colors_bar = [COUNTRY_COLORS[c] for c in countries]

    # df: country, paid_display_pct_ww, paid_search_pct_ww, paid_display_pct_us, paid_search_pct_us
    df_idx = df.set_index("country")

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.8))
    fig.patch.set_alpha(0)

    spec = [
        (axes[0], "전세계(WW) 기준 광고 설치 비중",
         "paid_display_pct_ww", "paid_search_pct_ww"),
        (axes[1], "미국 앱스토어(US) 기준 광고 설치 비중",
         "paid_display_pct_us", "paid_search_pct_us"),
    ]

    for ax, title, disp_col, srch_col in spec:
        disp_vals = [df_idx.loc[c, disp_col] if c in df_idx.index else 0 for c in countries]
        srch_vals = [df_idx.loc[c, srch_col] if c in df_idx.index else 0 for c in countries]

        x = np.arange(len(countries))
        w = 0.38

        bars1 = ax.bar(x - w/2, disp_vals, w, label="Paid Display",
                       color=colors_bar, edgecolor="white", linewidth=1.1, alpha=0.92)
        bars2 = ax.bar(x + w/2, srch_vals, w, label="Paid Search",
                       color=colors_bar, edgecolor="white", linewidth=1.1, alpha=0.5,
                       hatch="///")

        for bar, v in zip(bars1, disp_vals):
            if v and v > 0:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                        f"{v:.1f}%", ha="center", va="bottom", fontsize=9, fontweight="bold")
        for bar, v in zip(bars2, srch_vals):
            if v and v > 0:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                        f"{v:.1f}%", ha="center", va="bottom", fontsize=9, fontweight="bold")

        ax.set_title(title, fontsize=11, fontweight="bold", color="#1E2761", pad=8)
        ax.set_xticks(x)
        ax.set_xticklabels(xlabels, fontsize=9)
        ax.set_ylabel("다운로드 중 광고 설치 비중 (%)", fontsize=9)
        ymax = max([v for v in disp_vals + srch_vals if v] or [10])
        ax.set_ylim(0, ymax * 1.35)
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.set_facecolor("#F8FAFF")
        ax.spines[["top", "right"]].set_visible(False)

        legend_handles = [
            mpatches.Patch(facecolor="#888", edgecolor="white",
                           label="Paid Display — 배너·영상·플레이어블 광고"),
            mpatches.Patch(facecolor="#888", edgecolor="white", hatch="///",
                           alpha=0.5, label="Paid Search — 앱스토어 검색광고"),
        ]
        ax.legend(handles=legend_handles, loc="upper left", fontsize=7.5, framealpha=0.85)

    fig.tight_layout()
    return fig_to_bytes(fig)


def main():
    # 수집된 국가별 시장 데이터 로드
    csv_path = BASE / "ad_rate_by_country.csv"
    if not csv_path.exists():
        print("오류: ad_rate_by_country.csv 없음. collect_ad_by_country.py를 먼저 실행하세요.")
        return
    df_country = pd.read_csv(csv_path)
    print(f"국가별 ad rate 데이터:\n{df_country.to_string(index=False)}\n")

    print(f"▶ PPT 로드: {PPTX_IN.name}")
    prs = Presentation(str(PPTX_IN))
    print(f"  총 {len(prs.slides)}개 슬라이드")

    # 현재 Slide 24 (index 23) 확인 후 삭제
    slide24_title = ""
    for sh in prs.slides[23].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            slide24_title = sh.text_frame.text.strip()[:60]
            break
    print(f"  현재 Slide 24: '{slide24_title}'")
    delete_slide(prs, 23)
    print(f"  Slide 24 삭제  (현재 {len(prs.slides)}개)")

    # 새 Slide 24 생성
    print("  [Slide 24 재생성] 국가별(시장 기준) 광고집행율...")
    img_ad = make_ad_rate_country_chart(df_country)

    blank_layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(blank_layout)

    # 배경
    bg = new_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFF)

    # 제목
    add_title_box(new_slide, "광고 집행율 & Paid Install 비중 — 국가별 (시장 기준)")

    # 서브타이틀
    sub_tb = new_slide.shapes.add_textbox(
        Inches(0.3), Inches(0.78), Inches(9.1), Inches(0.28))
    sub_run = sub_tb.text_frame.paragraphs[0].add_run()
    sub_run.text = (
        "각 나라 앱스토어 TOP 25 게임 기준 (2026년 1월)  |  "
        "퍼블리셔 출신 불문, 해당 시장에서 인기 있는 게임들의 광고 의존도"
    )
    sub_run.font.size = Pt(8.5)
    sub_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # 차트
    new_slide.shapes.add_picture(io.BytesIO(img_ad),
        Inches(0.25), Inches(1.1), Inches(9.2), Inches(5.2))

    # 하단 노트
    note_tb = new_slide.shapes.add_textbox(
        Inches(0.25), Inches(6.4), Inches(9.2), Inches(0.55))
    note_tf = note_tb.text_frame; note_tf.word_wrap = True
    note_run = note_tf.paragraphs[0].add_run()
    note_run.text = (
        "WW = 해당 게임의 전세계 다운로드 중 광고 설치 비중  |  "
        "US = 미국 앱스토어 기준 광고 설치 비중\n"
        "Paid Display: 배너·영상·플레이어블 광고  |  "
        "Paid Search: Apple Search Ads 등 앱스토어 검색광고  |  Source: Sensor Tower API (2026-01)"
    )
    note_run.font.size = Pt(7.8)
    note_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # index 23에 삽입 (결론 슬라이드 바로 앞)
    insert_slide_at(prs, new_slide, 23)

    prs.save(str(PPTX_IN))
    print(f"\n▶ 저장 완료: {PPTX_IN.name}")
    print(f"  총 {len(prs.slides)}개 슬라이드")
    print("  Slide 23: 광고집행율 퍼블리셔 출신국별 (한국/일본/중화권/서구권)")
    print("  Slide 24: 광고집행율 국가별 시장 기준 (KR/JP/US/CN) — 2026-01 TOP 25")
    print("  Slide 25: 결론")


if __name__ == "__main__":
    main()
