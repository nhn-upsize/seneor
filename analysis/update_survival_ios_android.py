"""
신규 게임 3개월 생존율 iOS vs Android 비교 차트 생성 및 v3.pptx 업데이트

업데이트 대상 슬라이드 (제목으로 찾기):
  - 신규 게임 3개월 생존율  →  국가별 (iOS vs AOS 나란히)
  - 신규 게임 3개월 생존율  →  퍼블리셔별 (iOS vs AOS 나란히)
"""
import sys, io
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

rcParams["font.family"]        = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False

BASE = Path("C:/Users/NHN/Documents/sensortower_api")
PPTX = BASE / "MobileGame_Market_Analysis_2022-2026_v3.pptx"

CTR_COLORS  = {"KR":"#E74C3C", "JP":"#2ECC71", "US":"#3498DB"}
CTR_LABEL   = {"KR":"KR 한국", "JP":"JP 일본", "US":"US 미국"}
CTR_ORDER   = ["KR", "JP", "US"]

PUB_COLORS  = {"중화권":"#F39C12","서구권":"#3498DB","일본":"#2ECC71","한국":"#E74C3C","기타":"#9B59B6"}
PUB_ORDER   = ["일본","서구권","중화권","한국"]


# ── 데이터 로드 ───────────────────────────────────────────────────────────
ios_df = pd.read_csv(BASE / "slide1_v2.csv")
aod_df = pd.read_csv(BASE / "slide1_v2_android.csv")

# CN 제외 (국가별 차트)
ios_ctr = ios_df[ios_df["country"].isin(CTR_ORDER)]
aod_ctr = aod_df[aod_df["country"].isin(CTR_ORDER)]


def survival_by_country(df):
    """국가별 × period 생존율 DataFrame 반환"""
    g = df.groupby(["country","period"])["survived_3m"].agg(["sum","count"])
    g["rate"] = (g["sum"] / g["count"] * 100).round(1)
    return g.reset_index()


def survival_by_pub(df):
    """퍼블리셔별 × period 생존율 DataFrame 반환"""
    g = df.groupby(["publisher_group","period"])["survived_3m"].agg(["sum","count"])
    g["rate"] = (g["sum"] / g["count"] * 100).round(1)
    return g.reset_index()


def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="none", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def replace_or_add(slide, img_bytes, top_in=1.15, h_in=5.2):
    for sh in slide.shapes:
        if sh.shape_type == 13:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
            sh._element.getparent().remove(sh._element)
            slide.shapes.add_picture(io.BytesIO(img_bytes), l, t, w, h)
            return True
    slide.shapes.add_picture(io.BytesIO(img_bytes),
        Inches(0.3), Inches(top_in), Inches(9.1), Inches(h_in))
    return False


def find_slide(prs, keyword):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and keyword in sh.text_frame.text:
                return i, s
    return None, None


# ══════════════════════════════════════════════════════════════════════════
# 차트 1: 국가별 신규 게임 3개월 생존율 — iOS vs AOS 비교 (2×3 패널)
# ══════════════════════════════════════════════════════════════════════════
def make_survival_country_chart() -> bytes:
    """iOS / AOS × KR/JP/US × pre2025/post2025 → 2×3 그리드"""

    ios_s = survival_by_country(ios_ctr)
    aod_s = survival_by_country(aod_ctr)

    periods    = ["pre2025", "post2025"]
    per_labels = {"pre2025": "2022~2024", "post2025": "2025~"}

    fig, axes = plt.subplots(2, 3, figsize=(13, 7.5),
                             gridspec_kw={"hspace":0.5, "wspace":0.35})
    fig.patch.set_alpha(0)
    fig.suptitle("신규 게임 3개월 생존율 — 국가별 iOS vs Android 비교",
                 fontsize=13, fontweight="bold", color="#1E2761", y=1.01)

    platform_data = [("iOS",     ios_s,  "#2980B9", "solid"),
                     ("Android", aod_s,  "#27AE60", "dashed")]

    for row_idx, (plat, df_s, base_color, ls) in enumerate(platform_data):
        for col_idx, country in enumerate(CTR_ORDER):
            ax = axes[row_idx, col_idx]
            ax.set_facecolor("#F8FAFF")
            ax.set_title(f"📱 {plat}\n{CTR_LABEL[country]}",
                         fontsize=10, fontweight="bold",
                         color="#1E2761", pad=4)

            bars_rate  = []
            bars_n     = []
            xlbls      = []
            clrs       = []

            sub = df_s[df_s["country"] == country]
            for per in periods:
                row = sub[sub["period"] == per]
                if len(row) == 0:
                    bars_rate.append(0)
                    bars_n.append(0)
                else:
                    bars_rate.append(float(row["rate"].values[0]))
                    bars_n.append(int(row["count"].values[0]))
                clrs.append(CTR_COLORS[country])

            # n수치를 x축 레이블에 포함
            for per, n in zip(periods, bars_n):
                xlbls.append(f"{per_labels[per]}\n(n={n})")

            xp = np.arange(len(periods))
            w  = 0.5
            ax.bar(xp, bars_rate, w, color=clrs,
                   alpha=0.8 if plat == "iOS" else 0.55,
                   edgecolor="white", lw=1.2,
                   hatch="" if plat == "iOS" else "///")

            for xi, rate in zip(xp, bars_rate):
                if rate > 0:
                    ax.text(xi, rate + 1.5, f"{rate:.0f}%",
                            ha="center", va="bottom",
                            fontsize=11, fontweight="bold",
                            color=CTR_COLORS[country])

            ax.set_xticks(xp)
            ax.set_xticklabels(xlbls, fontsize=9)
            ax.set_ylim(0, 105)
            ax.set_ylabel("생존율 (%)", fontsize=8)
            ax.yaxis.grid(True, linestyle="--", alpha=0.35)
            ax.spines[["top","right"]].set_visible(False)

    # 범례 (상단)
    legend_patches = [
        mpatches.Patch(facecolor="#2980B9", alpha=0.8, label="iOS"),
        mpatches.Patch(facecolor="#27AE60", alpha=0.55, hatch="///", label="Android"),
    ]
    fig.legend(handles=legend_patches, loc="upper right",
               bbox_to_anchor=(0.98, 1.0), fontsize=9, framealpha=0.85)

    fig.text(0.5, -0.02,
             "iOS: App Store TOP 100 KR/JP/US  |  Android: Google Play TOP 100 KR/JP/US  |"
             "  신규진입 기준: 첫 차트 등장월 = 출시월  |  출처: Sensor Tower",
             ha="center", fontsize=7.5, color="#888")

    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 차트 2: 퍼블리셔별 신규 게임 3개월 생존율 — iOS vs AOS 비교 (2×2 패널)
# ══════════════════════════════════════════════════════════════════════════
def make_survival_pub_chart() -> bytes:
    """iOS / AOS × pre2025/post2025 → 2패널, 각 패널=퍼블리셔별 막대"""

    ios_s = survival_by_pub(ios_df)
    aod_s = survival_by_pub(aod_df)

    periods    = ["pre2025", "post2025"]
    per_labels = {"pre2025": "2022~2024", "post2025": "2025~"}

    fig, axes = plt.subplots(2, 2, figsize=(12, 8),
                             gridspec_kw={"hspace":0.55, "wspace":0.35})
    fig.patch.set_alpha(0)
    fig.suptitle("신규 게임 3개월 생존율 — 퍼블리셔별 iOS vs Android 비교",
                 fontsize=13, fontweight="bold", color="#1E2761", y=1.01)

    configs = [
        (axes[0,0], "iOS — 2022~2024 (pre2025)",    ios_s, "pre2025"),
        (axes[0,1], "iOS — 2025~ (post2025)",        ios_s, "post2025"),
        (axes[1,0], "Android — 2022~2024 (pre2025)", aod_s, "pre2025"),
        (axes[1,1], "Android — 2025~ (post2025)",    aod_s, "post2025"),
    ]

    for ax, title, df_s, period in configs:
        ax.set_facecolor("#F8FAFF")
        ax.set_title(title, fontsize=10, fontweight="bold", color="#1E2761", pad=5)

        sub = df_s[df_s["period"] == period]
        pubs   = [p for p in PUB_ORDER if p in sub["publisher_group"].values]
        rates  = []
        ns     = []
        clrs   = []

        for pub in pubs:
            row = sub[sub["publisher_group"] == pub]
            if len(row) == 0:
                rates.append(0)
                ns.append(0)
            else:
                rates.append(float(row["rate"].values[0]))
                ns.append(int(row["count"].values[0]))
            clrs.append(PUB_COLORS.get(pub, "#9B59B6"))

        # n수치를 x축 레이블에 포함
        xlbls_pub = [f"{p}\n(n={n})" for p, n in zip(pubs, ns)]

        xp = np.arange(len(pubs))
        w  = 0.5
        is_ios = "iOS" in title

        ax.bar(xp, rates, w, color=clrs,
               alpha=0.85 if is_ios else 0.55,
               edgecolor="white", lw=1.2,
               hatch="" if is_ios else "///")

        for xi, rate, clr in zip(xp, rates, clrs):
            if rate > 0:
                ax.text(xi, rate + 1.5, f"{rate:.0f}%",
                        ha="center", va="bottom",
                        fontsize=11, fontweight="bold",
                        color=clr)

        ax.set_xticks(xp)
        ax.set_xticklabels(xlbls_pub, fontsize=9, fontweight="bold")
        ax.set_ylim(0, 105)
        ax.set_ylabel("생존율 (%)", fontsize=8)
        ax.yaxis.grid(True, linestyle="--", alpha=0.35)
        ax.spines[["top","right"]].set_visible(False)

    # 범례
    legend_patches = [
        mpatches.Patch(facecolor="#888", alpha=0.85, label="iOS"),
        mpatches.Patch(facecolor="#888", alpha=0.55, hatch="///", label="Android"),
    ]
    fig.legend(handles=legend_patches, loc="upper right",
               bbox_to_anchor=(0.98, 1.0), fontsize=9, framealpha=0.85)

    fig.text(0.5, -0.02,
             "iOS: App Store TOP 100 KR/JP/US  |  Android: Google Play TOP 100 KR/JP/US  |"
             "  신규진입 기준: 첫 차트 등장월 = 출시월  |  중화권=CN·TW·HK, 서구권=US·EU 등  |  출처: Sensor Tower",
             ha="center", fontsize=7.5, color="#888")

    return fig_bytes(fig)


# ══════════════════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════════════════
def main():
    print("▶ 신규 게임 3개월 생존율 iOS vs AOS 차트 생성 중...")
    img_ctr = make_survival_country_chart()
    print("  ✔ 국가별 생존율 차트 (2×3 패널)")
    img_pub = make_survival_pub_chart()
    print("  ✔ 퍼블리셔별 생존율 차트 (2×2 패널)")

    prs = Presentation(str(PPTX))
    print(f"\n현재 슬라이드 수: {len(prs.slides)}장")

    # 슬라이드 검색 + 업데이트
    updates = [
        ("생존율 — 국가별",         img_ctr, "생존율 국가별 iOS vs AOS"),
        ("생존율 — 퍼블리셔",       img_pub, "생존율 퍼블리셔별 iOS vs AOS"),
    ]

    # 보조 키워드로 확장 검색
    keyword_fallbacks = {
        "생존율 — 국가별":    ["3개월 생존율 — 국가별", "생존율 — 국가별 ×", "신규 게임 3개월 생존율"],
        "생존율 — 퍼블리셔": ["3개월 생존율 — 퍼블리셔", "생존율 — 퍼블리셔 별", "생존율 — 퍼블리셔별"],
    }

    for primary_kw, img, label in updates:
        idx, slide = find_slide(prs, primary_kw)
        if slide is None:
            for fallback in keyword_fallbacks.get(primary_kw, []):
                idx, slide = find_slide(prs, fallback)
                if slide:
                    print(f"  ℹ fallback keyword 사용: '{fallback}'")
                    break
        if slide is None:
            print(f"  ⚠ [{label}] 슬라이드 찾기 실패 (keyword='{primary_kw}')")
            continue
        replaced = replace_or_add(slide, img)
        print(f"  ✔ [{label}] Slide {idx+1} — {'교체' if replaced else '추가'} 완료")

    prs.save(str(PPTX))
    print(f"\n▶ 저장 완료: {PPTX.name}")

    # 슬라이드 제목 목록 출력 (디버깅용)
    print("\n[슬라이드 제목 목록]")
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t and len(t) < 60:
                    print(f"  Slide {i+1}: {t}")
                    break


if __name__ == "__main__":
    main()
